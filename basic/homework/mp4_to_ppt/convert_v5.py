"""
24 Solar Terms - Horizontal Card Layout v5
横向卡片式布局设计
"""
import os
from pptx import Presentation

from design_system_v5 import (
    PALETTE, I, CW, CH, BW,
    cell_bg_unified, cell_border, add_cell_header_v5,
    add_rect, add_oval, add_rounded_rect, push_back
)
from atoms_flat import (
    sun_flat, cloud_flat, hill_flat, ground_flat,
    tree_round, flower_simple, grass_simple, leaf_simple,
    deer_simple, ox_simple, frog_simple, ladybug_simple, butterfly_simple,
    person_flat, kite_simple, umbrella_simple, water_wheel_simple,
    rain_simple, lotus_simple, seedling_simple, lightning_simple,
    bird_simple, dragonfly_simple, watermelon_simple, fan_simple,
    pumpkin_simple, wind_simple, chrysanthemum_simple, birds_flying_simple,
    food_storage_simple, pine_tree_simple, snowman_simple, dumpling_simple,
    plum_blossom_simple, ice_simple, snow_simple
)

C = PALETTE
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "24节气_v5.pptx")

def blank(prs): 
    return prs.slides.add_slide(prs.slide_layouts[6])

def cell_xy(col, row): 
    return col * CW, row * CH

# ============================================================================
# Slide 1: 立春 雨水 惊蛰 春分
# ============================================================================
def slide1(prs):
    s = blank(prs)
    
    # ── 立春 (0,0) - 春回大地 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳 - 右上
    sun_flat(s, cx + CW * 0.78, cy + I(0.25), int(I(0.20)), C["sun_yellow"])
    
    # 桃花树 - 右侧
    tree_round(s, cx + CW * 0.65, cy + CH * 0.25, scale=1.1, leaf_color=C["flower_pink"])
    
    # 小鹿 - 左下
    deer_simple(s, cx + CW * 0.15, cy + CH * 0.48, scale=1.3)
    
    # 小鸟 - 树枝上
    bird_simple(s, cx + CW * 0.58, cy + CH * 0.22, scale=0.8)
    
    # 远山
    hill_flat(s, cx - I(0.1), cy + CH * 0.55, int(I(1.5)), int(I(0.4)), C["spring_green"])
    
    # 房屋 - 远景
    house_w = I(0.5)
    house_h = I(0.35)
    house_x = cx + CW * 0.45
    house_y = cy + CH * 0.52
    # 房屋主体
    add_rect(s, house_x, house_y, house_w, house_h, C["cream_white"])
    # 屋顶
    roof_points = [(house_x, house_y), (house_x + house_w/2, house_y - I(0.15)), (house_x + house_w, house_y)]
    
    # 地面花草
    for fx in [0.1, 0.35, 0.6, 0.85]:
        flower_simple(s, cx + CW * fx, cy + CH * 0.78, scale=0.6, petal_color=C["flower_pink"])
    
    # 底部说明文字
    footer_text = "立春，标志着春天的开始，气温逐渐回升。"
    add_cell_header_v5(s, cx, cy, "立春", "Start of Spring", 
                       "春回大地，万物复苏。", 
                       "Spring begins. It gets warmer.",
                       footer_text, icon_type="sun")
    cell_border(s, 0, 0)
    
    # ── 雨水 (1,0) - 润物无声 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["spring_bg"], C["sky_blue"], C["spring_gnd"])
    
    # 雨云
    cloud_flat(s, cx + CW * 0.55, cy + I(0.12), I(0.9), I(0.35), color=C["cloud_gray"])
    
    # 雨滴
    rain_simple(s, cx + CW * 0.5, cy + I(0.40), I(0.8), I(0.5), n=12)
    
    # 打伞男孩 - 右中
    person_flat(s, cx + CW * 0.65, cy + CH * 0.42, C["rain_blue"], scale=1.0, pose="hold")
    umbrella_simple(s, cx + CW * 0.58, cy + CH * 0.25, scale=1.2)
    
    # 青蛙 - 池塘边
    frog_simple(s, cx + CW * 0.45, cy + CH * 0.65, scale=0.9)
    
    # 水坑
    for px in [0.2, 0.45, 0.75]:
        add_oval(s, cx + CW * px, cy + CH * 0.72, I(0.25), I(0.1), C["rain_blue"])
    
    # 花朵
    for fx in [0.55, 0.75, 0.9]:
        flower_simple(s, cx + CW * fx, cy + CH * 0.52, scale=0.7, petal_color=C["flower_pink"])
    
    # 草地
    grass_simple(s, cx + BW, cy + CH * 0.72, CW - BW*2, I(0.15), C["spring_gnd"])
    
    footer_text = "雨水，表示降水开始，雨量逐渐增多。"
    add_cell_header_v5(s, cx, cy, "雨水", "Rain Water",
                       "降雨增多，润物无声。",
                       "More rain, everything grows.",
                       footer_text, icon_type="rain")
    cell_border(s, 1, 0)
    
    # ── 惊蛰 (0,1) - 虫鸣渐起 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["warm_yellow"], C["sky_blue"], C["spring_gnd"])
    
    # 雷云
    cloud_flat(s, cx + CW * 0.5, cy + I(0.05), I(1.0), I(0.38), color=C["cloud_gray"])
    
    # 闪电
    lightning_simple(s, cx + CW * 0.55, cy + I(0.35), scale=0.9)
    
    # 瓢虫 - 左下大
    ladybug_simple(s, cx + CW * 0.12, cy + CH * 0.58, scale=1.6)
    
    # 蜜蜂/蝴蝶
    butterfly_simple(s, cx + CW * 0.45, cy + CH * 0.42, scale=0.8)
    
    # 青蛙
    frog_simple(s, cx + CW * 0.7, cy + CH * 0.6, scale=1.0)
    
    # 树木
    tree_round(s, cx + CW * 0.75, cy + CH * 0.35, scale=0.9, leaf_color=C["spring_green"])
    
    # 地面植物
    for sx in [0.15, 0.4, 0.65, 0.85]:
        seedling_simple(s, cx + CW * sx, cy + CH * 0.75, scale=0.7)
    
    footer_text = "惊蛰，春雷乍动，惊醒蛰伏的昆虫和动物。"
    add_cell_header_v5(s, cx, cy, "惊蛰", "Awakening of Insects",
                       "春雷惊蛰，虫鸣渐起。",
                       "Thunder wakes sleeping insects.",
                       footer_text, icon_type="thunder")
    cell_border(s, 0, 1)
    
    # ── 春分 (1,1) - 昼夜平分 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳
    sun_flat(s, cx + CW * 0.75, cy + I(0.20), int(I(0.18)), C["sun_yellow"])
    
    # 小女孩观察花朵
    person_flat(s, cx + CW * 0.15, cy + CH * 0.48, C["soft_coral"], scale=1.0, pose="hold")
    
    # 花朵阵列 - 右侧花田
    flower_colors = [C["flower_pink"], C["warm_yellow"], C["spring_green"], C["soft_coral"]]
    for i in range(8):
        fx = 0.45 + (i % 4) * 0.14
        fy = 0.52 + (i // 4) * 0.08
        flower_simple(s, cx + CW * fx, cy + CH * fy, scale=0.65, 
                     petal_color=flower_colors[i % 4])
    
    # 蝴蝶
    butterfly_simple(s, cx + CW * 0.85, cy + CH * 0.35, scale=0.7)
    
    # 太阳月亮平衡装置 (简化)
    pole_x = cx + CW * 0.82
    pole_y = cy + I(0.25)
    # 支柱
    add_rect(s, pole_x, pole_y, I(0.03), I(0.35), C["ink_brown"])
    # 横杆
    add_rect(s, pole_x - I(0.25), pole_y + I(0.05), I(0.53), I(0.02), C["ink_brown"])
    # 太阳盘
    add_oval(s, pole_x - I(0.30), pole_y - I(0.05), I(0.15), I(0.10), C["sun_yellow"])
    # 月亮盘
    add_oval(s, pole_x + I(0.15), pole_y - I(0.05), I(0.15), I(0.10), C["cream_white"])
    
    footer_text = "春分，昼夜平分，春意正浓，万物茂盛。"
    add_cell_header_v5(s, cx, cy, "春分", "Spring Equinox",
                       "昼夜平分，生机盎然。",
                       "Day and night are equal length.",
                       footer_text, icon_type="sun")
    cell_border(s, 1, 1)

# ============================================================================
# Slide 2: 清明 谷雨 立夏 小满
# ============================================================================
def slide2(prs):
    s = blank(prs)
    
    # ── 清明 (0,0) - 踏青扫墓 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳
    sun_flat(s, cx + CW * 0.78, cy + I(0.22), int(I(0.18)), C["sun_yellow"])
    
    # 柳树/树木
    tree_round(s, cx + CW * 0.65, cy + CH * 0.22, scale=1.0, leaf_color=C["spring_green"])
    
    # 放风筝的小孩
    kite_simple(s, cx + CW * 0.75, cy + I(0.15), scale=0.8)
    person_flat(s, cx + CW * 0.15, cy + CH * 0.48, C["soft_coral"], scale=0.9)
    
    # 远山
    hill_flat(s, cx + CW * 0.3, cy + CH * 0.45, I(1.2), I(0.35), C["spring_green"])
    
    # 地面花草
    for fx in [0.12, 0.35, 0.6, 0.85]:
        flower_simple(s, cx + CW * fx, cy + CH * 0.72, scale=0.6, petal_color=C["flower_pink"])
    
    footer_text = "清明，天气晴朗，草木繁茂，适合踏青郊游。"
    add_cell_header_v5(s, cx, cy, "清明", "Clear and Bright",
                       "天清气明，踏青扫墓。",
                       "Sunny day for spring outings.",
                       footer_text, icon_type="sun")
    cell_border(s, 0, 0)
    
    # ── 谷雨 (1,0) - 雨生百谷 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["spring_bg"], C["sky_blue"], C["spring_gnd"])
    
    # 雨滴
    rain_simple(s, cx + CW * 0.2, cy + I(0.25), I(0.7), I(0.5), n=15)
    
    # 稻田/幼苗
    for sx in [0.15, 0.35, 0.55, 0.75]:
        seedling_simple(s, cx + CW * sx, cy + CH * 0.65, scale=0.8)
    
    # 农民
    person_flat(s, cx + CW * 0.45, cy + CH * 0.42, C["wheat_gold"], scale=1.0)
    
    # 远山
    hill_flat(s, cx + CW * 0.2, cy + CH * 0.48, I(1.5), I(0.32), C["spring_green"])
    
    footer_text = "谷雨，雨生百谷，是播种移苗的好时节。"
    add_cell_header_v5(s, cx, cy, "谷雨", "Grain Rain",
                       "雨润百谷，万物生长。",
                       "Rain helps crops grow well.",
                       footer_text, icon_type="rain")
    cell_border(s, 1, 0)
    
    # ── 立夏 (0,1) - 夏日初临 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 太阳
    sun_flat(s, cx + CW * 0.75, cy + I(0.18), int(I(0.22)), C["sun_yellow"])
    
    # 大树
    tree_round(s, cx + CW * 0.6, cy + CH * 0.25, scale=1.2, leaf_color=C["spring_green"])
    
    # 荷花池
    add_oval(s, cx + CW * 0.15, cy + CH * 0.58, I(0.6), I(0.2), C["rain_blue"])
    lotus_simple(s, cx + CW * 0.25, cy + CH * 0.48, scale=0.7)
    lotus_simple(s, cx + CW * 0.45, cy + CH * 0.52, scale=0.6)
    
    # 人物
    person_flat(s, cx + CW * 0.72, cy + CH * 0.45, C["soft_coral"], scale=0.9)
    
    footer_text = "立夏，夏季开始，万物进入生长旺季。"
    add_cell_header_v5(s, cx, cy, "立夏", "Start of Summer",
                       "夏日初临，万物生长。",
                       "Summer comes. Plants grow fast.",
                       footer_text, icon_type="sun")
    cell_border(s, 0, 1)
    
    # ── 小满 (1,1) - 麦粒渐满 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 太阳
    sun_flat(s, cx + CW * 0.2, cy + I(0.15), int(I(0.20)), C["sun_yellow"])
    
    # 麦田
    hill_flat(s, cx + CW * 0.4, cy + CH * 0.38, I(1.0), I(0.4), C["wheat_gold"])
    
    # 水车
    from atoms_flat import water_wheel_simple
    water_wheel_simple(s, cx + CW * 0.65, cy + CH * 0.35, scale=0.9)
    
    # 河流
    add_oval(s, cx + CW * 0.45, cy + CH * 0.68, I(0.8), I(0.15), C["rain_blue"])
    
    footer_text = "小满，麦粒渐满，但尚未成熟。"
    add_cell_header_v5(s, cx, cy, "小满", "Grain Full",
                       "麦粒渐满，丰收在望。",
                       "Grains start to become full.",
                       footer_text, icon_type="seedling")
    cell_border(s, 1, 1)

# ============================================================================
# Slide 3: 芒种 夏至 小暑 大暑
# ============================================================================
def slide3(prs):
    s = blank(prs)
    
    # ── 芒种 (0,0) - 麦熟仲夏 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 太阳
    sun_flat(s, cx + CW * 0.75, cy + I(0.15), int(I(0.22)), C["sun_yellow"])
    
    # 麦田
    hill_flat(s, cx + CW * 0.25, cy + CH * 0.42, I(1.1), I(0.38), C["wheat_gold"])
    
    # 收割的农民
    person_flat(s, cx + CW * 0.62, cy + CH * 0.45, C["soft_coral"], scale=1.0)
    
    # 麦穗
    for sx in [0.12, 0.32, 0.52, 0.72]:
        seedling_simple(s, cx + CW * sx, cy + CH * 0.65, scale=0.9)
    
    footer_text = "芒种，麦类成熟，稻谷播种，忙收忙种。"
    add_cell_header_v5(s, cx, cy, "芒种", "Grain in Ear",
                       "麦熟仲夏，忙收忙种。",
                       "Wheat matures, rice is planted.",
                       footer_text, icon_type="seedling")
    cell_border(s, 0, 0)
    
    # ── 夏至 (1,0) - 昼最长夜最短 ─────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 大太阳 - 最高位置
    sun_flat(s, cx + CW * 0.5, cy + I(0.08), int(I(0.28)), C["sun_yellow"])
    
    # 荷花盛开
    add_oval(s, cx + CW * 0.2, cy + CH * 0.58, I(0.7), I(0.22), C["rain_blue"])
    lotus_simple(s, cx + CW * 0.35, cy + CH * 0.48, scale=1.0)
    lotus_simple(s, cx + CW * 0.55, cy + CH * 0.52, scale=0.8)
    
    # 荷叶
    from atoms_flat import leaf_simple
    leaf_simple(s, cx + CW * 0.25, cy + CH * 0.62, scale=0.8, color=C["spring_green"])
    leaf_simple(s, cx + CW * 0.60, cy + CH * 0.60, scale=0.7, color=C["spring_green"])
    
    # 蜻蜓
    from atoms_flat import dragonfly_simple
    dragonfly_simple(s, cx + CW * 0.72, cy + CH * 0.35, scale=0.8)
    
    footer_text = "夏至，北半球白昼最长，阳气最盛。"
    add_cell_header_v5(s, cx, cy, "夏至", "Summer Solstice",
                       "夏至已至，昼最长夜最短。",
                       "Longest day of the year.",
                       footer_text, icon_type="sun")
    cell_border(s, 1, 0)
    
    # ── 小暑 (0,1) - 初伏将至 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 烈日
    sun_flat(s, cx + CW * 0.2, cy + I(0.12), int(I(0.24)), C["sun_yellow"])
    
    # 热浪效果 - 用浅色椭圆表示
    for i, hx in enumerate([0.3, 0.5, 0.7]):
        add_oval(s, cx + CW * hx, cy + CH * (0.25 + i*0.08), I(0.15), I(0.06), 
                 rgb(0xF8, 0xE8, 0xC0))
    
    # 树荫下乘凉
    tree_round(s, cx + CW * 0.65, cy + CH * 0.15, scale=1.3, leaf_color=C["spring_green"])
    person_flat(s, cx + CW * 0.55, cy + CH * 0.52, C["rain_blue"], scale=0.9)
    
    # 西瓜
    from atoms_flat import watermelon_simple
    watermelon_simple(s, cx + CW * 0.25, cy + CH * 0.65, scale=0.8)
    
    # 扇子
    from atoms_flat import fan_simple
    fan_simple(s, cx + CW * 0.72, cy + CH * 0.45, scale=0.6)
    
    footer_text = "小暑，天气开始炎热，但还没到最热。"
    add_cell_header_v5(s, cx, cy, "小暑", "Minor Heat",
                       "初伏将至，暑气渐盛。",
                       "The weather becomes hot.",
                       footer_text, icon_type="sun")
    cell_border(s, 0, 1)
    
    # ── 大暑 (1,1) - 最热时节 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    # 烈日当空
    sun_flat(s, cx + CW * 0.75, cy + I(0.10), int(I(0.26)), C["sun_yellow"])
    
    # 池塘 - 消暑
    add_oval(s, cx + CW * 0.15, cy + CH * 0.58, I(0.8), I(0.25), C["rain_blue"])
    
    # 荷花
    lotus_simple(s, cx + CW * 0.30, cy + CH * 0.48, scale=0.9)
    lotus_simple(s, cx + CW * 0.55, cy + CH * 0.45, scale=0.7)
    
    # 青蛙
    frog_simple(s, cx + CW * 0.72, cy + CH * 0.65, scale=1.0)
    
    # 萤火虫 (用黄色小圆点表示)
    for fx, fy in [(0.45, 0.35), (0.65, 0.42), (0.82, 0.38), (0.55, 0.28)]:
        add_oval(s, cx + CW * fx, cy + CH * fy, I(0.05), I(0.05), C["sun_yellow"])
    
    footer_text = "大暑，一年中最热的时期，热浪滚滚。"
    add_cell_header_v5(s, cx, cy, "大暑", "Major Heat",
                       "最热时节，荷塘消暑。",
                       "The hottest time of year.",
                       footer_text, icon_type="sun")
    cell_border(s, 1, 1)

# ============================================================================
# Slide 4: 立秋 处暑 白露 秋分
# ============================================================================
def slide4(prs):
    s = blank(prs)
    
    # ── 立秋 (0,0) - 秋意渐起 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["autumn_bg"], C["autumn_sky"], C["autumn_gnd"])
    
    # 秋阳
    sun_flat(s, cx + CW * 0.78, cy + I(0.18), int(I(0.20)), C["sun_yellow"])
    
    # 金黄树叶的树
    tree_round(s, cx + CW * 0.6, cy + CH * 0.22, scale=1.1, leaf_color=C["autumn_amber"])
    
    # 落叶
    for lx, ly in [(0.25, 0.65), (0.45, 0.72), (0.70, 0.68), (0.85, 0.75)]:
        leaf_simple(s, cx + CW * lx, cy + CH * ly, scale=0.5, color=C["autumn_amber"])
    
    # 南瓜/果实
    from atoms_flat import pumpkin_simple
    pumpkin_simple(s, cx + CW * 0.2, cy + CH * 0.65, scale=0.7)
    
    # 凉风
    from atoms_flat import wind_simple
    wind_simple(s, cx + CW * 0.72, cy + CH * 0.35, scale=0.8)
    
    footer_text = "立秋，秋季开始，暑去凉来，树叶渐黄。"
    add_cell_header_v5(s, cx, cy, "立秋", "Start of Autumn",
                       "秋意渐起，暑去凉来。",
                       "Autumn begins. Cool weather arrives.",
                       footer_text, icon_type="seedling")
    cell_border(s, 0, 0)
    
    # ── 处暑 (1,0) - 暑气终止 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["autumn_bg"], C["autumn_sky"], C["autumn_gnd"])
    
    # 温和的阳光
    sun_flat(s, cx + CW * 0.75, cy + I(0.22), int(I(0.18)), C["sun_yellow"])
    
    # 稻田金黄
    hill_flat(s, cx + CW * 0.2, cy + CH * 0.45, I(1.3), I(0.35), C["wheat_gold"])
    
    # 农民收割
    person_flat(s, cx + CW * 0.55, cy + CH * 0.42, C["soft_coral"], scale=1.0)
    
    # 麦穗
    for sx in [0.15, 0.35, 0.55, 0.75]:
        seedling_simple(s, cx + CW * sx, cy + CH * 0.68, scale=0.8)
    
    footer_text = "处暑，炎热暑天结束，气温逐渐下降。"
    add_cell_header_v5(s, cx, cy, "处暑", "End of Heat",
                       "暑气终止，凉风送爽。",
                       "Hot summer days are over.",
                       footer_text, icon_type="cloud")
    cell_border(s, 1, 0)
    
    # ── 白露 (0,1) - 晨露晶莹 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["autumn_bg"], C["autumn_sky"], C["autumn_gnd"])
    
    # 清晨薄雾感 - 浅白色覆盖
    add_oval(s, cx + CW * 0.3, cy + CH * 0.15, I(0.8), I(0.25), rgb(0xF0, 0xF0, 0xF0))
    
    # 露珠 (蓝色小圆点)
    for dx, dy in [(0.22, 0.62), (0.38, 0.58), (0.55, 0.65), (0.72, 0.60), (0.85, 0.68)]:
        add_oval(s, cx + CW * dx, cy + CH * dy, I(0.06), I(0.06), C["rain_blue"])
    
    # 植物叶片
    leaf_simple(s, cx + CW * 0.3, cy + CH * 0.45, scale=0.9, color=C["spring_green"])
    leaf_simple(s, cx + CW * 0.65, cy + CH * 0.42, scale=0.8, color=C["spring_green"])
    
    # 芦苇/草
    grass_simple(s, cx + BW, cy + CH * 0.70, CW - BW*2, I(0.12), C["autumn_gnd"])
    
    footer_text = "白露，夜间水汽凝结成露，天气转凉。"
    add_cell_header_v5(s, cx, cy, "白露", "White Dew",
                       "晨露晶莹，天气转凉。",
                       "Dew appears on plants at night.",
                       footer_text, icon_type="rain")
    cell_border(s, 0, 1)
    
    # ── 秋分 (1,1) - 昼夜平分 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["autumn_bg"], C["autumn_sky"], C["autumn_gnd"])
    
    # 平衡的天平 (同春分但秋色)
    pole_x = cx + CW * 0.82
    pole_y = cy + I(0.28)
    add_rect(s, pole_x, pole_y, I(0.03), I(0.30), C["ink_brown"])
    add_rect(s, pole_x - I(0.22), pole_y + I(0.04), I(0.47), I(0.02), C["ink_brown"])
    add_oval(s, pole_x - I(0.26), pole_y - I(0.04), I(0.13), I(0.09), C["sun_yellow"])
    add_oval(s, pole_x + I(0.13), pole_y - I(0.04), I(0.13), I(0.09), C["cream_white"])
    
    # 秋月/菊花
    from atoms_flat import chrysanthemum_simple
    chrysanthemum_simple(s, cx + CW * 0.2, cy + CH * 0.55, scale=0.9)
    chrysanthemum_simple(s, cx + CW * 0.4, cy + CH * 0.62, scale=0.7)
    
    # 落叶
    for lx, ly in [(0.55, 0.72), (0.75, 0.68), (0.88, 0.75)]:
        leaf_simple(s, cx + CW * lx, cy + CH * ly, scale=0.5, color=C["autumn_amber"])
    
    footer_text = "秋分，昼夜平分，秋色宜人，丹桂飘香。"
    add_cell_header_v5(s, cx, cy, "秋分", "Autumn Equinox",
                       "昼夜平分，秋色正浓。",
                       "Day and night are equal length.",
                       footer_text, icon_type="sun")
    cell_border(s, 1, 1)

# ============================================================================
# Slide 5: 寒露 霜降 立冬 小雪
# ============================================================================
def slide5(prs):
    s = blank(prs)
    
    # ── 寒露 (0,0) - 露水更凉 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["autumn_bg"], C["autumn_sky"], C["autumn_gnd"])
    
    # 较弱的阳光
    sun_flat(s, cx + CW * 0.75, cy + I(0.20), int(I(0.16)), C["sun_yellow"])
    
    # 更大的露珠
    for dx, dy in [(0.18, 0.58), (0.35, 0.65), (0.52, 0.55), (0.70, 0.62), (0.85, 0.70)]:
        add_oval(s, cx + CW * dx, cy + CH * dy, I(0.08), I(0.08), C["rain_blue"])
    
    # 渐黄的草
    grass_simple(s, cx + BW, cy + CH * 0.68, CW - BW*2, I(0.15), C["autumn_gnd"])
    
    # 大雁南飞 (简化为三角形)
    from atoms_flat import birds_flying_simple
    birds_flying_simple(s, cx + CW * 0.45, cy + CH * 0.25, scale=0.7)
    
    footer_text = "寒露，气温更低，露水更凉，将要凝结。"
    add_cell_header_v5(s, cx, cy, "寒露", "Cold Dew",
                       "露气寒冷，大雁南飞。",
                       "Dew turns cold. Geese fly south.",
                       footer_text, icon_type="rain")
    cell_border(s, 0, 0)
    
    # ── 霜降 (1,0) - 初霜出现 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["autumn_bg"], C["winter_sky"], C["autumn_gnd"])
    
    # 白色霜层覆盖效果
    add_oval(s, cx + CW * 0.2, cy + CH * 0.58, I(0.7), I(0.20), rgb(0xF0, 0xF8, 0xFF))
    
    # 霜叶红于二月花
    leaf_simple(s, cx + CW * 0.25, cy + CH * 0.45, scale=1.0, color=C["soft_coral"])
    leaf_simple(s, cx + CW * 0.55, cy + CH * 0.52, scale=0.8, color=C["autumn_amber"])
    leaf_simple(s, cx + CW * 0.75, cy + CH * 0.42, scale=0.9, color=rgb(0xC8, 0x50, 0x50))
    
    # 落霜的地面
    for sx in [0.15, 0.35, 0.55, 0.75]:
        add_oval(s, cx + CW * sx, cy + CH * 0.72, I(0.12), I(0.05), rgb(0xF8, 0xF8, 0xFF))
    
    footer_text = "霜降，初霜出现，秋意最浓，气温骤降。"
    add_cell_header_v5(s, cx, cy, "霜降", "Frost's Descent",
                       "初霜出现，秋尽冬来。",
                       "First frost appears.",
                       footer_text, icon_type="cloud")
    cell_border(s, 1, 0)
    
    # ── 立冬 (0,1) - 冬季开始 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 苍白太阳
    sun_flat(s, cx + CW * 0.75, cy + I(0.18), int(I(0.18)), rgb(0xE8, 0xE0, 0xD0))
    
    # 枯树
    tree_round(s, cx + CW * 0.6, cy + CH * 0.25, scale=1.0, leaf_color=rgb(0xA8, 0x90, 0x70))
    
    # 落叶堆积
    for lx, ly in [(0.2, 0.72), (0.4, 0.78), (0.6, 0.75), (0.8, 0.80)]:
        leaf_simple(s, cx + CW * lx, cy + CH * ly, scale=0.6, color=C["autumn_amber"])
    
    # 准备过冬的动物/食物储存
    from atoms_flat import food_storage_simple
    food_storage_simple(s, cx + CW * 0.25, cy + CH * 0.55, scale=0.8)
    
    footer_text = "立冬，冬季开始，万物收藏，规避寒冷。"
    add_cell_header_v5(s, cx, cy, "立冬", "Start of Winter",
                       "冬之始，万物收藏。",
                       "Winter begins. Nature stores energy.",
                       footer_text, icon_type="cloud")
    cell_border(s, 0, 1)
    
    # ── 小雪 (1,1) - 初雪飘落 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 雪花飘落
    from atoms_flat import snow_simple
    snow_simple(s, cx + CW * 0.1, cy + I(0.15), I(0.8), I(0.4), n=8)
    
    # 积薄雪
    add_oval(s, cx + CW * 0.15, cy + CH * 0.65, I(0.25), I(0.08), rgb(0xF8, 0xF8, 0xFF))
    add_oval(s, cx + CW * 0.55, cy + CH * 0.68, I(0.35), I(0.10), rgb(0xF8, 0xF8, 0xFF))
    
    # 冬季树木
    tree_round(s, cx + CW * 0.65, cy + CH * 0.25, scale=0.9, leaf_color=rgb(0xA8, 0xA0, 0x98))
    
    footer_text = "小雪，开始降雪，但雪量不大。"
    add_cell_header_v5(s, cx, cy, "小雪", "Minor Snow",
                       "初雪飘落，寒意渐浓。",
                       "Light snow begins to fall.",
                       footer_text, icon_type="cloud")
    cell_border(s, 1, 1)

# ============================================================================
# Slide 6: 大雪 冬至 小寒 大寒
# ============================================================================
def slide6(prs):
    s = blank(prs)
    
    # ── 大雪 (0,0) - 雪量增多 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 大雪纷飞
    snow_simple(s, cx + CW * 0.1, cy + I(0.10), I(0.85), I(0.45), n=15)
    
    # 厚积雪
    add_oval(s, cx + CW * 0.1, cy + CH * 0.62, I(0.8), I(0.18), rgb(0xF5, 0xF8, 0xFF))
    
    # 雪松
    from atoms_flat import pine_tree_simple
    pine_tree_simple(s, cx + CW * 0.65, cy + CH * 0.25, scale=1.1)
    
    # 雪人
    from atoms_flat import snowman_simple
    snowman_simple(s, cx + CW * 0.25, cy + CH * 0.55, scale=0.9)
    
    footer_text = "大雪，雪量增多，地面可能积雪，天气寒冷。"
    add_cell_header_v5(s, cx, cy, "大雪", "Major Snow",
                       "雪量增多，银装素裹。",
                       "Heavy snow falls.",
                       footer_text, icon_type="cloud")
    cell_border(s, 0, 0)
    
    # ── 冬至 (1,0) - 昼最短夜最长 ─────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 低挂的苍白太阳
    sun_flat(s, cx + CW * 0.5, cy + CH * 0.38, int(I(0.16)), rgb(0xE0, 0xD8, 0xC8))
    
    # 雪地
    add_oval(s, cx + CW * 0.1, cy + CH * 0.68, I(0.85), I(0.15), rgb(0xF0, 0xF5, 0xFA))
    
    # 冬季松树
    pine_tree_simple(s, cx + CW * 0.25, cy + CH * 0.35, scale=0.9)
    pine_tree_simple(s, cx + CW * 0.75, cy + CH * 0.32, scale=1.0)
    
    # 饺子/汤圆 (冬至食物)
    from atoms_flat import dumpling_simple
    dumpling_simple(s, cx + CW * 0.45, cy + CH * 0.58, scale=0.6)
    dumpling_simple(s, cx + CW * 0.55, cy + CH * 0.62, scale=0.5)
    dumpling_simple(s, cx + CW * 0.50, cy + CH * 0.52, scale=0.55)
    
    footer_text = "冬至，白昼最短，黑夜最长，数九寒天开始。"
    add_cell_header_v5(s, cx, cy, "冬至", "Winter Solstice",
                       "昼最短夜最长，数九开始。",
                       "Shortest day of the year.",
                       footer_text, icon_type="sun")
    cell_border(s, 1, 0)
    
    # ── 小寒 (0,1) - 初寒未极 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 寒风效果
    wind_simple(s, cx + CW * 0.65, cy + CH * 0.25, scale=1.0)
    
    # 积雪覆盖
    add_oval(s, cx + CW * 0.12, cy + CH * 0.65, I(0.8), I(0.18), rgb(0xF0, 0xF5, 0xFA))
    
    # 松树
    pine_tree_simple(s, cx + CW * 0.65, cy + CH * 0.28, scale=1.0)
    
    # 梅花 (冬季开花)
    from atoms_flat import plum_blossom_simple
    plum_blossom_simple(s, cx + CW * 0.25, cy + CH * 0.45, scale=0.8)
    plum_blossom_simple(s, cx + CW * 0.40, cy + CH * 0.35, scale=0.6)
    
    footer_text = "小寒，天气寒冷但未到极点，梅花初绽。"
    add_cell_header_v5(s, cx, cy, "小寒", "Minor Cold",
                       "初寒未极，梅花吐蕊。",
                       "Cold but not the coldest.",
                       footer_text, icon_type="cloud")
    cell_border(s, 0, 1)
    
    # ── 大寒 (1,1) - 一年最冷 ─────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["winter_bg"], C["winter_sky"], C["winter_gnd"])
    
    # 极弱的阳光
    sun_flat(s, cx + CW * 0.78, cy + I(0.25), int(I(0.14)), rgb(0xD8, 0xD0, 0xC0))
    
    # 厚雪覆盖
    add_oval(s, cx + CW * 0.08, cy + CH * 0.60, I(0.9), I(0.25), rgb(0xF5, 0xF8, 0xFF))
    
    # 冰雪覆盖的松树
    pine_tree_simple(s, cx + CW * 0.6, cy + CH * 0.22, scale=1.1)
    
    # 雪花
    snow_simple(s, cx + CW * 0.15, cy + I(0.12), I(0.7), I(0.35), n=12)
    
    # 冰上活动提示
    from atoms_flat import ice_simple
    ice_simple(s, cx + CW * 0.3, cy + CH * 0.55, scale=0.7)
    
    footer_text = "大寒，一年中最冷的时期，冰天雪地。"
    add_cell_header_v5(s, cx, cy, "大寒", "Major Cold",
                       "一年最冷，静待春归。",
                       "The coldest time of year.",
                       footer_text, icon_type="cloud")
    cell_border(s, 1, 1)

def rgb(r,g,b):
    from pptx.dml.color import RGBColor
    return RGBColor(r,g,b)

# ============================================================================
# 主函数
# ============================================================================
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    
    prs.save(OUT)
    print(f"Saved: {OUT}")

if __name__ == "__main__":
    from design_system_v5 import SW, SH
    main()
