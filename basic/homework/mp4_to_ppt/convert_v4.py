"""
24 Solar Terms - Scandinavian Flat Design v4
统一视觉系统：标题安全区 + 对角线构图 + 扁平插画
"""
import os, math
from pptx import Presentation
from pptx.util import Inches

from design_system import (
    PALETTE, I, add_rect, add_oval, add_text, add_cell_header,
    push_back, SAFE_MARGIN_LEFT, SAFE_MARGIN_TOP, HORIZON_Y_RATIO,
    Composition
)
from atoms_flat import (
    sun_flat, cloud_flat, hill_flat, ground_flat,
    tree_round, flower_simple, grass_simple,
    deer_simple, ox_simple, frog_simple, ladybug_simple, butterfly_simple,
    person_flat, house_simple, kite_simple, umbrella_simple, water_wheel_simple,
    rain_simple, lotus_simple, seedling_simple, lightning_simple,
    bird_simple
)

C = PALETTE

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "24节气_v4.pptx")
SW = Inches(13.33); SH = Inches(7.50)
CW = SW // 2;       CH = SH // 2
BW = int(Inches(0.03))   # 更细边框

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])
def cell_xy(col, row): return col*CW, row*CH

# ============================================================================
# Cell 背景系统 - 统一地平线
# ============================================================================
def cell_bg_unified(slide, col, row, bg_color, sky_color=None, ground_color=None):
    """统一 cell 背景 - 固定地平线位置"""
    cx, cy = cell_xy(col, row)
    
    # 背景色
    bg = add_rect(slide, cx, cy, CW, CH, bg_color)
    push_back(slide, bg)
    
    # 天空区域
    if sky_color:
        horizon_y = cy + int(CH * HORIZON_Y_RATIO)
        sky = add_rect(slide, cx, cy, CW, horizon_y - cy, sky_color)
        push_back(slide, sky)
    
    # 地面区域
    if ground_color:
        horizon_y = cy + int(CH * HORIZON_Y_RATIO)
        ground_h = cy + CH - horizon_y
        ground = add_rect(slide, cx, horizon_y, CW, ground_h, ground_color)
        push_back(slide, ground)
    
    return cx, cy

def cell_border(slide, col, row):
    """统一白色边框"""
    cx, cy = cell_xy(col, row)
    for rx,ry,rw,rh in [(cx,cy,CW,BW),(cx,cy+CH-BW,CW,BW),
                         (cx,cy,BW,CH),(cx+CW-BW,cy,BW,CH)]:
        add_rect(slide, rx, ry, rw, rh, C["frost_white"])

# ============================================================================
# Slide 1: 立春 雨水 惊蛰 春分
# 构图模式：对角线构图 (左上→右下)
# 标题区：左上安全区
# 视觉锚点：右上太阳 + 右下主体
# ============================================================================
def slide1(prs):
    s = blank(prs)
    comp = Composition()
    
    # ── 立春 (0,0) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳 - 统一位置：右上
    sun_x, sun_y = comp.get_sun_position(cx, cy, CW, CH)
    sun_flat(s, sun_x, sun_y, int(I(0.22)))
    
    # 远景山丘 (对角线左下→右上，避开标题区)
    hill_flat(s, cx-int(I(0.15)), cy+int(CH*0.45), int(I(2.0)), int(I(0.55)), 
              C["spring_green"])
    
    # 牛 - 右下主体 (对角线终点)
    ox_simple(s, cx+int(I(1.20)), cy+int(CH*0.58), scale=1.6)
    
    # 人物骑牛 - 简化位置
    person_flat(s, cx+int(I(1.48)), cy+int(CH*0.42), C["soft_coral"], scale=0.70)
    
    # 风筝 - 右上天空 (靠近太阳但避开标题)
    kite_simple(s, cx+int(I(1.80)), cy+int(I(0.25)), scale=0.85)
    
    # 前景树 - 左侧中景
    tree_round(s, cx+int(I(0.25)), cy+int(CH*0.38), scale=1.3, leaf_color=C["spring_green"])
    
    # 地面小花点缀
    for fx in [0.15, 0.55, 1.05, 1.55, 2.0]:
        flower_simple(s, cx+int(I(fx)), cy+int(CH*0.72), scale=0.7, petal_color=C["flower_pink"])
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "立春", "Start of Spring", "Spring begins. It gets warmer.", comp)
    cell_border(s, 0, 0)
    
    # ── 雨水 (1,0) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳 (弱化显示，多云天气)
    sun_flat(s, cx+CW-int(I(0.42)), cy+int(I(0.12)), int(I(0.15)), C["warm_yellow"])
    
    # 雨云 - 右上区域 (避开标题)
    cloud_flat(s, cx+int(I(1.50)), cy+int(I(0.08)), int(I(0.80)), int(I(0.35)), 
               color=C["frost_white"])
    
    # 雨滴 - 仅在云下区域
    rain_simple(s, cx+int(I(1.20)), cy+int(I(0.38)), int(I(0.90)), int(I(0.45)), n=10)
    
    # 远景山
    hill_flat(s, cx+int(I(0.30)), cy+int(CH*0.48), int(I(1.80)), int(I(0.38)), 
              C["spring_green"])
    
    # 雨伞人物 - 中右主体
    umbrella_simple(s, cx+int(I(1.10)), cy+int(CH*0.35), scale=1.4)
    person_flat(s, cx+int(I(0.95)), cy+int(CH*0.48), C["rain_blue"], scale=1.0)
    
    # 池塘与荷花 - 右下角
    add_oval(s, cx+int(I(1.50)), cy+int(CH*0.62), int(I(0.75)), int(I(0.22)), C["rain_blue"])
    lotus_simple(s, cx+int(I(1.55)), cy+int(CH*0.52), scale=1.0)
    lotus_simple(s, cx+int(I(1.95)), cy+int(CH*0.56), scale=0.75)
    
    # 青蛙 - 池塘边
    frog_simple(s, cx+int(I(1.40)), cy+int(CH*0.62), scale=0.80)
    
    # 地面雨坑
    for px in [0.20, 0.70, 1.30, 1.85]:
        add_oval(s, cx+int(I(px)), cy+int(CH*0.75), int(I(0.18)), int(I(0.08)), C["rain_blue"])
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "雨水", "Rain Water", "More rain helps everything grow.", comp)
    cell_border(s, 1, 0)
    
    # ── 惊蛰 (0,1) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["warm_yellow"], C["sky_mint"], C["spring_gnd"])
    
    # 太阳半露
    sun_flat(s, cx+CW-int(I(0.35)), cy+int(I(0.10)), int(I(0.18)))
    
    # 雷云 - 右上 (避开标题区)
    cloud_flat(s, cx+int(I(0.80)), cy-int(I(0.02)), int(I(1.20)), int(I(0.42)), 
               color=C["light_gray"])
    cloud_flat(s, cx+int(I(1.40)), cy+int(I(0.05)), int(I(0.70)), int(I(0.32)), 
               color=C["warm_gray"])
    
    # 闪电 - 右侧
    lightning_simple(s, cx+int(I(1.35)), cy+int(I(0.32)), scale=1.0)
    lightning_simple(s, cx+int(I(1.75)), cy+int(I(0.38)), scale=0.7)
    
    # 远景山丘
    hill_flat(s, cx-int(I(0.10)), cy+int(CH*0.42), int(I(1.60)), int(I(0.40)), C["spring_green"])
    hill_flat(s, cx+int(I(0.90)), cy+int(CH*0.46), int(I(1.50)), int(I(0.35)), C["spring_green"])
    
    # 瓢虫 - 左下主体 (大)
    ladybug_simple(s, cx+int(I(0.20)), cy+int(CH*0.58), scale=1.8)
    
    # 青蛙 - 右下
    frog_simple(s, cx+int(I(1.30)), cy+int(CH*0.56), scale=1.2)
    
    # 蝴蝶 - 空中飞舞
    butterfly_simple(s, cx+int(I(1.00)), cy+int(CH*0.32), scale=1.0)
    butterfly_simple(s, cx+int(I(1.60)), cy+int(CH*0.38), scale=0.7)
    
    # 地面苏醒的小生物暗示 - 蚯蚓/虫子
    add_oval(s, cx+int(I(0.75)), cy+int(CH*0.68), int(I(0.14)), int(I(0.06)), C["ink_brown"])
    add_oval(s, cx+int(I(0.78)), cy+int(CH*0.66), int(I(0.06)), int(I(0.06)), C["ink_brown"])
    
    # 幼苗
    for sx in [0.10, 0.35, 0.65, 1.00, 1.35, 1.70, 2.0]:
        seedling_simple(s, cx+int(I(sx)), cy+int(CH*0.72), scale=0.9)
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "惊蛰", "Awakening of Insects", "Thunder wakes sleeping insects.", comp)
    cell_border(s, 0, 1)
    
    # ── 春分 (1,1) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["spring_bg"], C["spring_sky"], C["spring_gnd"])
    
    # 太阳 - 右上黄金位置
    sun_x, sun_y = comp.get_sun_position(cx, cy, CW, CH)
    sun_flat(s, sun_x, sun_y, int(I(0.24)))
    
    # 平衡的象征 - 天平 (右中)
    pole_x = cx + int(I(1.85))
    pole_y = cy + int(I(0.25))
    add_rect(s, pole_x, pole_y, int(I(0.04)), int(I(0.45)), C["ink_brown"])
    add_rect(s, pole_x-int(I(0.28)), pole_y+int(I(0.02)), int(I(0.60)), int(I(0.03)), C["ink_brown"])
    # 太阳盘
    add_oval(s, pole_x-int(I(0.32)), pole_y-int(I(0.02)), int(I(0.18)), int(I(0.12)), C["warm_yellow"])
    # 月亮盘
    add_oval(s, pole_x+int(I(0.14)), pole_y-int(I(0.02)), int(I(0.18)), int(I(0.12)), C["frost_white"])
    
    # 远景花田山丘
    hill_flat(s, cx-int(I(0.05)), cy+int(CH*0.38), int(I(2.20)), int(I(0.42)), C["spring_green"])
    
    # 前景花田 - 密集花朵
    flower_colors = [C["flower_pink"], C["soft_coral"], C["warm_yellow"], C["spring_green"]]
    for i in range(10):
        fx = 0.08 + i * 0.22
        fy = 0.55 + (i % 2) * 0.06
        flower_simple(s, cx+int(I(fx)), cy+int(CH*fy), scale=0.8, petal_color=flower_colors[i % 4])
    
    # 人物在花田 - 左侧
    person_flat(s, cx+int(I(0.15)), cy+int(CH*0.48), C["soft_coral"], scale=1.0, pose="hold")
    
    # 蝴蝶飞舞
    butterfly_simple(s, cx+int(I(0.70)), cy+int(CH*0.30), scale=1.1)
    butterfly_simple(s, cx+int(I(1.80)), cy+int(CH*0.35), scale=0.9)
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "春分", "Spring Equinox", "Day and night are equal length.", comp)
    cell_border(s, 1, 1)

# ============================================================================
# Slide 2: 清明 谷雨 立夏 小满
# ============================================================================
def slide2(prs):
    s = blank(prs)
    comp = Composition()
    
    # ── 清明 (0,0) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 0, C["sky_mint"], C["spring_sky"], C["spring_gnd"])
    
    sun_flat(s, cx+CW-int(I(0.40)), cy+int(I(0.15)), int(I(0.20)))
    
    # 樱花/杏花背景
    hill_flat(s, cx+int(I(0.10)), cy+int(CH*0.32), int(I(1.80)), int(I(0.65)), C["flower_pink"])
    
    # 树
    tree_round(s, cx+int(I(0.45)), cy+int(CH*0.25), scale=1.0, leaf_color=C["spring_green"])
    
    # 房屋
    house_simple(s, cx+int(I(0.80)), cy+int(CH*0.28), scale=0.9)
    
    # 放风筝
    kite_simple(s, cx+int(I(0.20)), cy+int(I(0.15)), scale=0.9)
    person_flat(s, cx+int(I(0.25)), cy+int(CH*0.55), C["soft_coral"], scale=0.85)
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "清明", "Clear and Bright", "Sunny day for spring outings.", comp)
    cell_border(s, 0, 0)
    
    # ── 谷雨 (1,0) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 0, C["spring_bg"], C["rain_blue"], C["spring_gnd"])
    
    # 雨滴
    rain_simple(s, cx+int(I(0.15)), cy+int(I(0.25)), CW-int(I(0.30)), int(I(0.70)), n=15)
    
    # 稻田
    hill_flat(s, cx+int(I(0.05)), cy+int(CH*0.45), int(I(2.10)), int(I(0.35)), C["spring_green"])
    
    # 农民
    person_flat(s, cx+int(I(1.00)), cy+int(CH*0.38), C["wheat_gold"], scale=0.90)
    
    # 幼苗田
    for sx in [0.12, 0.32, 0.52, 0.72, 1.00, 1.25, 1.50, 1.75, 2.0]:
        seedling_simple(s, cx+int(I(sx)), cy+int(CH*0.62), scale=0.8)
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["spring_gnd"])
    add_cell_header(s, cx, cy, "谷雨", "Grain Rain", "Rain helps crops grow well.", comp)
    cell_border(s, 1, 0)
    
    # ── 立夏 (0,1) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 0, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    sun_flat(s, cx+CW-int(I(0.42)), cy+int(I(0.18)), int(I(0.22)), color=C["warm_yellow"])
    
    # 大树
    tree_round(s, cx+int(I(1.35)), cy+int(CH*0.28), scale=1.2, leaf_color=C["spring_green"])
    
    # 池塘
    add_oval(s, cx+int(I(0.08)), cy+int(CH*0.58), int(I(0.90)), int(I(0.28)), C["rain_blue"])
    lotus_simple(s, cx+int(I(0.25)), cy+int(CH*0.52), scale=0.8)
    lotus_simple(s, cx+int(I(0.65)), cy+int(CH*0.55), scale=0.7)
    
    # 人物观荷
    person_flat(s, cx+int(I(1.00)), cy+int(CH*0.45), C["soft_coral"], scale=0.90)
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["summer_gnd"])
    add_cell_header(s, cx, cy, "立夏", "Start of Summer", "Summer comes. Plants grow fast.", comp)
    cell_border(s, 0, 1)
    
    # ── 小满 (1,1) ─────────────────────────────────────────────────────────
    cx, cy = cell_bg_unified(s, 1, 1, C["summer_bg"], C["summer_sky"], C["summer_gnd"])
    
    sun_flat(s, cx+int(I(0.35)), cy+int(I(0.12)), int(I(0.24)))
    
    # 麦田
    hill_flat(s, cx+int(I(0.50)), cy+int(CH*0.38), int(I(1.30)), int(I(0.50)), C["wheat_gold"])
    
    # 水车
    water_wheel_simple(s, cx+CW-int(I(0.55)), cy+int(CH*0.42), scale=1.0)
    
    # 河流
    add_oval(s, cx+int(I(0.75)), cy+int(CH*0.68), int(I(1.20)), int(I(0.20)), C["rain_blue"])
    
    grass_simple(s, cx+BW, cy+int(CH*HORIZON_Y_RATIO), CW-BW*2, int(I(0.20)), C["summer_gnd"])
    add_cell_header(s, cx, cy, "小满", "Grain Full", "Grains start to become full.", comp)
    cell_border(s, 1, 1)

# ============================================================================
# 主函数
# ============================================================================
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    
    slide1(prs)
    slide2(prs)
    # TODO: slides 3-6 for remaining 16 solar terms
    
    prs.save(OUT)
    print(f"Saved: {OUT}")

if __name__ == "__main__":
    main()
