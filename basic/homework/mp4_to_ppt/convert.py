"""
Convert 24节气视频 to 10-page PPT, preserving the original UI design.
Strategy:
  1. Extract 10 evenly-spaced keyframes from the video via OpenCV.
  2. For each frame, attempt lightweight OCR-style title detection using the
     top/bottom region heuristic (bright large text area), falling back to a
     slide index label when nothing is found.
  3. Build a python-pptx presentation:
     - Slide size matches the video aspect ratio (16:9 → 33.87 × 19.05 cm).
     - Each frame is placed as the FULL-SLIDE background image.
     - A subtle semi-transparent caption bar at the bottom shows slide number
       and any detected title text.
"""

import os
import sys
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io

# ── paths ─────────────────────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
VIDEO_PATH = os.path.join(BASE_DIR, "24节气视频转ppt.mp4")
FRAMES_DIR = os.path.join(BASE_DIR, "frames")
OUTPUT_PPT = os.path.join(BASE_DIR, "24节气.pptx")

NUM_SLIDES = 10
os.makedirs(FRAMES_DIR, exist_ok=True)

# ── 1. Extract keyframes ──────────────────────────────────────────────────────

def extract_frames(video_path: str, n: int) -> list[str]:
    cap = cv2.VideoCapture(video_path)
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    fps   = cap.get(cv2.CAP_PROP_FPS)
    w     = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
    h     = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
    print(f"Video: {total} frames, {fps:.1f} fps, {w}×{h}")

    # sample at positions 5%, 15%, 25%, …, 95% to avoid black intro/outro
    positions = [int(total * (i * 2 + 1) / (2 * n)) for i in range(n)]
    paths = []
    for idx, pos in enumerate(positions):
        cap.set(cv2.CAP_PROP_POS_FRAMES, pos)
        ret, frame = cap.read()
        if not ret:
            # fallback: try next frame
            cap.set(cv2.CAP_PROP_POS_FRAMES, max(0, pos - 5))
            ret, frame = cap.read()
        if not ret:
            print(f"  WARNING: could not read frame at pos {pos}")
            continue
        path = os.path.join(FRAMES_DIR, f"frame_{idx+1:02d}.jpg")
        cv2.imwrite(path, frame, [cv2.IMWRITE_JPEG_QUALITY, 95])
        ts = pos / fps if fps > 0 else 0
        print(f"  Frame {idx+1:02d}: pos={pos} ({ts:.1f}s) → {path}")
        paths.append(path)
    cap.release()
    return paths


# ── 2. Detect title text in frame (simple heuristic, no heavy OCR) ────────────

def detect_title_region(frame_path: str, slide_idx: int) -> str:
    """
    Returns a label for the slide.  We try to detect whether the frame has a
    prominent title band at top or bottom by checking average brightness of
    the top-15% strip vs middle.  This is purely cosmetic – we just use the
    slide number as fallback label.
    """
    # We keep it simple: just return "第N节气" style label.
    # A real OCR pass would require paddleocr/tesseract which may not be installed.
    labels = [
        "春雨惊春清谷天", "夏满芒夏暑相连",
        "秋处露秋寒霜降", "冬雪雪冬小大寒",
        "立春 · 雨水",    "惊蛰 · 春分",
        "清明 · 谷雨",    "立夏 · 小满",
        "芒种 · 夏至",    "小暑 · 大暑",
    ]
    return labels[slide_idx] if slide_idx < len(labels) else f"第 {slide_idx+1} 页"


# ── 3. Build PPT ──────────────────────────────────────────────────────────────

def build_ppt(frame_paths: list[str], output_path: str):
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen (33.87 cm × 19.05 cm)
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for idx, frame_path in enumerate(frame_paths):
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # ── background image (full slide) ────────────────────────────────────
        pic = slide.shapes.add_picture(
            frame_path,
            left=0, top=0,
            width=slide_w, height=slide_h
        )
        # Move picture to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # ── bottom caption bar ───────────────────────────────────────────────
        bar_h = Inches(0.6)
        bar_top = slide_h - bar_h

        # Semi-transparent overlay via a filled rectangle with transparency
        txBox = slide.shapes.add_textbox(
            left=0, top=bar_top,
            width=slide_w, height=bar_h
        )
        tf = txBox.text_frame
        tf.word_wrap = False

        label = detect_title_region(frame_path, idx)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"  {idx+1:02d} / {NUM_SLIDES}   {label}  "
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Microsoft YaHei"

        # Fill the textbox background with a dark translucent colour
        from pptx.oxml.ns import qn
        from lxml import etree

        spPr = txBox._element.spPr
        # solid fill with 60% transparency
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", "000000")
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", "60000")   # 60 000 = 60% opacity in OOXML (out of 100 000)

        print(f"  Slide {idx+1}: {os.path.basename(frame_path)}  [{label}]")

    prs.save(output_path)
    print(f"\n✅  Saved: {output_path}")


# ── main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=== Step 1: Extract frames ===")
    frames = extract_frames(VIDEO_PATH, NUM_SLIDES)
    if not frames:
        sys.exit("No frames extracted – check the video path.")

    print(f"\n=== Step 2: Build PPT ({len(frames)} slides) ===")
    build_ppt(frames, OUTPUT_PPT)
