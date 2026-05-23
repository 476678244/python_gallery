"""
24 Solar Terms – Handmade-Book Style PPT  (v2)
================================================
Recreates the UI design of the hand-crafted children's book entirely with
python-pptx shapes / text-boxes / coloured rectangles – NO raw image pasting.

10 slides follow the story arc:
  1. Cover
  2. Contents + Origin (intro)
  3. The Origin (full)
  4. 立春 Start of Spring  /  雨水 Rain Water
  5. 惊蛰 Awakening of Insects  /  春分 Spring Equinox
  6. 清明 Clear & Bright  /  谷雨 Grain Rain
  7. 立夏 Start of Summer  /  小满 Grain Full
  8. 芒种 Grain in Ear  /  夏至 Summer Solstice
  9. 小暑 Minor Heat  /  大暑 Major Heat
 10. 立秋 Start of Autumn  /  处暑 End of Heat
  (bonus 10 = pass-down + influence epilogue, merged as slide 10 to keep 10 total)

Wait – video has:
  Cover → Contents/Origin → Origin full → Spring pair → ClearBright/GrainRain
  → Summer pair → GrainInEar/SummerSolstice → MinorHeat/MajorHeat
  → Autumn/EndHeat → ColdDew/FrostDescent → StartWinter/MinorSnow
  → MajorSnow/WinterSolstice → PassDown + Influence (epilogue)

We condense into 10 slides:
  1 Cover  2 TOC+Origin  3 立春+雨水  4 惊蛰+春分
  5 清明+谷雨  6 立夏+小满  7 芒种+夏至
  8 小暑+大暑  9 立秋→霜降 (autumn group)
 10 立冬→冬至+Epilogue
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "24节气_v2.pptx")

# ── Palette ──────────────────────────────────────────────────────────────────
C = dict(
    sky_blue   = RGBColor(0x4F, 0xB3, 0xE8),
    book_blue  = RGBColor(0x3A, 0x9A, 0xD4),
    lime_green = RGBColor(0xB8, 0xD8, 0x6A),
    pale_green = RGBColor(0xC8, 0xE6, 0x9A),
    warm_yellow= RGBColor(0xF5, 0xD3, 0x60),
    pale_yellow= RGBColor(0xF9, 0xEC, 0xA0),
    orange     = RGBColor(0xF0, 0x8C, 0x30),
    burnt_org  = RGBColor(0xD4, 0x6A, 0x1A),
    purple     = RGBColor(0x6A, 0x3A, 0x8C),
    lavender   = RGBColor(0xB0, 0x7A, 0xD0),
    deep_blue  = RGBColor(0x2A, 0x5C, 0x9A),
    powder_blue= RGBColor(0xAA, 0xCC, 0xE8),
    pink       = RGBColor(0xF0, 0xA0, 0xB0),
    rose       = RGBColor(0xE8, 0x60, 0x70),
    cream      = RGBColor(0xFA, 0xF5, 0xE4),
    off_white  = RGBColor(0xF5, 0xF0, 0xE0),
    ink        = RGBColor(0x1A, 0x1A, 0x2E),
    dark_brown = RGBColor(0x4A, 0x30, 0x10),
    mid_grey   = RGBColor(0xBB, 0xBB, 0xBB),
    grass_green= RGBColor(0x5A, 0xA8, 0x50),
    dark_green = RGBColor(0x2E, 0x7A, 0x3C),
    red        = RGBColor(0xCC, 0x33, 0x33),
    teal       = RGBColor(0x3A, 0xA8, 0x9A),
)

# ── Helpers ───────────────────────────────────────────────────────────────────
W = Inches(13.33)   # slide width
H = Inches(7.5)     # slide height


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill: RGBColor, alpha=None, line=None, line_w=None):
    """Add a coloured rectangle."""
    shp = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if alpha is not None:
        # set transparency via XML
        sp = shp._element
        spPr = sp.spPr
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgb = solidFill.find(qn("a:srgbClr"))
            if srgb is not None:
                a = etree.SubElement(srgb, qn("a:alpha"))
                a.set("val", str(int(alpha * 100000)))
    if line:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp


def ellipse(slide, x, y, w, h, fill: RGBColor, line=None, line_w=None):
    """Add an ellipse (oval)."""
    shp = slide.shapes.add_shape(9, x, y, w, h)   # MSO_CONNECTOR_TYPE 9 = oval
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=None,
            align=PP_ALIGN.LEFT, font="Microsoft YaHei", italic=False,
            wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = font
    if color:
        run.font.color.rgb = color
    return txb


def send_to_back(slide, shp):
    sp = shp._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    spTree.insert(2, sp)


# ── Decorative helpers ────────────────────────────────────────────────────────

def draw_tree_bare(slide, cx, cy, scale=1.0, color=None):
    """Simple bare-branch tree silhouette using thin rectangles."""
    c = color or C["dark_brown"]
    trunk_w = int(Inches(0.08) * scale)
    trunk_h = int(Inches(0.5)  * scale)
    # trunk
    rect(slide, cx - trunk_w//2, cy, trunk_w, trunk_h, c)
    # branches
    for dx, dy, bw, bh in [
        (-Inches(0.25)*scale, cy - Inches(0.35)*scale, Inches(0.35)*scale, Inches(0.06)*scale),
        ( Inches(0.05)*scale, cy - Inches(0.25)*scale, Inches(0.30)*scale, Inches(0.06)*scale),
        (-Inches(0.15)*scale, cy - Inches(0.5 )*scale, Inches(0.22)*scale, Inches(0.05)*scale),
    ]:
        rect(slide, int(cx + dx), int(dy), int(bw), int(bh), c)


def draw_cloud(slide, cx, cy, w, h, color=None):
    c = color or RGBColor(0xFF, 0xFF, 0xFF)
    ellipse(slide, cx, cy, w, int(h * 0.6), c)
    ellipse(slide, int(cx + w * 0.2), int(cy - h * 0.25), int(w * 0.45), int(h * 0.5), c)
    ellipse(slide, int(cx + w * 0.5), int(cy - h * 0.1), int(w * 0.4), int(h * 0.45), c)


def draw_sun(slide, cx, cy, r, color=None):
    c = color or C["warm_yellow"]
    ellipse(slide, int(cx - r), int(cy - r), r*2, r*2, c)


def draw_snowflake_dot(slide, cx, cy, r):
    ellipse(slide, int(cx-r), int(cy-r), r*2, r*2, RGBColor(0xFF,0xFF,0xFF))


def page_bg(slide, color: RGBColor):
    bg = rect(slide, 0, 0, W, H, color)
    send_to_back(slide, bg)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    """Sky-blue cover with yellow title block, lotus circle decoration."""
    s = blank(prs)
    page_bg(s, C["sky_blue"])

    # Wood-table feel – brown bottom strip
    rect(s, 0, int(H*0.82), W, int(H*0.18), C["dark_brown"])

    # Central book body
    book_x, book_y = int(W*0.20), int(H*0.08)
    book_w, book_h = int(W*0.60), int(H*0.72)
    rect(s, book_x, book_y, book_w, book_h, C["book_blue"],
         line=C["dark_brown"], line_w=Pt(2))

    # Yellow title band
    band_y = book_y + int(book_h * 0.06)
    band_h = int(book_h * 0.28)
    rect(s, book_x + int(book_w*0.05), band_y,
         int(book_w*0.90), band_h, C["warm_yellow"])

    textbox(s,
            book_x + int(book_w*0.05), band_y + int(band_h*0.02),
            int(book_w*0.90), int(band_h*0.48),
            "24 Solar Terms", size=38, bold=True, color=C["ink"],
            align=PP_ALIGN.CENTER, font="Georgia")
    textbox(s,
            book_x + int(book_w*0.05), band_y + int(band_h*0.48),
            int(book_w*0.90), int(band_h*0.50),
            "二十四节气", size=36, bold=True, color=C["ink"],
            align=PP_ALIGN.CENTER)

    # Author line
    textbox(s,
            book_x + int(book_w*0.10), band_y + band_h + int(book_h*0.04),
            int(book_w*0.80), int(book_h*0.08),
            "Written & Illustrated by Sunny and Mum",
            size=13, italic=True, color=RGBColor(0xFF,0xFF,0xFF),
            align=PP_ALIGN.CENTER, font="Georgia")

    # Yellow circle (sun / lotus wheel)
    circ_r = int(book_h * 0.22)
    circ_cx = book_x + book_w // 2
    circ_cy = book_y + int(book_h * 0.68)
    ellipse(s, circ_cx - circ_r, circ_cy - circ_r,
            circ_r*2, circ_r*2, C["warm_yellow"])
    # Lotus petals (simple ovals arranged radially)
    import math
    for ang in range(0, 360, 60):
        rad = math.radians(ang)
        px = int(circ_cx + (circ_r*0.6) * math.cos(rad))
        py = int(circ_cy + (circ_r*0.6) * math.sin(rad))
        pr = int(circ_r * 0.22)
        ellipse(s, px - pr, py - pr, pr*2, pr*2, C["rose"])

    # Tiny season icons in corners
    for cx2, cy2, c2 in [
        (int(book_x + book_w*0.08), int(book_y + book_h*0.50), C["grass_green"]),
        (int(book_x + book_w*0.85), int(book_y + book_h*0.50), C["orange"]),
    ]:
        ellipse(s, cx2-int(Inches(0.18)), cy2-int(Inches(0.18)),
                int(Inches(0.36)), int(Inches(0.36)), c2)


def slide_02_toc_origin(prs):
    """Light-green spread: Contents on left, Origin text on right."""
    s = blank(prs)
    page_bg(s, C["pale_green"])

    # ── Left half: Contents ──────────────────────────────────────────────────
    rect(s, int(W*0.02), int(H*0.06), int(W*0.46), int(H*0.88),
         RGBColor(0xD8, 0xED, 0xB0), line=C["grass_green"], line_w=Pt(1.5))

    textbox(s, int(W*0.05), int(H*0.10), int(W*0.40), int(H*0.10),
            "Contents", size=28, bold=True, color=C["dark_green"],
            font="Georgia")

    items = [
        ("1.", "Origin 起源"),
        ("2.", "Characteristics 特征"),
        ("3.", "Pass Down 传承"),
        ("4.", "Influence 影响"),
    ]
    for i, (num, text) in enumerate(items):
        y = int(H*(0.24 + i*0.14))
        textbox(s, int(W*0.06), y, int(W*0.06), int(H*0.10),
                num, size=20, bold=True, color=C["rose"])
        textbox(s, int(W*0.12), y, int(W*0.34), int(H*0.10),
                text, size=20, color=C["ink"])

    # Lotus decoration bottom-left
    ellipse(s, int(W*0.06), int(H*0.72), int(W*0.12), int(H*0.18), C["lime_green"])
    ellipse(s, int(W*0.10), int(H*0.64), int(W*0.10), int(H*0.14), C["rose"])

    # ── Right half: Origin intro text ─────────────────────────────────────────
    rect(s, int(W*0.52), int(H*0.06), int(W*0.46), int(H*0.88),
         RGBColor(0xD8, 0xED, 0xB0), line=C["grass_green"], line_w=Pt(1.5))

    textbox(s, int(W*0.54), int(H*0.09), int(W*0.42), int(H*0.08),
            "The Origin  起源", size=22, bold=True, color=C["dark_green"],
            font="Georgia")

    body = (
        "The Twenty-Four Solar Terms were invented by ancient Chinese people.\n\n"
        "Long long ago, people lived by farming. They watched the sun, the weather "
        "and plant changes, and summed up the Twenty-Four Solar Terms little by little.\n\n"
        "Solar Terms tell us when to plant seeds and when to harvest crops, making "
        "farming easier. Now, we still use Solar Terms to know the seasons and weather."
    )
    textbox(s, int(W*0.54), int(H*0.19), int(W*0.42), int(H*0.60),
            body, size=14, color=C["ink"], font="Georgia", wrap=True)

    # Small snowman + girl deco
    # snowman body
    ellipse(s, int(W*0.76), int(H*0.72), int(W*0.06), int(H*0.10), RGBColor(0xEE,0xEE,0xFF))
    ellipse(s, int(W*0.77), int(H*0.65), int(W*0.045), int(H*0.08), RGBColor(0xEE,0xEE,0xFF))
    # hat
    rect(s, int(W*0.775), int(H*0.62), int(W*0.035), int(H*0.04), C["deep_blue"])
    # dog
    ellipse(s, int(W*0.86), int(H*0.74), int(W*0.07), int(H*0.12), C["warm_yellow"])
    ellipse(s, int(W*0.875), int(H*0.70), int(W*0.045), int(H*0.07), C["warm_yellow"])


def slide_03_spring_lichun_yushui(prs):
    """Warm yellow spread: 立春 Start of Spring / 雨水 Rain Water."""
    s = blank(prs)
    page_bg(s, C["warm_yellow"])

    # Autumn-orange tree trunks (bare trees flank both pages like the original)
    draw_tree_bare(s, int(W*0.05), int(H*0.20), scale=2.5, color=C["burnt_org"])
    draw_tree_bare(s, int(W*0.48), int(H*0.20), scale=2.5, color=C["burnt_org"])
    draw_tree_bare(s, int(W*0.95), int(H*0.20), scale=2.5, color=C["burnt_org"])

    # Ground strips
    rect(s, 0, int(H*0.75), W, int(H*0.08), C["lime_green"])
    rect(s, 0, int(H*0.83), W, int(H*0.17), C["pale_yellow"])

    # Left page content (立春)
    textbox(s, int(W*0.05), int(H*0.05), int(W*0.20), int(H*0.12),
            "立春", size=30, bold=True, color=C["ink"])
    textbox(s, int(W*0.05), int(H*0.17), int(W*0.32), int(H*0.08),
            "Start of Spring", size=18, bold=True, color=C["dark_brown"],
            font="Georgia")
    textbox(s, int(W*0.05), int(H*0.26), int(W*0.38), int(H*0.08),
            "Spring begins. It gets warmer.", size=13, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Ox body (rounded rectangle as body)
    ellipse(s, int(W*0.08), int(H*0.44), int(W*0.22), int(H*0.24), C["burnt_org"])
    # Ox head
    ellipse(s, int(W*0.25), int(H*0.40), int(W*0.09), int(H*0.12), C["burnt_org"])
    # Saddle decoration
    rect(s,  int(W*0.12), int(H*0.44), int(W*0.10), int(H*0.12), C["rose"])
    # Legs
    for lx in [0.10, 0.16, 0.22, 0.26]:
        rect(s, int(W*lx), int(H*0.68), int(W*0.02), int(H*0.08), C["dark_brown"])
    # Tail
    rect(s, int(W*0.08), int(H*0.46), int(W*0.02), int(H*0.08), C["burnt_org"])

    # Right page content (雨水)
    textbox(s, int(W*0.62), int(H*0.05), int(W*0.15), int(H*0.12),
            "雨水", size=30, bold=True, color=C["ink"])
    textbox(s, int(W*0.62), int(H*0.18), int(W*0.32), int(H*0.08),
            "Rain Water", size=18, bold=True, color=C["dark_brown"],
            font="Georgia")
    textbox(s, int(W*0.62), int(H*0.27), int(W*0.34), int(H*0.08),
            "More rain. Everything grows.", size=13, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Rain drops
    import math
    for i in range(12):
        rx = int(W*(0.55 + (i % 4)*0.08))
        ry = int(H*(0.38 + (i // 4)*0.12))
        ellipse(s, rx, ry, int(Inches(0.08)), int(Inches(0.13)), C["deep_blue"])

    # Umbrella
    umb_cx = int(W*0.78)
    umb_cy = int(H*0.55)
    ellipse(s, umb_cx - int(Inches(0.5)), umb_cy - int(Inches(0.3)),
            int(Inches(1.0)), int(Inches(0.5)), C["deep_blue"])
    rect(s, umb_cx - int(Inches(0.03)), umb_cy,
         int(Inches(0.06)), int(Inches(0.5)), C["dark_brown"])

    # Swirling wind / water lines
    for wy in [0.35, 0.45, 0.55]:
        rect(s, int(W*0.52), int(H*wy), int(W*0.08), int(Inches(0.04)),
             C["teal"], alpha=0.6)


def slide_04_jingzhe_chunfen(prs):
    """Light blue-green: 惊蛰 Awakening of Insects / 春分 Spring Equinox."""
    s = blank(prs)
    page_bg(s, RGBColor(0xD0, 0xE8, 0xD0))

    # Page divider
    rect(s, int(W*0.49), int(H*0.05), int(Inches(0.04)), int(H*0.90),
         C["grass_green"])

    # Ground
    rect(s, 0, int(H*0.78), W, int(H*0.22), C["lime_green"])
    ellipse(s, int(W*0.05), int(H*0.70), int(W*0.35), int(H*0.20), C["pale_green"])
    ellipse(s, int(W*0.52), int(H*0.70), int(W*0.40), int(H*0.20), C["pale_green"])

    # 惊蛰 – left page
    textbox(s, int(W*0.05), int(H*0.06), int(W*0.20), int(H*0.10),
            "惊蛰", size=30, bold=True, color=C["ink"])
    textbox(s, int(W*0.05), int(H*0.18), int(W*0.38), int(H*0.07),
            "Awakening of Insects", size=17, bold=True,
            color=C["dark_green"], font="Georgia")
    textbox(s, int(W*0.05), int(H*0.26), int(W*0.40), int(H*0.10),
            "Thunder wakes the sleeping insects.\nSpring warmth fills the earth.",
            size=13, italic=True, color=C["ink"], font="Georgia")

    # Lightning bolt shape (simple triangle + rect)
    bolt_x, bolt_y = int(W*0.32), int(H*0.38)
    rect(s, bolt_x, bolt_y, int(Inches(0.12)), int(Inches(0.38)), C["warm_yellow"])
    rect(s, bolt_x - int(Inches(0.06)), bolt_y + int(Inches(0.20)),
         int(Inches(0.18)), int(Inches(0.12)), C["warm_yellow"])

    # Bug (ladybug style)
    bug_cx, bug_cy = int(W*0.22), int(H*0.56)
    ellipse(s, bug_cx, bug_cy, int(Inches(0.28)), int(Inches(0.22)), C["red"])
    ellipse(s, bug_cx + int(Inches(0.11)), bug_cy - int(Inches(0.10)),
            int(Inches(0.14)), int(Inches(0.14)), C["ink"])
    for dx2, dy2 in [(0.06, 0.04), (0.06, 0.12), (0.16, 0.04), (0.16, 0.12)]:
        ellipse(s, bug_cx + int(Inches(dx2)), bug_cy + int(Inches(dy2)),
                int(Inches(0.05)), int(Inches(0.05)), C["ink"])

    # Seedling
    rect(s, int(W*0.10), int(H*0.58), int(Inches(0.05)), int(Inches(0.22)),
         C["grass_green"])
    ellipse(s, int(W*0.07), int(H*0.54), int(Inches(0.18)), int(Inches(0.10)),
            C["dark_green"])
    ellipse(s, int(W*0.10), int(H*0.50), int(Inches(0.14)), int(Inches(0.10)),
            C["grass_green"])

    # 春分 – right page
    textbox(s, int(W*0.54), int(H*0.06), int(W*0.20), int(H*0.10),
            "春分", size=30, bold=True, color=C["ink"])
    textbox(s, int(W*0.54), int(H*0.18), int(W*0.40), int(H*0.07),
            "Spring Equinox", size=17, bold=True, color=C["dark_green"],
            font="Georgia")
    textbox(s, int(W*0.54), int(H*0.26), int(W*0.42), int(H*0.10),
            "Day and night are equal length.\nFlowers bloom everywhere.",
            size=13, italic=True, color=C["ink"], font="Georgia")

    # Sun
    draw_sun(s, int(W*0.80), int(H*0.22), int(Inches(0.45)))

    # Flowers row
    for fx in [0.55, 0.63, 0.71, 0.80, 0.88]:
        stem_x = int(W*fx)
        rect(s, stem_x, int(H*0.62), int(Inches(0.04)), int(Inches(0.22)),
             C["grass_green"])
        ellipse(s, stem_x - int(Inches(0.08)), int(H*0.56),
                int(Inches(0.20)), int(Inches(0.16)), C["rose"])
        ellipse(s, stem_x - int(Inches(0.04)), int(H*0.57),
                int(Inches(0.12)), int(Inches(0.12)), C["warm_yellow"])

    # Butterfly
    fly_x, fly_y = int(W*0.73), int(H*0.40)
    ellipse(s, fly_x,              fly_y, int(Inches(0.22)), int(Inches(0.14)),
            RGBColor(0xFF,0xA0,0xD0))
    ellipse(s, fly_x + int(Inches(0.16)), fly_y, int(Inches(0.22)), int(Inches(0.14)),
            RGBColor(0xFF,0xA0,0xD0))
    ellipse(s, fly_x + int(Inches(0.07)), fly_y - int(Inches(0.05)),
            int(Inches(0.08)), int(Inches(0.20)), C["ink"])


def slide_05_qingming_guyu(prs):
    """Grey-blue: 清明 Clear & Bright / 谷雨 Grain Rain."""
    s = blank(prs)
    page_bg(s, RGBColor(0xCC, 0xD8, 0xE4))

    # Ground
    rect(s, 0, int(H*0.78), W, int(H*0.22), C["pale_green"])

    # Divider
    rect(s, int(W*0.49), int(H*0.05), int(Inches(0.04)), int(H*0.90), C["mid_grey"])

    # 清明 – left ──────────────────────────────────────────────────────────────
    textbox(s, int(W*0.03), int(H*0.08), int(W*0.12), int(H*0.22),
            "清\n明", size=28, bold=True, color=C["ink"])
    textbox(s, int(W*0.14), int(H*0.08), int(W*0.32), int(H*0.08),
            "Clear and Bright", size=18, bold=True, color=C["ink"], font="Georgia")
    textbox(s, int(W*0.14), int(H*0.17), int(W*0.32), int(H*0.08),
            "Sunny day for outings.", size=13, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Hill / mountain background
    ellipse(s, int(W*0.05), int(H*0.35), int(W*0.35), int(H*0.50), C["pink"])
    ellipse(s, int(W*0.02), int(H*0.42), int(W*0.22), int(H*0.42), C["powder_blue"])

    # House on hill
    rect(s, int(W*0.18), int(H*0.36), int(W*0.12), int(H*0.12), C["rose"])
    rect(s, int(W*0.16), int(H*0.30), int(W*0.16), int(H*0.08), C["burnt_org"])  # roof
    # windows
    rect(s, int(W*0.20), int(H*0.39), int(W*0.03), int(H*0.04), C["warm_yellow"])
    rect(s, int(W*0.25), int(H*0.39), int(W*0.03), int(H*0.04), C["warm_yellow"])

    # Tree on hill
    rect(s, int(W*0.12), int(H*0.38), int(Inches(0.07)), int(Inches(0.35)),
         C["dark_brown"])
    ellipse(s, int(W*0.08), int(H*0.28), int(W*0.10), int(H*0.14), C["grass_green"])

    # Girl flying kite
    # kite
    kite_x, kite_y = int(W*0.08), int(H*0.12)
    # diamond kite using rotated rect approximation (two triangles = two rects)
    ellipse(s, kite_x, kite_y, int(Inches(0.25)), int(Inches(0.32)),
            RGBColor(0xDD,0xCC,0x22))
    rect(s, kite_x + int(Inches(0.10)), kite_y + int(Inches(0.05)),
         int(Inches(0.06)), int(Inches(0.20)), RGBColor(0x44,0x88,0xFF))
    # kite string
    rect(s, kite_x + int(Inches(0.12)), kite_y + int(Inches(0.32)),
         int(Inches(0.02)), int(Inches(0.5)), C["dark_brown"])

    # Girl body
    girl_x, girl_y = int(W*0.16), int(H*0.55)
    ellipse(s, girl_x, girl_y - int(H*0.07), int(Inches(0.2)), int(Inches(0.2)),
            RGBColor(0xFF,0xCC,0x99))  # head
    rect(s, girl_x, girl_y, int(Inches(0.20)), int(Inches(0.28)),
         RGBColor(0xFF,0x88,0xAA))    # body
    rect(s, girl_x, girl_y + int(Inches(0.28)), int(Inches(0.08)), int(Inches(0.18)),
         RGBColor(0x55,0x44,0xFF))    # leg L
    rect(s, girl_x + int(Inches(0.12)), girl_y + int(Inches(0.28)),
         int(Inches(0.08)), int(Inches(0.18)), RGBColor(0x55,0x44,0xFF))

    # 谷雨 – right ────────────────────────────────────────────────────────────
    textbox(s, int(W*0.86), int(H*0.25), int(W*0.12), int(H*0.22),
            "谷\n雨", size=28, bold=True, color=C["ink"])
    textbox(s, int(W*0.52), int(H*0.08), int(W*0.32), int(H*0.08),
            "Grain Rain", size=18, bold=True, color=C["ink"], font="Georgia")
    textbox(s, int(W*0.52), int(H*0.17), int(W*0.32), int(H*0.08),
            "Rain helps crops grow well.", size=13, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Rain grid / crop field
    for row in range(5):
        for col in range(6):
            rect(s,
                 int(W*(0.54 + col*0.055)),
                 int(H*(0.30 + row*0.09)),
                 int(W*0.03), int(Inches(0.06)), C["grass_green"])

    # Farmer boy
    farmer_x, farmer_y = int(W*0.76), int(H*0.52)
    ellipse(s, farmer_x, farmer_y - int(H*0.07), int(Inches(0.22)), int(Inches(0.22)),
            RGBColor(0xFF,0xCC,0x99))
    rect(s, farmer_x - int(Inches(0.01)), farmer_y,
         int(Inches(0.24)), int(Inches(0.28)), RGBColor(0x55,0x99,0xFF))
    rect(s, farmer_x - int(Inches(0.01)), farmer_y + int(Inches(0.28)),
         int(Inches(0.08)), int(Inches(0.18)), RGBColor(0x22,0x44,0xCC))
    rect(s, farmer_x + int(Inches(0.14)), farmer_y + int(Inches(0.28)),
         int(Inches(0.08)), int(Inches(0.18)), RGBColor(0x22,0x44,0xCC))
    # rain drops above him
    for i in range(6):
        rx = farmer_x + int(Inches(i * 0.12 - 0.2))
        ry = farmer_y - int(H*0.12) + int(Inches(i % 2 * 0.08))
        ellipse(s, rx, ry, int(Inches(0.05)), int(Inches(0.10)), C["deep_blue"])


def slide_06_lixia_xiaoman(prs):
    """Green: 立夏 Start of Summer / 小满 Grain Full."""
    s = blank(prs)
    page_bg(s, C["pale_green"])

    # Sky top
    rect(s, 0, 0, W, int(H*0.30), RGBColor(0xA8, 0xD8, 0xF0))
    # Ground / water
    rect(s, 0, int(H*0.72), W, int(H*0.28), RGBColor(0x5A, 0xA8, 0x70))
    ellipse(s, int(W*0.05), int(H*0.60), int(W*0.30), int(H*0.30),
            RGBColor(0x60, 0xB8, 0xD8))  # pond

    # Divider
    rect(s, int(W*0.49), int(H*0.05), int(Inches(0.04)), int(H*0.90),
         C["grass_green"])

    # Draw_sun
    draw_sun(s, int(W*0.88), int(H*0.10), int(Inches(0.40)))

    # 立夏 – left
    textbox(s, int(W*0.03), int(H*0.05), int(W*0.08), int(H*0.22),
            "立\n夏", size=26, bold=True, color=C["ink"])
    textbox(s, int(W*0.12), int(H*0.05), int(W*0.34), int(H*0.07),
            "Start of Summer", size=17, bold=True, color=C["dark_green"],
            font="Georgia")
    textbox(s, int(W*0.12), int(H*0.14), int(W*0.34), int(H*0.08),
            "Summer comes. Plants grow fast.", size=12, italic=True,
            color=C["ink"], font="Georgia")

    # Tree
    rect(s, int(W*0.30), int(H*0.28), int(Inches(0.10)), int(Inches(0.60)),
         C["dark_brown"])
    ellipse(s, int(W*0.22), int(H*0.12), int(W*0.18), int(H*0.22), C["grass_green"])
    ellipse(s, int(W*0.26), int(H*0.08), int(W*0.12), int(H*0.16), C["dark_green"])

    # Boy with butterfly net
    boy_x, boy_y = int(W*0.20), int(H*0.52)
    ellipse(s, boy_x, boy_y - int(H*0.08), int(Inches(0.22)), int(Inches(0.22)),
            RGBColor(0xFF,0xCC,0x88))
    rect(s, boy_x, boy_y, int(Inches(0.22)), int(Inches(0.26)),
         RGBColor(0x44,0xAA,0xFF))
    # net pole
    rect(s, boy_x + int(Inches(0.22)), boy_y - int(H*0.12),
         int(Inches(0.04)), int(Inches(0.55)), C["dark_brown"])
    ellipse(s, boy_x + int(Inches(0.18)), boy_y - int(H*0.20),
            int(Inches(0.22)), int(Inches(0.18)),
            RGBColor(0xCC, 0xEE, 0xFF))  # net

    # Lotus in pond
    ellipse(s, int(W*0.08), int(H*0.66), int(W*0.10), int(H*0.08),
            C["dark_green"])
    ellipse(s, int(W*0.13), int(H*0.62), int(W*0.07), int(H*0.06),
            C["rose"])

    # 小满 – right
    textbox(s, int(W*0.52), int(H*0.05), int(W*0.08), int(H*0.16),
            "小满", size=26, bold=True, color=C["ink"])
    textbox(s, int(W*0.62), int(H*0.05), int(W*0.34), int(H*0.07),
            "Grain Full", size=17, bold=True, color=C["dark_green"],
            font="Georgia")
    textbox(s, int(W*0.52), int(H*0.14), int(W*0.44), int(H*0.08),
            "Grains start to become full.", size=12, italic=True,
            color=C["ink"], font="Georgia")

    # Water wheel
    wheel_x, wheel_y = int(W*0.72), int(H*0.28)
    ellipse(s, wheel_x, wheel_y, int(Inches(0.7)), int(Inches(0.7)), C["orange"])
    ellipse(s, wheel_x + int(Inches(0.25)), wheel_y + int(Inches(0.25)),
            int(Inches(0.20)), int(Inches(0.20)),
            RGBColor(0xAA, 0x55, 0x00))
    # Spokes
    for ang2 in [0, 45, 90, 135]:
        import math
        r2 = math.radians(ang2)
        sx = int(wheel_x + Inches(0.35) + Inches(0.28) * math.cos(r2))
        sy = int(wheel_y + Inches(0.35) + Inches(0.28) * math.sin(r2))
        rect(s, min(sx, wheel_x + int(Inches(0.35))),
             min(sy, wheel_y + int(Inches(0.35))),
             abs(sx - wheel_x - int(Inches(0.35))) + int(Inches(0.03)),
             abs(sy - wheel_y - int(Inches(0.35))) + int(Inches(0.03)),
             C["dark_brown"])

    # Grain stalks
    for gx in [0.55, 0.62, 0.70, 0.82, 0.88]:
        stem_x2 = int(W*gx)
        rect(s, stem_x2, int(H*0.55), int(Inches(0.04)), int(Inches(0.28)),
             C["grass_green"])
        ellipse(s, stem_x2 - int(Inches(0.05)), int(H*0.46),
                int(Inches(0.14)), int(Inches(0.12)), C["warm_yellow"])


def slide_07_mangzhong_xiazhi(prs):
    """Yellow: 芒种 Grain in Ear / 夏至 Summer Solstice."""
    s = blank(prs)
    page_bg(s, C["pale_yellow"])

    # Sky
    rect(s, 0, 0, W, int(H*0.25), RGBColor(0xA0, 0xCC, 0xF0))
    # Ground
    rect(s, 0, int(H*0.72), W, int(H*0.28), C["warm_yellow"])

    # Divider
    rect(s, int(W*0.49), int(H*0.05), int(Inches(0.04)), int(H*0.90),
         C["burnt_org"])

    # 芒种 – left
    textbox(s, int(W*0.03), int(H*0.04), int(W*0.08), int(H*0.22),
            "芒\n种", size=26, bold=True, color=C["ink"])
    textbox(s, int(W*0.12), int(H*0.04), int(W*0.34), int(H*0.07),
            "Grain in Ear", size=17, bold=True, color=C["dark_brown"],
            font="Georgia")
    textbox(s, int(W*0.12), int(H*0.13), int(W*0.34), int(H*0.08),
            "Best time for planting and harvesting.",
            size=12, italic=True, color=C["ink"], font="Georgia")

    # Birds (V shapes as tiny rects)
    for bx, by in [(0.18, 0.06), (0.22, 0.04), (0.30, 0.05), (0.36, 0.07)]:
        rect(s, int(W*bx), int(H*by), int(Inches(0.12)), int(Inches(0.04)),
             C["ink"])

    # Wheat field (wave of red curves approximated by tall thin rects)
    for col in range(18):
        h_frac = 0.36 + 0.04 * ((col % 3) - 1)
        rect(s, int(W*(0.03 + col*0.025)), int(H*h_frac),
             int(W*0.022), int(H*0.30), C["orange"])

    # Egret / crane bird
    ellipse(s, int(W*0.12), int(H*0.60), int(W*0.06), int(H*0.14),
            RGBColor(0xF0, 0xF0, 0xF0))
    ellipse(s, int(W*0.13), int(H*0.56), int(W*0.04), int(H*0.06),
            RGBColor(0xF0, 0xF0, 0xF0))  # head
    rect(s, int(W*0.15), int(H*0.58), int(Inches(0.14)), int(Inches(0.04)),
         RGBColor(0xF5, 0xAA, 0x20))   # beak

    # 夏至 – right
    textbox(s, int(W*0.80), int(H*0.04), int(W*0.18), int(H*0.08),
            "Summer Solstice", size=14, bold=True, color=C["dark_brown"],
            font="Georgia")
    textbox(s, int(W*0.80), int(H*0.13), int(W*0.10), int(H*0.22),
            "夏\n至", size=26, bold=True, color=C["ink"])
    textbox(s, int(W*0.52), int(H*0.13), int(W*0.26), int(H*0.08),
            "Longest day of the year.",
            size=12, italic=True, color=C["ink"], font="Georgia")

    # Pond water (blue patch)
    ellipse(s, int(W*0.52), int(H*0.30), int(W*0.28), int(H*0.34),
            RGBColor(0x88, 0xCC, 0xEE))

    # Lotus leaves on water
    for lx2, ly2 in [(0.54, 0.36), (0.62, 0.28), (0.70, 0.38)]:
        ellipse(s, int(W*lx2), int(H*ly2),
                int(W*0.10), int(H*0.10), C["dark_green"])

    # Lotus flower
    ellipse(s, int(W*0.62), int(H*0.22), int(W*0.06), int(H*0.08), C["lavender"])
    for dx3, dy3 in [(-0.04, 0.02), (0.04, -0.02), (0, -0.04), (-0.02, -0.04)]:
        ellipse(s, int(W*(0.64+dx3)), int(H*(0.24+dy3)),
                int(W*0.04), int(H*0.06), C["pink"])

    # Frog at edge
    ellipse(s, int(W*0.78), int(H*0.58), int(W*0.05), int(H*0.06),
            C["grass_green"])
    ellipse(s, int(W*0.77), int(H*0.55), int(W*0.03), int(H*0.04),
            C["grass_green"])
    for ex in [0.76, 0.82]:
        ellipse(s, int(W*ex), int(H*0.54), int(Inches(0.08)), int(Inches(0.06)),
                C["grass_green"])


def slide_08_xiaoshu_dashu(prs):
    """Purple-blue: 小暑 Minor Heat / 大暑 Major Heat."""
    s = blank(prs)
    page_bg(s, C["purple"])

    # Sky gradient effect: blue top wave
    ellipse(s, 0, int(H*0.30), W, int(H*0.55), C["teal"])
    # Ground teal strip
    rect(s, 0, int(H*0.75), W, int(H*0.25), C["teal"])

    # Divider
    rect(s, int(W*0.49), int(H*0.05), int(Inches(0.04)), int(H*0.90),
         C["lavender"])

    # 小暑 – left
    textbox(s, int(W*0.03), int(H*0.26), int(W*0.12), int(H*0.08),
            "小暑", size=26, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    textbox(s, int(W*0.03), int(H*0.35), int(W*0.22), int(H*0.08),
            "Minor Heat", size=16, bold=True,
            color=RGBColor(0xFF,0xFF,0xFF), font="Georgia")
    textbox(s, int(W*0.03), int(H*0.44), int(W*0.43), int(H*0.08),
            "It becomes very hot.", size=12, italic=True,
            color=RGBColor(0xEE,0xEE,0xFF), font="Georgia")

    # Tree (bare with leaves)
    rect(s, int(W*0.06), int(H*0.30), int(Inches(0.10)), int(Inches(0.52)),
         C["burnt_org"])
    ellipse(s, int(W*0.02), int(H*0.16), int(W*0.14), int(H*0.18), C["dark_green"])
    ellipse(s, int(W*0.06), int(H*0.10), int(W*0.10), int(H*0.14), C["grass_green"])
    # Rabbit under tree
    ellipse(s, int(W*0.04), int(H*0.72), int(W*0.06), int(H*0.08),
            RGBColor(0xF0, 0xF0, 0xF0))
    ellipse(s, int(W*0.05), int(H*0.68), int(W*0.04), int(H*0.06),
            RGBColor(0xF0, 0xF0, 0xF0))

    # Girl holding watermelon
    girl2_x, girl2_y = int(W*0.25), int(H*0.46)
    ellipse(s, girl2_x, girl2_y - int(H*0.08), int(Inches(0.22)), int(Inches(0.22)),
            RGBColor(0xFF,0xCC,0x99))
    # Hair
    ellipse(s, girl2_x - int(Inches(0.02)), girl2_y - int(H*0.09),
            int(Inches(0.26)), int(Inches(0.12)), C["ink"])
    # Dress
    rect(s, girl2_x - int(Inches(0.01)), girl2_y,
         int(Inches(0.24)), int(Inches(0.30)), C["rose"])
    # Watermelon
    ellipse(s, girl2_x + int(Inches(0.14)), girl2_y + int(Inches(0.05)),
            int(Inches(0.26)), int(Inches(0.20)), C["grass_green"])
    ellipse(s, girl2_x + int(Inches(0.18)), girl2_y + int(Inches(0.08)),
            int(Inches(0.14)), int(Inches(0.12)), C["red"])

    # 大暑 – right
    textbox(s, int(W*0.86), int(H*0.26), int(W*0.12), int(H*0.22),
            "大\n暑", size=26, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    textbox(s, int(W*0.52), int(H*0.05), int(W*0.28), int(H*0.07),
            "Major Heat", size=16, bold=True,
            color=RGBColor(0xFF,0xFF,0xFF), font="Georgia")
    textbox(s, int(W*0.52), int(H*0.14), int(W*0.32), int(H*0.08),
            "The hottest time of year.", size=12, italic=True,
            color=RGBColor(0xEE,0xEE,0xFF), font="Georgia")

    # Tent / canopy
    rect(s, int(W*0.64), int(H*0.08), int(W*0.24), int(H*0.22),
         C["warm_yellow"])
    rect(s, int(W*0.62), int(H*0.06), int(W*0.28), int(H*0.04),
         C["deep_blue"])   # awning bar
    # Sun behind tent
    draw_sun(s, int(W*0.90), int(H*0.05), int(Inches(0.35)))

    # Watermelons on ground
    for wx2 in [0.53, 0.62, 0.73]:
        ellipse(s, int(W*wx2), int(H*0.76), int(W*0.06), int(H*0.08),
                C["grass_green"])
        ellipse(s, int(W*(wx2+0.01)), int(H*0.77), int(W*0.04), int(H*0.05),
                C["red"])


def slide_09_autumn_group(prs):
    """Orange-warm: 立秋+处暑+寒露+霜降 autumn group."""
    s = blank(prs)
    page_bg(s, C["warm_yellow"])

    # Sky strip top
    rect(s, 0, 0, W, int(H*0.20), RGBColor(0xFF, 0xDD, 0x88))
    # Ground
    rect(s, 0, int(H*0.78), W, int(H*0.22), RGBColor(0xFF, 0xCC, 0x66))
    # Snowy hill at bottom centre
    ellipse(s, int(W*0.30), int(H*0.68), int(W*0.40), int(H*0.20),
            RGBColor(0xFF, 0xF5, 0xEE))

    # Sun small
    draw_sun(s, int(W*0.50), int(H*0.10), int(Inches(0.28)))

    # Section labels ─ 4 mini sections ───────────────────────────────────────
    # Top-left: 立秋
    textbox(s, int(W*0.03), int(H*0.04), int(W*0.20), int(H*0.07),
            "立秋  Start of Autumn", size=16, bold=True, color=C["ink"])
    textbox(s, int(W*0.03), int(H*0.12), int(W*0.35), int(H*0.07),
            "Autumn begins. It cools down.", size=12, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Top-right: 处暑
    textbox(s, int(W*0.58), int(H*0.04), int(W*0.20), int(H*0.07),
            "处暑  End of Heat", size=16, bold=True, color=C["ink"])
    textbox(s, int(W*0.58), int(H*0.12), int(W*0.40), int(H*0.07),
            "Hot days are over.", size=12, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Bottom-left: 寒露
    textbox(s, int(W*0.03), int(H*0.56), int(W*0.20), int(H*0.07),
            "寒露  Cold Dew", size=15, bold=True, color=C["ink"])
    textbox(s, int(W*0.03), int(H*0.64), int(W*0.35), int(H*0.07),
            "Dew turns cold.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Bottom-right: 霜降
    textbox(s, int(W*0.58), int(H*0.56), int(W*0.20), int(H*0.07),
            "霜降  Frost's Descent", size=15, bold=True, color=C["ink"])
    textbox(s, int(W*0.58), int(H*0.64), int(W*0.40), int(H*0.07),
            "Frost comes. It gets cold.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Dividers
    rect(s, int(W*0.49), int(H*0.02), int(Inches(0.04)), int(H*0.96), C["burnt_org"])
    rect(s, 0, int(H*0.50), W, int(Inches(0.04)), C["burnt_org"])

    # Autumn trees (orange leaf)
    for tx in [0.08, 0.20]:
        rect(s, int(W*tx), int(H*0.30), int(Inches(0.08)), int(Inches(0.36)),
             C["dark_brown"])
        ellipse(s, int(W*(tx-0.04)), int(H*0.18), int(W*0.14), int(H*0.16),
                C["burnt_org"])
        ellipse(s, int(W*(tx-0.01)), int(H*0.14), int(W*0.10), int(H*0.12),
                C["warm_yellow"])

    # Painting girl (left panel)
    pg_x, pg_y = int(W*0.12), int(H*0.27)
    ellipse(s, pg_x, pg_y - int(H*0.08), int(Inches(0.22)), int(Inches(0.22)),
            RGBColor(0xFF,0xCC,0x99))
    rect(s, pg_x, pg_y, int(Inches(0.22)), int(Inches(0.28)), C["lavender"])
    # Easel
    rect(s, pg_x + int(Inches(0.26)), pg_y - int(H*0.06),
         int(Inches(0.16)), int(Inches(0.22)), RGBColor(0xFF, 0xFF, 0xFF))
    rect(s, pg_x + int(Inches(0.24)), pg_y + int(Inches(0.16)),
         int(Inches(0.05)), int(Inches(0.20)), C["dark_brown"])
    rect(s, pg_x + int(Inches(0.38)), pg_y + int(Inches(0.16)),
         int(Inches(0.05)), int(Inches(0.20)), C["dark_brown"])

    # Farmer with basket (right panel top)
    f2_x, f2_y = int(W*0.62), int(H*0.22)
    ellipse(s, f2_x, f2_y - int(H*0.07), int(Inches(0.22)), int(Inches(0.22)),
            RGBColor(0xFF,0xCC,0x99))
    rect(s, f2_x, f2_y, int(Inches(0.22)), int(Inches(0.28)), C["burnt_org"])
    # Hat
    ellipse(s, f2_x - int(Inches(0.04)), f2_y - int(H*0.10),
            int(Inches(0.30)), int(Inches(0.08)), C["warm_yellow"])
    rect(s, f2_x + int(Inches(0.04)), f2_y - int(H*0.16),
         int(Inches(0.14)), int(Inches(0.08)), C["warm_yellow"])
    # Basket
    ellipse(s, f2_x + int(Inches(0.24)), f2_y + int(Inches(0.06)),
            int(Inches(0.24)), int(Inches(0.20)), C["burnt_org"])
    # Harvest items in basket
    ellipse(s, f2_x + int(Inches(0.26)), f2_y + int(Inches(0.04)),
            int(Inches(0.10)), int(Inches(0.10)), C["orange"])
    ellipse(s, f2_x + int(Inches(0.32)), f2_y + int(Inches(0.04)),
            int(Inches(0.10)), int(Inches(0.10)), C["red"])

    # Scarecrow (right bottom)
    sc_x, sc_y = int(W*0.80), int(H*0.60)
    ellipse(s, sc_x, sc_y - int(H*0.07), int(Inches(0.22)), int(Inches(0.22)),
            C["warm_yellow"])  # head
    rect(s, sc_x + int(Inches(0.03)), sc_y, int(Inches(0.16)), int(Inches(0.24)),
         C["burnt_org"])
    rect(s, sc_x - int(Inches(0.16)), sc_y + int(Inches(0.04)),
         int(Inches(0.50)), int(Inches(0.06)), C["burnt_org"])  # arms

    # Cotton plant (left bottom)
    for stem_y in [0.58, 0.66]:
        rect(s, int(W*0.08), int(H*stem_y), int(Inches(0.04)), int(Inches(0.14)),
             C["grass_green"])
        ellipse(s, int(W*0.05), int(H*(stem_y-0.04)), int(Inches(0.14)),
                int(Inches(0.10)), RGBColor(0xFF,0xFF,0xFF))

    # Apple tree (right bottom)
    rect(s, int(W*0.70), int(H*0.56), int(Inches(0.08)), int(Inches(0.30)),
         C["dark_brown"])
    ellipse(s, int(W*0.65), int(H*0.44), int(W*0.14), int(H*0.16), C["lavender"])
    for ax, ay in [(0.67, 0.43), (0.72, 0.41), (0.76, 0.44)]:
        ellipse(s, int(W*ax), int(H*ay), int(Inches(0.10)), int(Inches(0.11)), C["red"])


def slide_10_winter_epilogue(prs):
    """Blue: 立冬→冬至 + Epilogue pass-down / influence."""
    s = blank(prs)
    page_bg(s, C["powder_blue"])

    # Split into 3 panels: left=立冬小雪, mid=大雪冬至, right=Epilogue
    # Left panel bg
    rect(s, int(W*0.01), int(H*0.04), int(W*0.32), int(H*0.92),
         RGBColor(0xCC, 0xDD, 0xF5), line=C["deep_blue"], line_w=Pt(1))
    # Mid panel bg
    rect(s, int(W*0.35), int(H*0.04), int(W*0.32), int(H*0.92),
         RGBColor(0xDD, 0xEE, 0xFF), line=C["deep_blue"], line_w=Pt(1))
    # Right panel bg (epilogue) – cream
    rect(s, int(W*0.69), int(H*0.04), int(W*0.30), int(H*0.92),
         C["cream"], line=C["burnt_org"], line_w=Pt(1))

    # ── Left: 立冬 + 小雪 ──────────────────────────────────────────────────────
    textbox(s, int(W*0.03), int(H*0.06), int(W*0.28), int(H*0.07),
            "立冬  Start of Winter", size=14, bold=True, color=C["ink"])
    textbox(s, int(W*0.03), int(H*0.14), int(W*0.28), int(H*0.06),
            "Winter comes.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Family portrait (3 figures)
    for fi, (fx3, fy3, fh, fc) in enumerate([
        (0.03, 0.30, 0.26, RGBColor(0xAA,0xCC,0xFF)),
        (0.09, 0.28, 0.28, RGBColor(0xFF,0xAA,0xCC)),
        (0.16, 0.32, 0.22, RGBColor(0xAA,0xCC,0xFF)),
    ]):
        bx3 = int(W*fx3)
        by3 = int(H*fy3)
        body_h = int(H*fh)
        ellipse(s, bx3, by3 - int(H*0.08), int(Inches(0.22)), int(Inches(0.22)),
                RGBColor(0xFF,0xCC,0x99))
        rect(s, bx3, by3, int(Inches(0.22)), body_h, fc)

    # Food on table
    rect(s, int(W*0.04), int(H*0.62), int(W*0.22), int(Inches(0.06)),
         C["dark_brown"])  # table top
    for fx4, fc4 in [(0.05, C["warm_yellow"]), (0.10, C["rose"]),
                     (0.15, C["orange"]), (0.19, C["dark_green"])]:
        ellipse(s, int(W*fx4), int(H*0.56), int(W*0.04), int(W*0.04), fc4)

    textbox(s, int(W*0.03), int(H*0.73), int(W*0.28), int(H*0.07),
            "小雪  Minor Snow", size=14, bold=True, color=C["ink"])
    textbox(s, int(W*0.03), int(H*0.81), int(W*0.28), int(H*0.06),
            "Light snow falls.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")
    # Snow trees (right of left panel)
    rect(s, int(W*0.25), int(H*0.40), int(Inches(0.07)), int(Inches(0.40)),
         C["dark_brown"])
    ellipse(s, int(W*0.22), int(H*0.26), int(W*0.10), int(H*0.18), C["grass_green"])
    ellipse(s, int(W*0.23), int(H*0.22), int(W*0.08), int(H*0.10),
            RGBColor(0xFF,0xFF,0xFF))  # snow on tree
    # snowflakes
    for snx, sny in [(0.10,0.22),(0.16,0.18),(0.24,0.29),(0.08,0.30)]:
        draw_snowflake_dot(s, int(W*snx), int(H*sny), int(Inches(0.05)))

    # ── Mid: 大雪 + 冬至 ────────────────────────────────────────────────────────
    textbox(s, int(W*0.37), int(H*0.06), int(W*0.28), int(H*0.07),
            "大雪  Major Snow", size=14, bold=True, color=C["ink"])
    textbox(s, int(W*0.37), int(H*0.14), int(W*0.28), int(H*0.06),
            "Heavy snow covers the ground.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")

    # Snowy house
    rect(s, int(W*0.38), int(H*0.38), int(W*0.16), int(H*0.18),
         RGBColor(0xAA, 0xBB, 0xDD))
    rect(s, int(W*0.36), int(H*0.30), int(W*0.20), int(H*0.10),
         RGBColor(0x88, 0x99, 0xCC))  # roof
    ellipse(s, int(W*0.35), int(H*0.26), int(W*0.22), int(H*0.10),
            RGBColor(0xFF, 0xFF, 0xFF))  # snow on roof
    # Windows
    rect(s, int(W*0.40), int(H*0.41), int(W*0.04), int(H*0.05),
         C["warm_yellow"])
    rect(s, int(W*0.47), int(H*0.41), int(W*0.04), int(H*0.05),
         C["warm_yellow"])
    # Snowman
    ellipse(s, int(W*0.50), int(H*0.40), int(W*0.06), int(H*0.10),
            RGBColor(0xEE,0xEE,0xFF))
    ellipse(s, int(W*0.51), int(H*0.34), int(W*0.045), int(H*0.08),
            RGBColor(0xEE,0xEE,0xFF))
    rect(s, int(W*0.515), int(H*0.31), int(W*0.025), int(H*0.04), C["deep_blue"])
    # More snowflakes
    for snx2, sny2 in [(0.40,0.22),(0.44,0.18),(0.52,0.24),(0.57,0.20),(0.60,0.28)]:
        draw_snowflake_dot(s, int(W*snx2), int(H*sny2), int(Inches(0.05)))

    textbox(s, int(W*0.37), int(H*0.60), int(W*0.28), int(H*0.07),
            "冬至  Winter Solstice", size=14, bold=True, color=C["ink"])
    textbox(s, int(W*0.37), int(H*0.68), int(W*0.28), int(H*0.06),
            "Longest night of the year.", size=11, italic=True,
            color=C["dark_brown"], font="Georgia")
    # Dining scene
    rect(s, int(W*0.38), int(H*0.76), int(W*0.18), int(Inches(0.04)),
         C["warm_yellow"])  # table
    # Dumplings
    for dx3 in [0.39, 0.43, 0.47, 0.51]:
        ellipse(s, int(W*dx3), int(H*0.73), int(W*0.03), int(Inches(0.06)),
                RGBColor(0xFF,0xFF,0xDD))
    # Girl at table
    g3_x, g3_y = int(W*0.57), int(H*0.64)
    ellipse(s, g3_x, g3_y - int(H*0.06), int(Inches(0.18)), int(Inches(0.18)),
            RGBColor(0xFF,0xCC,0x99))
    rect(s, g3_x, g3_y, int(Inches(0.18)), int(Inches(0.22)),
         RGBColor(0xFF,0x88,0xAA))

    # ── Right: Epilogue ───────────────────────────────────────────────────────
    textbox(s, int(W*0.71), int(H*0.06), int(W*0.26), int(H*0.07),
            "二十四节气的传承", size=14, bold=True, color=C["ink"])
    body_left = (
        "The 24 Solar Terms are the wisdom of ancient Chinese people. "
        "They have been passed down from generation to generation.\n\n"
        "It is China's World Cultural Heritage that we should inherit and carry forward."
    )
    textbox(s, int(W*0.71), int(H*0.15), int(W*0.26), int(H*0.32),
            body_left, size=11, color=C["ink"], font="Georgia", wrap=True)

    textbox(s, int(W*0.71), int(H*0.50), int(W*0.26), int(H*0.07),
            "二十四节气的影响", size=14, bold=True, color=C["ink"])
    body_right = (
        "The 24 Solar Terms guide farming and show the changes of seasons. "
        "They influence our food, clothes and daily life, and are closely connected with nature."
    )
    textbox(s, int(W*0.71), int(H*0.58), int(W*0.26), int(H*0.30),
            body_right, size=11, color=C["ink"], font="Georgia", wrap=True)

    # Small butterfly + dog decorations
    ellipse(s, int(W*0.83), int(H*0.46), int(Inches(0.18)), int(Inches(0.12)),
            C["lavender"])
    ellipse(s, int(W*0.87), int(H*0.46), int(Inches(0.18)), int(Inches(0.12)),
            C["lavender"])
    ellipse(s, int(W*0.85), int(H*0.44), int(Inches(0.06)), int(Inches(0.16)),
            C["ink"])

    ellipse(s, int(W*0.88), int(H*0.86), int(Inches(0.20)), int(Inches(0.14)),
            C["warm_yellow"])
    ellipse(s, int(W*0.89), int(H*0.82), int(Inches(0.14)), int(Inches(0.12)),
            C["warm_yellow"])


# ── Main ─────────────────────────────────────────────────────────────────────

def build():
    prs = prs_new()
    print("Building slides…")
    builders = [
        ("01 Cover",              slide_01_cover),
        ("02 Contents + Origin",  slide_02_toc_origin),
        ("03 立春 + 雨水",         slide_03_spring_lichun_yushui),
        ("04 惊蛰 + 春分",         slide_04_jingzhe_chunfen),
        ("05 清明 + 谷雨",         slide_05_qingming_guyu),
        ("06 立夏 + 小满",         slide_06_lixia_xiaoman),
        ("07 芒种 + 夏至",         slide_07_mangzhong_xiazhi),
        ("08 小暑 + 大暑",         slide_08_xiaoshu_dashu),
        ("09 立秋 → 霜降",         slide_09_autumn_group),
        ("10 立冬 → 冬至 + Epilogue", slide_10_winter_epilogue),
    ]
    for name, fn in builders:
        print(f"  {name}")
        fn(prs)
    prs.save(OUT)
    print(f"\n✅  Saved: {OUT}")


if __name__ == "__main__":
    build()
