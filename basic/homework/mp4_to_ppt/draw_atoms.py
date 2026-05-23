"""Shared drawing atoms for the 24 Solar Terms PPT."""
import math, os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(r,g,b): return RGBColor(r,g,b)
def I(v): return Inches(v)

C = dict(
    ink=rgb(0x1A,0x1A,0x2E), white=rgb(0xFF,0xFF,0xFF), cream=rgb(0xFA,0xF5,0xE4),
    off_white=rgb(0xF5,0xF0,0xE0), skin=rgb(0xFF,0xCC,0x99), hair=rgb(0x22,0x11,0x00),
    brown=rgb(0x8B,0x45,0x13), dark_brown=rgb(0x4A,0x30,0x10),
    ox_body=rgb(0xC8,0x80,0x40), red=rgb(0xCC,0x33,0x33), rose=rgb(0xE8,0x60,0x70),
    pink=rgb(0xF4,0xA0,0xB8), orange=rgb(0xF0,0x8C,0x30), warm_yellow=rgb(0xF5,0xD3,0x60),
    pale_yellow=rgb(0xF9,0xEC,0xA0), gold=rgb(0xDD,0xAA,0x00), grass=rgb(0x5A,0xA8,0x50),
    dark_green=rgb(0x2E,0x7A,0x3C), deep_blue=rgb(0x2A,0x5C,0x9A),
    sky_blue=rgb(0x4A,0xA8,0xD8), teal=rgb(0x3A,0xA8,0x9A), purple=rgb(0x6A,0x3A,0x8C),
    lavender=rgb(0xB0,0x7A,0xD0), grey=rgb(0xBB,0xBB,0xBB), wheat=rgb(0xF0,0xD0,0x70),
    burnt_org=rgb(0xC8,0x60,0x18), spring_bg=rgb(0xC8,0xE8,0xA8),
    spring_sky=rgb(0xA8,0xD4,0xF4), spring_gnd=rgb(0x88,0xC0,0x60),
    summer_bg=rgb(0xB8,0xE0,0x88), summer_sky=rgb(0x78,0xC4,0xF0),
    summer_gnd=rgb(0x50,0xA0,0x50), pond=rgb(0x60,0xB4,0xD8),
    autumn_bg=rgb(0xF4,0xD0,0x70), autumn_sky=rgb(0xFF,0xE8,0x80),
    autumn_gnd=rgb(0xD4,0x88,0x20), winter_bg=rgb(0xCC,0xDD,0xF0),
    winter_sky=rgb(0xAA,0xCC,0xE8), winter_gnd=rgb(0xE8,0xF4,0xFF),
)

def add_rect(slide, x, y, w, h, fill, lc=None, lp=0):
    s = slide.shapes.add_shape(1,int(x),int(y),int(w),int(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb=lc; s.line.width=Pt(lp)
    else: s.line.fill.background()
    return s

def add_oval(slide, x, y, w, h, fill, lc=None, lp=0):
    s = slide.shapes.add_shape(9,int(x),int(y),int(w),int(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb=lc; s.line.width=Pt(lp)
    else: s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, size=14, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, font="Microsoft YaHei", wrap=True):
    tb = slide.shapes.add_textbox(int(x),int(y),int(w),int(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.name=font
    if color: r.font.color.rgb = color
    return tb

def push_back(slide, shp):
    sp = shp._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)

# ── atoms ────────────────────────────────────────────────────────────────────
def sun(slide, cx, cy, r, color=None):
    c = color or C["warm_yellow"]
    add_oval(slide, cx-r, cy-r, r*2, r*2, c)
    for a in range(0,360,45):
        rad=math.radians(a)
        add_oval(slide, cx+int((r+I(0.09))*math.cos(rad))-int(I(0.05)),
                 cy+int((r+I(0.09))*math.sin(rad))-int(I(0.05)),
                 int(I(0.10)), int(I(0.10)), c)

def grass_strip(slide, x, y, w, h, color=None):
    c = color or C["spring_gnd"]
    add_rect(slide, x, y, w, h, c)
    step = int(I(0.16))
    for i in range(int(w)//step):
        gx = x+i*step+int(I(0.03))
        add_rect(slide, gx, y-int(I(0.09)), int(I(0.03)), int(I(0.11)), C["dark_green"])
        add_rect(slide, gx+int(I(0.05)), y-int(I(0.13)), int(I(0.03)), int(I(0.13)), C["dark_green"])

def tree_bare(slide, cx, cy, scale=1.0):
    tw=int(I(0.10*scale)); th=int(I(0.52*scale))
    add_rect(slide, cx-tw//2, cy, tw, th, C["dark_brown"])
    branches=[(-I(0.26),cy-I(0.20),I(0.28),I(0.06)),
              ( I(0.02),cy-I(0.30),I(0.26),I(0.06)),
              (-I(0.14),cy-I(0.42),I(0.20),I(0.05)),
              (-I(0.22),cy-I(0.52),I(0.14),I(0.04)),
              ( I(0.00),cy-I(0.46),I(0.16),I(0.04))]
    for dx,dy,bw,bh in branches:
        add_rect(slide, cx+int(dx*scale), int(dy*scale if scale!=1 else dy),
                 int(bw*scale), int(bh*scale), C["dark_brown"])

def tree_leafy(slide, cx, cy, scale=1.0, lc=None, tc=None):
    lc=lc or C["grass"]; tc=tc or C["dark_brown"]
    add_rect(slide, cx-int(I(0.05*scale)), cy+int(I(0.14*scale)),
             int(I(0.10*scale)), int(I(0.50*scale)), tc)
    add_oval(slide, cx-int(I(0.30*scale)), cy-int(I(0.06*scale)),
             int(I(0.60*scale)), int(I(0.40*scale)), lc)
    add_oval(slide, cx-int(I(0.22*scale)), cy-int(I(0.22*scale)),
             int(I(0.44*scale)), int(I(0.34*scale)), C["dark_green"])
    add_oval(slide, cx-int(I(0.14*scale)), cy-int(I(0.32*scale)),
             int(I(0.28*scale)), int(I(0.22*scale)), lc)

def tree_pine(slide, cx, cy, scale=1.0):
    add_rect(slide, cx-int(I(0.05*scale)), cy+int(I(0.16*scale)),
             int(I(0.10*scale)), int(I(0.30*scale)), C["dark_brown"])
    for dy2,sw2,sh2 in [(I(0.02),I(0.50),I(0.22)),(-I(0.14),I(0.38),I(0.20)),
                         (-I(0.28),I(0.26),I(0.18))]:
        add_oval(slide, cx-int(sw2*scale//2), cy+int(dy2*scale),
                 int(sw2*scale), int(sh2*scale), C["dark_green"])

def person(slide, bx, by, body_color, skirt=False, hat=False, hat_color=None,
           scale=1.0, pose="stand", season="spring"):
    """Draw a child figure.
    pose:   stand | wave | walk | sit | cheer | hold
    season: spring | summer | autumn | winter  (affects clothing thickness)
    """
    hac = hat_color or C["warm_yellow"]
    s = scale
    _i = lambda v: int(I(v * s))

    # ── seasonal clothing adjustments ──
    # body width/height multipliers for thick/thin clothing
    if season == "winter":
        bw_mul, bh_mul = 1.30, 1.15   # thick puffy coat
        arm_w_mul = 1.5                # thick sleeves
        leg_color = C["dark_brown"]
        show_scarf = True
        auto_hat = True
    elif season == "autumn":
        bw_mul, bh_mul = 1.12, 1.05   # light jacket
        arm_w_mul = 1.2
        leg_color = C["dark_brown"]
        show_scarf = False
        auto_hat = False
    elif season == "summer":
        bw_mul, bh_mul = 0.92, 0.80   # thin T-shirt / tank top
        arm_w_mul = 0.8               # bare arms
        leg_color = C["skin"]          # shorts → skin-color legs
        show_scarf = False
        auto_hat = False
    else:  # spring
        bw_mul, bh_mul = 1.0, 1.0
        arm_w_mul = 1.0
        leg_color = C["dark_brown"]
        show_scarf = False
        auto_hat = False

    # ── head ──
    hw = _i(0.24); hh = _i(0.24)
    head_y = by - hh - _i(0.02)
    add_oval(slide, bx, head_y, hw, hh, C["skin"])
    # hair
    add_oval(slide, bx - _i(0.02), head_y - _i(0.04),
             hw + _i(0.04), int(hh * 0.6), C["hair"])
    # eyes
    add_oval(slide, bx + _i(0.06), head_y + _i(0.08), _i(0.03), _i(0.03), C["ink"])
    add_oval(slide, bx + _i(0.15), head_y + _i(0.08), _i(0.03), _i(0.03), C["ink"])
    # smile
    add_rect(slide, bx + _i(0.10), head_y + _i(0.18), _i(0.04), _i(0.02), C["red"])
    # cheeks (blush)
    add_oval(slide, bx + _i(0.02), head_y + _i(0.13), _i(0.04), _i(0.03), C["pink"])
    add_oval(slide, bx + _i(0.18), head_y + _i(0.13), _i(0.04), _i(0.03), C["pink"])

    # ── hat (explicit or winter auto) ──
    if hat or auto_hat:
        if season == "winter":
            # woolly beanie
            add_oval(slide, bx - _i(0.04), head_y - _i(0.10),
                     hw + _i(0.08), _i(0.14), hac)
            # pom-pom
            add_oval(slide, bx + _i(0.06), head_y - _i(0.16), _i(0.10), _i(0.10), C["white"])
        else:
            add_oval(slide, bx - _i(0.06), head_y - _i(0.04),
                     hw + _i(0.12), _i(0.08), hac)
            add_oval(slide, bx + _i(0.02), head_y - _i(0.14),
                     _i(0.18), _i(0.14), hac)

    # ── scarf (winter) ──
    if show_scarf:
        add_rect(slide, bx - _i(0.02), by - _i(0.06), hw + _i(0.04), _i(0.08), C["red"])
        add_rect(slide, bx + _i(0.04), by + _i(0.02), _i(0.06), _i(0.14), C["red"])

    # ── body ──
    bw = int(_i(0.24) * bw_mul); bh = int(_i(0.30) * bh_mul)
    body_x = bx - int((bw - _i(0.24)) / 2)  # center wider body on head
    if skirt and season != "winter":
        add_rect(slide, body_x, by, bw, int(bh * 0.55), body_color)
        add_oval(slide, body_x - _i(0.04), by + int(bh * 0.40),
                 bw + _i(0.08), int(bh * 0.50), body_color)
    elif season == "winter":
        # puffy coat with rounded bottom
        add_rect(slide, body_x, by, bw, int(bh * 0.80), body_color)
        add_oval(slide, body_x - _i(0.02), by + int(bh * 0.60),
                 bw + _i(0.04), int(bh * 0.40), body_color)
        # coat buttons
        for bi in range(3):
            add_oval(slide, body_x + bw // 2 - _i(0.02), by + _i(0.06) + bi * _i(0.08),
                     _i(0.04), _i(0.04), C["white"])
    else:
        add_rect(slide, body_x, by, bw, bh, body_color)

    # ── summer detail: collar / neckline ──
    if season == "summer":
        add_oval(slide, bx + _i(0.04), by - _i(0.02), _i(0.16), _i(0.06), C["skin"])

    # ── arms (pose-dependent) ──
    aw = int(_i(0.05) * arm_w_mul); ah = _i(0.18)
    arm_color = C["skin"] if season == "summer" else body_color

    if pose == "wave":
        # left arm down, right arm up (waving)
        add_rect(slide, body_x - _i(0.04), by + _i(0.02), aw, ah, arm_color)
        add_rect(slide, body_x + bw, by - _i(0.12), aw, ah, arm_color)
        # hand at top of wave
        add_oval(slide, body_x + bw - _i(0.02), by - _i(0.18), _i(0.08), _i(0.08), C["skin"])
    elif pose == "cheer":
        # both arms raised up in V
        add_rect(slide, body_x - _i(0.08), by - _i(0.10), aw, ah, arm_color)
        add_rect(slide, body_x + bw + _i(0.02), by - _i(0.10), aw, ah, arm_color)
        add_oval(slide, body_x - _i(0.10), by - _i(0.16), _i(0.08), _i(0.08), C["skin"])
        add_oval(slide, body_x + bw + _i(0.02), by - _i(0.16), _i(0.08), _i(0.08), C["skin"])
    elif pose == "walk":
        # arms slightly angled outward
        add_rect(slide, body_x - _i(0.08), by + _i(0.04), aw, int(ah * 0.9), arm_color)
        add_rect(slide, body_x + bw + _i(0.02), by, aw, int(ah * 0.9), arm_color)
    elif pose == "hold":
        # arms forward (holding something)
        add_rect(slide, body_x + _i(0.02), by + _i(0.06), aw, int(ah * 0.7), arm_color)
        add_rect(slide, body_x + bw - _i(0.06), by + _i(0.06), aw, int(ah * 0.7), arm_color)
    elif pose == "sit":
        # arms resting on sides, shorter
        add_rect(slide, body_x - _i(0.04), by + _i(0.04), aw, int(ah * 0.65), arm_color)
        add_rect(slide, body_x + bw, by + _i(0.04), aw, int(ah * 0.65), arm_color)
    else:  # stand (default, but slightly angled out for liveliness)
        add_rect(slide, body_x - _i(0.05), by + _i(0.02), aw, ah, arm_color)
        add_rect(slide, body_x + bw, by + _i(0.02), aw, ah, arm_color)

    # ── legs (pose-dependent) ──
    leg_h = _i(0.22)
    shoe_c = C["dark_brown"]

    if season == "summer":
        leg_h_vis = int(leg_h * 0.55)  # shorts: shorter pants, skin legs
        # shorts
        add_rect(slide, body_x + _i(0.01), by + bh, int(bw * 0.42), _i(0.08), body_color)
        add_rect(slide, body_x + int(bw * 0.55), by + bh, int(bw * 0.42), _i(0.08), body_color)
        # skin legs below shorts
        add_rect(slide, body_x + _i(0.02), by + bh + _i(0.08), int(bw * 0.35), leg_h_vis, C["skin"])
        add_rect(slide, body_x + int(bw * 0.58), by + bh + _i(0.08), int(bw * 0.35), leg_h_vis, C["skin"])
        # sandals
        shoe_y = by + bh + _i(0.08) + leg_h_vis - _i(0.02)
        add_oval(slide, body_x, shoe_y, _i(0.12), _i(0.06), C["warm_yellow"])
        add_oval(slide, body_x + int(bw * 0.55), shoe_y, _i(0.12), _i(0.06), C["warm_yellow"])
    elif pose == "walk":
        # legs spread apart (stride)
        add_rect(slide, body_x - _i(0.02), by + bh, int(bw * 0.38), leg_h, leg_color)
        add_rect(slide, body_x + int(bw * 0.62), by + bh, int(bw * 0.38), leg_h, leg_color)
        add_oval(slide, body_x - _i(0.04), by + bh + leg_h - _i(0.02), _i(0.12), _i(0.07), shoe_c)
        add_oval(slide, body_x + int(bw * 0.58), by + bh + leg_h - _i(0.02), _i(0.12), _i(0.07), shoe_c)
    elif pose == "sit":
        # legs bent forward (seated)
        add_rect(slide, body_x + _i(0.02), by + bh, int(bw * 0.40), int(leg_h * 0.5), leg_color)
        add_rect(slide, body_x + int(bw * 0.55), by + bh, int(bw * 0.40), int(leg_h * 0.5), leg_color)
        # feet forward
        add_oval(slide, body_x + _i(0.02), by + bh + int(leg_h * 0.45), _i(0.12), _i(0.07), shoe_c)
        add_oval(slide, body_x + int(bw * 0.55), by + bh + int(leg_h * 0.45), _i(0.12), _i(0.07), shoe_c)
    else:
        # normal standing legs (slightly apart for liveliness)
        add_rect(slide, body_x + _i(0.01), by + bh, int(bw * 0.38), leg_h, leg_color)
        add_rect(slide, body_x + int(bw * 0.55), by + bh, int(bw * 0.38), leg_h, leg_color)
        add_oval(slide, body_x - _i(0.01), by + bh + leg_h - _i(0.02), _i(0.12), _i(0.07), shoe_c)
        add_oval(slide, body_x + int(bw * 0.50), by + bh + leg_h - _i(0.02), _i(0.12), _i(0.07), shoe_c)

    # ── winter boots (chunkier) ──
    if season == "winter" and pose not in ("sit",):
        boot_y = by + bh + leg_h - _i(0.04)
        add_oval(slide, body_x - _i(0.02), boot_y, _i(0.14), _i(0.09), C["dark_brown"])
        add_oval(slide, body_x + int(bw * 0.50), boot_y, _i(0.14), _i(0.09), C["dark_brown"])

def ox(slide, ox_x, oy, scale=1.0):
    # body
    add_oval(slide, ox_x, oy, int(I(0.55*scale)), int(I(0.26*scale)), C["ox_body"])
    # irregular white patches
    add_oval(slide, ox_x+int(I(0.10*scale)), oy+int(I(0.02*scale)), int(I(0.12*scale)), int(I(0.08*scale)), C["white"])
    add_oval(slide, ox_x+int(I(0.30*scale)), oy+int(I(0.12*scale)), int(I(0.14*scale)), int(I(0.10*scale)), C["white"])
    # head
    add_oval(slide, ox_x+int(I(0.46*scale)), oy-int(I(0.06*scale)),
             int(I(0.18*scale)), int(I(0.16*scale)), C["ox_body"])
    # eyes
    add_oval(slide, ox_x+int(I(0.58*scale)), oy-int(I(0.02*scale)), int(I(0.03*scale)), int(I(0.03*scale)), C["ink"])
    # horns
    add_oval(slide, ox_x+int(I(0.48*scale)), oy-int(I(0.18*scale)),
             int(I(0.05*scale)), int(I(0.12*scale)), C["dark_brown"])
    add_oval(slide, ox_x+int(I(0.56*scale)), oy-int(I(0.18*scale)),
             int(I(0.05*scale)), int(I(0.12*scale)), C["dark_brown"])
    # saddle
    add_rect(slide, ox_x+int(I(0.16*scale)), oy+int(I(0.02*scale)),
             int(I(0.16*scale)), int(I(0.14*scale)), C["red"])
    add_oval(slide, ox_x+int(I(0.15*scale)), oy, int(I(0.18*scale)), int(I(0.05*scale)), C["gold"])
    # legs
    for lx in [0.06,0.15,0.33,0.42]:
        add_rect(slide, ox_x+int(I(lx*scale)), oy+int(I(0.22*scale)),
                 int(I(0.06*scale)), int(I(0.18*scale)), C["dark_brown"])
    add_oval(slide, ox_x-int(I(0.04*scale)), oy+int(I(0.05*scale)),
             int(I(0.08*scale)), int(I(0.18*scale)), C["ox_body"])

def lotus(slide, cx, cy, scale=1.0):
    add_oval(slide, cx-int(I(0.22*scale)), cy+int(I(0.05*scale)),
             int(I(0.44*scale)), int(I(0.18*scale)), C["dark_green"])
    add_oval(slide, cx-int(I(0.14*scale)), cy+int(I(0.01*scale)),
             int(I(0.28*scale)), int(I(0.14*scale)), C["grass"])
    add_rect(slide, cx-int(I(0.02*scale)), cy-int(I(0.18*scale)),
             int(I(0.04*scale)), int(I(0.24*scale)), C["grass"])
    for a in range(0,360,60):
        rad=math.radians(a)
        add_oval(slide, cx+int(I(0.14*scale)*math.cos(rad))-int(I(0.08*scale)),
                 cy-int(I(0.10*scale))+int(I(0.14*scale)*math.sin(rad))-int(I(0.10*scale)),
                 int(I(0.16*scale)), int(I(0.14*scale)), C["pink"])
    add_oval(slide, cx-int(I(0.07*scale)), cy-int(I(0.17*scale)),
             int(I(0.14*scale)), int(I(0.12*scale)), C["warm_yellow"])

def kite(slide, kx, ky, scale=1.0):
    add_oval(slide, kx, ky, int(I(0.22*scale)), int(I(0.28*scale)), C["warm_yellow"])
    add_rect(slide, kx+int(I(0.05*scale)), ky, int(I(0.12*scale)), int(I(0.28*scale)),
             rgb(0x44,0x88,0xFF))
    for i in range(3):
        add_oval(slide, kx+int(I(0.06*scale)), ky+int(I((0.30+i*0.12)*scale)),
                 int(I(0.10*scale)), int(I(0.06*scale)),
                 C["rose"] if i%2==0 else C["warm_yellow"])

def snowman(slide, cx, cy, scale=1.0):
    add_oval(slide, cx-int(I(0.18*scale)), cy, int(I(0.36*scale)), int(I(0.30*scale)), rgb(0xEE,0xEE,0xFF))
    add_oval(slide, cx-int(I(0.14*scale)), cy-int(I(0.26*scale)), int(I(0.28*scale)), int(I(0.26*scale)), rgb(0xF0,0xF0,0xFF))
    add_oval(slide, cx-int(I(0.12*scale)), cy-int(I(0.44*scale)), int(I(0.24*scale)), int(I(0.22*scale)), rgb(0xF8,0xF8,0xFF))
    add_rect(slide, cx-int(I(0.14*scale)), cy-int(I(0.58*scale)), int(I(0.28*scale)), int(I(0.04*scale)), C["deep_blue"])
    add_rect(slide, cx-int(I(0.10*scale)), cy-int(I(0.70*scale)), int(I(0.20*scale)), int(I(0.14*scale)), C["deep_blue"])
    add_rect(slide, cx-int(I(0.14*scale)), cy-int(I(0.30*scale)), int(I(0.28*scale)), int(I(0.05*scale)), C["red"])
    for bi in range(3):
        add_oval(slide, cx-int(I(0.03*scale)), cy-int(I((0.22-bi*0.08)*scale)),
                 int(I(0.06*scale)), int(I(0.06*scale)), C["grey"])
    add_rect(slide, cx+int(I(0.14*scale)), cy-int(I(0.28*scale)),
             int(I(0.04*scale)), int(I(0.38*scale)), C["brown"])
    add_oval(slide, cx+int(I(0.10*scale)), cy+int(I(0.06*scale)),
             int(I(0.12*scale)), int(I(0.20*scale)), C["warm_yellow"])

def rain_drops(slide, x, y, w, h, n=12):
    for i in range(n):
        rx=x+int(w*(i%4)/4)+int(I(0.04)); ry=y+int(h*(i//4)/max(1,n//4))
        add_oval(slide, rx, ry, int(I(0.05)), int(I(0.10)), C["deep_blue"])

def wheat_field(slide, x, y, w, h, rows=3, cols=6, color=None):
    c=color or C["wheat"]; sw=w//cols; sh=h//rows
    for r in range(rows):
        for c2 in range(cols):
            wx=x+c2*sw+int(sw*0.35); wy=y+r*sh
            add_rect(slide, wx, wy+int(sh*0.30), int(I(0.03)), int(sh*0.70), C["grass"])
            add_oval(slide, wx-int(I(0.05)), wy+int(sh*0.04), int(I(0.11)), int(sh*0.30), c)

def wavy_hills(slide, x, y, w, color, n=4):
    step=w//n
    for i in range(n):
        add_oval(slide, x+i*step-step//3, y, int(step*1.0), int(I(0.55)), color)

def house(slide, hx, hy, scale=1.0):
    add_rect(slide, hx, hy, int(I(0.32*scale)), int(I(0.24*scale)), C["rose"])
    add_rect(slide, hx-int(I(0.04*scale)), hy-int(I(0.14*scale)),
             int(I(0.40*scale)), int(I(0.16*scale)), C["burnt_org"])
    add_rect(slide, hx+int(I(0.12*scale)), hy+int(I(0.10*scale)),
             int(I(0.08*scale)), int(I(0.14*scale)), C["dark_brown"])
    add_rect(slide, hx+int(I(0.03*scale)), hy+int(I(0.04*scale)),
             int(I(0.07*scale)), int(I(0.07*scale)), C["warm_yellow"])
    add_rect(slide, hx+int(I(0.22*scale)), hy+int(I(0.04*scale)),
             int(I(0.07*scale)), int(I(0.07*scale)), C["warm_yellow"])

def frog(slide, fx, fy, scale=1.0):
    add_oval(slide, fx, fy, int(I(0.26*scale)), int(I(0.18*scale)), C["grass"])
    add_oval(slide, fx+int(I(0.05*scale)), fy-int(I(0.12*scale)),
             int(I(0.16*scale)), int(I(0.14*scale)), C["grass"])
    add_oval(slide, fx+int(I(0.04*scale)), fy-int(I(0.14*scale)),
             int(I(0.06*scale)), int(I(0.06*scale)), C["white"])
    add_oval(slide, fx+int(I(0.12*scale)), fy-int(I(0.14*scale)),
             int(I(0.06*scale)), int(I(0.06*scale)), C["white"])
    add_oval(slide, fx+int(I(0.06*scale)), fy-int(I(0.12*scale)),
             int(I(0.03*scale)), int(I(0.03*scale)), C["ink"])
    add_oval(slide, fx+int(I(0.14*scale)), fy-int(I(0.12*scale)),
             int(I(0.03*scale)), int(I(0.03*scale)), C["ink"])
    add_oval(slide, fx-int(I(0.10*scale)), fy+int(I(0.08*scale)),
             int(I(0.16*scale)), int(I(0.07*scale)), C["dark_green"])
    add_oval(slide, fx+int(I(0.18*scale)), fy+int(I(0.08*scale)),
             int(I(0.16*scale)), int(I(0.07*scale)), C["dark_green"])

def egret(slide, ex, ey, scale=1.0):
    add_oval(slide, ex, ey, int(I(0.14*scale)), int(I(0.36*scale)), rgb(0xF0,0xF0,0xF0))
    add_oval(slide, ex+int(I(0.04*scale)), ey-int(I(0.12*scale)),
             int(I(0.10*scale)), int(I(0.12*scale)), rgb(0xF0,0xF0,0xF0))
    add_rect(slide, ex+int(I(0.12*scale)), ey-int(I(0.06*scale)),
             int(I(0.16*scale)), int(I(0.04*scale)), C["gold"])
    add_oval(slide, ex-int(I(0.10*scale)), ey+int(I(0.04*scale)),
             int(I(0.18*scale)), int(I(0.10*scale)), rgb(0xF0,0xF0,0xF0))

def watermelon(slide, wx, wy, scale=1.0):
    add_oval(slide, wx, wy, int(I(0.34*scale)), int(I(0.26*scale)), C["grass"])
    add_oval(slide, wx+int(I(0.05*scale)), wy+int(I(0.04*scale)),
             int(I(0.24*scale)), int(I(0.18*scale)), C["red"])
    for sx in [0.08,0.16,0.24]:
        add_rect(slide, wx+int(I(sx*scale)), wy, int(I(0.03*scale)), int(I(0.26*scale)), C["dark_green"])

def scarecrow(slide, sx, sy, scale=1.0):
    add_rect(slide, sx+int(I(0.08*scale)), sy, int(I(0.05*scale)), int(I(0.60*scale)), C["brown"])
    add_rect(slide, sx, sy+int(I(0.10*scale)), int(I(0.30*scale)), int(I(0.05*scale)), C["brown"])
    add_oval(slide, sx+int(I(0.02*scale)), sy-int(I(0.20*scale)),
             int(I(0.22*scale)), int(I(0.22*scale)), C["warm_yellow"])
    add_oval(slide, sx-int(I(0.02*scale)), sy-int(I(0.22*scale)),
             int(I(0.30*scale)), int(I(0.08*scale)), C["burnt_org"])
    add_rect(slide, sx+int(I(0.06*scale)), sy-int(I(0.36*scale)),
             int(I(0.14*scale)), int(I(0.16*scale)), C["burnt_org"])
    add_rect(slide, sx, sy+int(I(0.15*scale)), int(I(0.30*scale)), int(I(0.22*scale)), C["burnt_org"])

def snowy_house(slide, hx, hy, scale=1.0):
    add_rect(slide, hx, hy, int(I(0.40*scale)), int(I(0.28*scale)), rgb(0xAA,0xBB,0xDD))
    add_rect(slide, hx-int(I(0.04*scale)), hy-int(I(0.16*scale)),
             int(I(0.48*scale)), int(I(0.18*scale)), rgb(0x88,0x99,0xCC))
    add_oval(slide, hx-int(I(0.06*scale)), hy-int(I(0.20*scale)),
             int(I(0.52*scale)), int(I(0.12*scale)), C["white"])
    add_rect(slide, hx+int(I(0.06*scale)), hy+int(I(0.06*scale)),
             int(I(0.10*scale)), int(I(0.10*scale)), C["warm_yellow"])
    add_rect(slide, hx+int(I(0.24*scale)), hy+int(I(0.06*scale)),
             int(I(0.10*scale)), int(I(0.10*scale)), C["warm_yellow"])
    add_rect(slide, hx+int(I(0.14*scale)), hy+int(I(0.12*scale)),
             int(I(0.12*scale)), int(I(0.16*scale)), C["deep_blue"])

def cotton_plant(slide, cpx, cpy, scale=1.0):
    add_rect(slide, cpx+int(I(0.08*scale)), cpy,
             int(I(0.04*scale)), int(I(0.50*scale)), C["grass"])
    for branch_y,branch_dx in [(0.05,0.00),(0.18,-0.10),(0.30,0.08),(0.42,-0.06)]:
        bx=cpx+int(I((0.08+branch_dx)*scale)); by=cpy+int(I(branch_y*scale))
        add_rect(slide, bx, by, int(I(0.14*scale)), int(I(0.03*scale)), C["grass"])
        add_oval(slide,
                 bx+int(I(0.08*scale)) if branch_dx<0 else bx-int(I(0.04*scale)),
                 by-int(I(0.08*scale)),
                 int(I(0.14*scale)), int(I(0.14*scale)), C["white"])

def persimmon_tree(slide, ptx, pty, scale=1.0):
    add_rect(slide, ptx+int(I(0.08*scale)), pty+int(I(0.16*scale)),
             int(I(0.08*scale)), int(I(0.40*scale)), C["dark_brown"])
    add_oval(slide, ptx, pty, int(I(0.24*scale)), int(I(0.20*scale)), C["lavender"])
    for px2,py2 in [(0.02,0.00),(0.10,-0.06),(0.16,0.02),(0.06,0.08)]:
        add_oval(slide, ptx+int(I(px2*scale)), pty+int(I(py2*scale)),
                 int(I(0.08*scale)), int(I(0.10*scale)), C["orange"])

def butterfly(slide, bfx, bfy, scale=1.0):
    add_oval(slide, bfx, bfy, int(I(0.20*scale)), int(I(0.12*scale)), C["lavender"])
    add_oval(slide, bfx+int(I(0.16*scale)), bfy,
             int(I(0.20*scale)), int(I(0.12*scale)), C["pink"])
    add_rect(slide, bfx+int(I(0.07*scale)), bfy-int(I(0.06*scale)),
             int(I(0.06*scale)), int(I(0.20*scale)), C["ink"])

def seedling(slide, sx, sy, scale=1.0):
    add_rect(slide, sx+int(I(0.04*scale)), sy, int(I(0.04*scale)), int(I(0.26*scale)), C["grass"])
    add_oval(slide, sx, sy-int(I(0.06*scale)), int(I(0.14*scale)), int(I(0.10*scale)), C["dark_green"])
    add_oval(slide, sx+int(I(0.06*scale)), sy-int(I(0.12*scale)),
             int(I(0.10*scale)), int(I(0.08*scale)), C["grass"])

def ladybug(slide, lx, ly, scale=1.0):
    add_oval(slide, lx, ly, int(I(0.20*scale)), int(I(0.15*scale)), C["red"])
    add_oval(slide, lx+int(I(0.06*scale)), ly-int(I(0.08*scale)),
             int(I(0.08*scale)), int(I(0.08*scale)), C["ink"])
    for dx,dy in [(0.04,0.03),(0.04,0.09),(0.12,0.03),(0.12,0.09)]:
        add_oval(slide, lx+int(I(dx*scale)), ly+int(I(dy*scale)),
                 int(I(0.04*scale)), int(I(0.04*scale)), C["ink"])
    add_rect(slide, lx+int(I(0.09*scale)), ly, int(I(0.02*scale)), int(I(0.15*scale)), C["ink"])

def water_wheel(slide, wwx, wwy, scale=1.0):
    r=int(I(0.32*scale))
    add_oval(slide, wwx-r, wwy-r, r*2, r*2, C["orange"], lc=C["burnt_org"], lp=2)
    add_oval(slide, wwx-int(I(0.10*scale)), wwy-int(I(0.10*scale)),
             int(I(0.20*scale)), int(I(0.20*scale)), C["dark_brown"])
    for angle in range(0,360,45):
        rad=math.radians(angle)
        spx=wwx+int(r*0.5*math.cos(rad)); spy=wwy+int(r*0.5*math.sin(rad))
        add_rect(slide, min(wwx,spx)-1, min(wwy,spy)-1,
                 abs(wwx-spx)+3, abs(wwy-spy)+3, C["burnt_org"])
    for angle in range(0,360,45):
        rad=math.radians(angle+22)
        px=wwx+int((r-int(I(0.08*scale)))*math.cos(rad))
        py=wwy+int((r-int(I(0.08*scale)))*math.sin(rad))
        add_rect(slide, px, py, int(I(0.12*scale)), int(I(0.06*scale)), C["brown"])

def snow_flakes(slide, ax, ay, aw, ah, n=8):
    import random; random.seed(42)
    for _ in range(n):
        fx=ax+int(random.random()*aw); fy=ay+int(random.random()*ah)
        r2=int(I(0.04))
        add_oval(slide, fx-r2, fy-r2, r2*2, r2*2, C["white"])

def umbrella(slide, cx, cy, scale=1.0):
    add_oval(slide, cx-int(I(0.32*scale)), cy-int(I(0.12*scale)),
             int(I(0.64*scale)), int(I(0.30*scale)), C["deep_blue"])
    add_rect(slide, cx-int(I(0.03*scale)), cy, int(I(0.06*scale)), int(I(0.42*scale)), C["dark_brown"])

def firefly(slide, ffx, ffy, scale=1.0):
    add_oval(slide, ffx, ffy, int(I(0.10*scale)), int(I(0.06*scale)), C["warm_yellow"])
    add_oval(slide, ffx+int(I(0.02*scale)), ffy-int(I(0.04*scale)),
             int(I(0.06*scale)), int(I(0.06*scale)), C["ink"])
    add_oval(slide, ffx-int(I(0.08*scale)), ffy-int(I(0.02*scale)),
             int(I(0.12*scale)), int(I(0.06*scale)), rgb(0xCC,0xFF,0xDD))
    add_oval(slide, ffx+int(I(0.08*scale)), ffy-int(I(0.02*scale)),
             int(I(0.12*scale)), int(I(0.06*scale)), rgb(0xCC,0xFF,0xDD))
