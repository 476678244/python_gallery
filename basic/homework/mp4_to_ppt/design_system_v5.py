"""
Scandinavian Flat Design System v5 - Horizontal Card Layout
横向卡片式布局设计系统
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb(r,g,b): return RGBColor(r,g,b)
def I(v): return Inches(v)

# ============================================================================
# 1. 色板系统 - Pastel Nature Palette (参考图片风格)
# ============================================================================
PALETTE = dict(
    # 主色调 - 更柔和的 pastel
    spring_green   = rgb(0xA8, 0xC8, 0x88),  # 草绿
    warm_yellow    = rgb(0xF0, 0xD8, 0x78),  # 暖黄
    soft_coral     = rgb(0xE8, 0x98, 0x88),  # 珊瑚
    rain_blue      = rgb(0x90, 0xB8, 0xD0),  # 雨蓝
    flower_pink    = rgb(0xF0, 0xA8, 0xB8),  # 花粉
    sky_blue       = rgb(0xA8, 0xC8, 0xE8),  # 天蓝
    wheat_gold     = rgb(0xD8, 0xC0, 0x78),  # 麦金
    autumn_amber   = rgb(0xD8, 0xA8, 0x68),  # 秋琥珀
    winter_blue    = rgb(0xB8, 0xC8, 0xD8),  # 冬蓝
    cream_white    = rgb(0xF8, 0xF5, 0xF0),  # 米白
    
    # 中性色
    paper_white    = rgb(0xFA, 0xF8, 0xF5),  # 纸白背景
    ink_dark       = rgb(0x2A, 0x35, 0x30),  # 深墨绿
    ink_brown      = rgb(0x4A, 0x40, 0x35),  # 深棕
    warm_gray      = rgb(0x88, 0x80, 0x78),  # 暖灰
    light_gray     = rgb(0xC0, 0xBC, 0xB8),  # 浅灰
    
    # 季节背景 - 更柔和的色调
    spring_bg      = rgb(0xE8, 0xF0, 0xD8),  # 春背景
    spring_sky     = rgb(0xD8, 0xE8, 0xD0),  # 春天空
    spring_gnd     = rgb(0xB8, 0xD0, 0xA8),  # 春地面
    
    summer_bg      = rgb(0xF0, 0xF5, 0xE0),  # 夏背景
    summer_sky     = rgb(0xC8, 0xE0, 0xF0),  # 夏天空
    summer_gnd     = rgb(0xA0, 0xC8, 0x90),  # 夏地面
    
    autumn_bg      = rgb(0xF5, 0xE8, 0xD0),  # 秋背景
    autumn_sky     = rgb(0xF8, 0xE8, 0xC8),  # 秋天空
    autumn_gnd     = rgb(0xD0, 0xB8, 0x88),  # 秋地面
    
    winter_bg      = rgb(0xE8, 0xF0, 0xF5),  # 冬背景
    winter_sky     = rgb(0xD8, 0xE8, 0xF0),  # 冬天空
    winter_gnd     = rgb(0xE0, 0xE8, 0xF5),  # 冬地面
    
    # UI 强调色
    label_green    = rgb(0x88, 0xB8, 0x78),  # 标签绿
    accent_green   = rgb(0x90, 0xC0, 0x80),  # 强调绿
    sun_yellow     = rgb(0xF8, 0xE0, 0x60),  # 太阳黄
    cloud_gray     = rgb(0xC8, 0xD0, 0xD8),  # 云灰
)

# ============================================================================
# 2. 布局系统 - 2x2 网格
# ============================================================================
SW = Inches(13.33)  # 幻灯片宽度
SH = Inches(7.50)   # 幻灯片高度
CW = SW / 2         # Cell 宽度 (6.665")
CH = SH / 2         # Cell 高度 (3.75")
BW = int(Inches(0.03))  # 边框宽度

# 安全边距
MARGIN_LEFT = I(0.35)
MARGIN_TOP = I(0.30)

# ============================================================================
# 3. 字体系统 - 参考图片风格
# ============================================================================
FONTS = dict(
    chinese = "Source Han Sans SC",
    english = "Nunito",
    backup_zh = "Microsoft YaHei",
    backup_en = "Segoe UI",
)

TYPOGRAPHY = dict(
    # 标签 "二十四节气"
    label = dict(
        size = 9,
        bold = False,
        color = PALETTE["cream_white"],
        font = FONTS["chinese"],
    ),
    # 中文节气标题 - 大且粗
    title_zh = dict(
        size = 36,
        bold = True,
        color = PALETTE["ink_dark"],
        font = FONTS["chinese"],
    ),
    # 强调线颜色
    accent_line_color = PALETTE["accent_green"],
    accent_line_width = Pt(3),
    # 英文节气名
    title_en = dict(
        size = 11,
        bold = False,
        italic = False,
        color = PALETTE["ink_brown"],
        font = FONTS["english"],
    ),
    # 中文描述
    desc_zh = dict(
        size = 10,
        bold = False,
        italic = False,
        color = PALETTE["warm_gray"],
        font = FONTS["chinese"],
    ),
    # 英文描述 (斜体)
    desc_en = dict(
        size = 9,
        bold = False,
        italic = True,
        color = PALETTE["warm_gray"],
        font = FONTS["english"],
    ),
    # 底部说明文字
    footer = dict(
        size = 8,
        bold = False,
        italic = False,
        color = PALETTE["warm_gray"],
        font = FONTS["chinese"],
    ),
)

# ============================================================================
# 4. 绘图辅助函数
# ============================================================================
def add_rect(slide, x, y, w, h, fill, stroke=False, stroke_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if stroke:
        s.line.color.rgb = stroke_color or PALETTE["light_gray"]
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_rounded_rect(slide, x, y, w, h, fill, corner_radius=Pt(6)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def add_oval(slide, x, y, w, h, fill, stroke=False):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if stroke:
        s.line.color.rgb = PALETTE["light_gray"]
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, style_key="title_zh", align=PP_ALIGN.LEFT, color=None):
    style = TYPOGRAPHY[style_key]
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(style["size"])
    r.font.bold = style["bold"]
    r.font.italic = style.get("italic", False)
    r.font.name = style["font"]
    r.font.color.rgb = color or style["color"]
    return tb

def push_back(slide, shp):
    sp = shp._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)

# ============================================================================
# 5. 统一 Cell 头部组件 (参考图片风格)
# ============================================================================
def add_cell_header_v5(slide, cx, cy, zh_title, en_title, zh_desc, en_desc, footer_text, icon_type="sun"):
    """
    添加统一的 cell 标题区 - 参考图片风格
    
    布局：
    - 顶部绿色圆角标签 "二十四节气"
    - 大标题 + 绿色下划线
    - 英文名称 + 中文描述 (斜体英文)
    - 底部图标 + 说明文字
    """
    margin_left = MARGIN_LEFT
    margin_top = MARGIN_TOP
    
    # 1. 顶部绿色标签 "二十四节气"
    label_w = I(0.95)
    label_h = I(0.22)
    label = add_rounded_rect(slide, cx + margin_left, cy + margin_top, label_w, label_h, PALETTE["label_green"])
    push_back(slide, label)
    
    # 标签文字
    add_text(slide, cx + margin_left, cy + margin_top + I(0.02), label_w, label_h, 
             "二十四节气", "label", align=PP_ALIGN.CENTER, color=PALETTE["cream_white"])
    
    # 2. 中文大标题
    title_y = cy + margin_top + I(0.35)
    add_text(slide, cx + margin_left, title_y, CW - I(0.5), I(0.55), 
             zh_title, "title_zh")
    
    # 3. 绿色强调线 (在标题下方)
    line_y = title_y + I(0.50)
    line_w = I(0.45)
    line_h = I(0.04)
    accent_line = add_rect(slide, cx + margin_left, line_y, line_w, line_h, 
                           PALETTE["accent_green"])
    
    # 4. 英文名称 (常规)
    en_y = line_y + I(0.12)
    add_text(slide, cx + margin_left, en_y, CW - I(0.5), I(0.20), 
             en_title, "title_en")
    
    # 5. 中文描述
    desc_y = en_y + I(0.22)
    add_text(slide, cx + margin_left, desc_y, CW - I(0.5), I(0.18), 
             zh_desc, "desc_zh")
    
    # 6. 英文描述 (斜体)
    en_desc_y = desc_y + I(0.18)
    add_text(slide, cx + margin_left, en_desc_y, CW - I(0.5), I(0.30), 
             en_desc, "desc_en")
    
    # 7. 底部说明栏
    footer_y = cy + CH - I(0.40)
    footer_icon_size = I(0.18)
    
    # 底部图标
    draw_footer_icon(slide, cx + margin_left, footer_y, footer_icon_size, icon_type)
    
    # 底部文字
    add_text(slide, cx + margin_left + footer_icon_size + I(0.08), footer_y + I(0.02), 
             CW - I(0.8), I(0.20), footer_text, "footer")

def draw_footer_icon(slide, x, y, size, icon_type):
    """绘制底部小图标"""
    if icon_type == "sun":
        # 太阳图标
        sun = add_oval(slide, x, y, size, size, PALETTE["sun_yellow"])
    elif icon_type == "rain":
        # 雨滴图标
        rain = add_oval(slide, x + size*0.3, y + size*0.2, size*0.4, size*0.6, PALETTE["rain_blue"])
    elif icon_type == "thunder":
        # 闪电图标 (黄色三角形)
        lightning = add_rect(slide, x + size*0.35, y, size*0.3, size, PALETTE["warm_yellow"])
    elif icon_type == "flower":
        # 花朵图标
        flower = add_oval(slide, x + size*0.2, y + size*0.2, size*0.6, size*0.6, PALETTE["flower_pink"])
    elif icon_type == "seedling":
        # 叶子/幼苗图标
        leaf = add_oval(slide, x + size*0.3, y, size*0.4, size*0.7, PALETTE["spring_green"])
    elif icon_type == "cloud":
        # 云朵图标
        cloud = add_oval(slide, x, y + size*0.2, size, size*0.5, PALETTE["cloud_gray"])
    else:
        # 默认圆点
        default = add_oval(slide, x + size*0.25, y + size*0.25, size*0.5, size*0.5, PALETTE["spring_green"])

# ============================================================================
# 6. Cell 背景系统
# ============================================================================
def cell_bg_unified(slide, col, row, bg_color, sky_color=None, ground_color=None):
    """统一 cell 背景 - 固定地平线位置"""
    cx = col * CW
    cy = row * CH
    
    # 背景色
    bg = add_rect(slide, cx, cy, CW, CH, bg_color)
    push_back(slide, bg)
    
    # 天空区域 (地平线以上)
    if sky_color:
        horizon_y = cy + CH * 0.72  # 地平线在 72% 位置
        sky = add_rect(slide, cx, cy, CW, horizon_y - cy, sky_color)
        push_back(slide, sky)
    
    # 地面区域 (地平线以下)
    if ground_color:
        horizon_y = cy + CH * 0.72
        ground_h = cy + CH - horizon_y
        ground = add_rect(slide, cx, horizon_y, CW, ground_h, ground_color)
        push_back(slide, ground)
    
    return cx, cy

def cell_border(slide, col, row):
    """统一白色边框"""
    cx = col * CW
    cy = row * CH
    for rx, ry, rw, rh in [(cx, cy, CW, BW), (cx, cy + CH - BW, CW, BW),
                            (cx, cy, BW, CH), (cx + CW - BW, cy, BW, CH)]:
        add_rect(slide, rx, ry, rw, rh, PALETTE["cream_white"])
