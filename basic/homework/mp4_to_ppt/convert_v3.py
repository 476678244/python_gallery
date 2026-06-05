"""
24 Solar Terms – Handmade Book PPT  v3
6 slides × 4 solar terms, 2×2 grid layout.
All illustrations drawn with python-pptx shapes only.
"""
import os, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── shared atoms ─────────────────────────────────────────────────────────────
from draw_atoms import (
    C, rgb, I, add_rect, add_oval, add_text, push_back,
    sun, grass_strip, tree_bare, tree_leafy, tree_pine,
    person, ox, lotus, kite, snowman, rain_drops, wheat_field,
    wavy_hills, house, frog, egret, watermelon, scarecrow,
    snowy_house, cotton_plant, persimmon_tree, butterfly,
    seedling, ladybug, water_wheel, snow_flakes, umbrella, firefly,
)

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "24节气_v3.pptx")
SW = Inches(13.33); SH = Inches(7.50)
CW = SW // 2;       CH = SH // 2
BW = int(Inches(0.04))   # border width

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])
def cell_xy(col, row): return col*CW, row*CH

# ── cell helpers ─────────────────────────────────────────────────────────────
def cell_bg(slide, col, row, bg, sky=None, gnd=None):
    cx, cy = cell_xy(col, row)
    b = add_rect(slide, cx, cy, CW, CH, bg); push_back(slide, b)
    if sky:
        sk = add_rect(slide, cx, cy, CW, int(CH*0.38), sky); push_back(slide, sk)
    if gnd:
        gd = add_rect(slide, cx, cy+CH-int(I(0.24)), CW, int(I(0.24)), gnd)
        push_back(slide, gd)

def cell_border(slide, col, row):
    cx, cy = cell_xy(col, row)
    for rx,ry,rw,rh in [(cx,cy,CW,BW),(cx,cy+CH-BW,CW,BW),
                         (cx,cy,BW,CH),(cx+CW-BW,cy,BW,CH)]:
        add_rect(slide, rx, ry, rw, rh, C["white"])

def cell_header(slide, col, row, zh, en, sub,
                zc=None, ec=None):
    cx, cy = cell_xy(col, row)
    zc = zc or C["ink"]; ec = ec or C["dark_brown"]
    # 往下挪一点，增加顶部边距 (Top Margin)
    top_pad = int(I(0.28))
    add_text(slide, cx+int(I(0.12)), cy+top_pad,
             int(I(1.10)), int(I(0.48)), zh, size=32, bold=True, color=zc)
    add_text(slide, cx+int(I(0.12)), cy+top_pad+int(I(0.52)),
             CW-int(I(0.24)), int(I(0.22)), en,
             size=14, bold=True, color=ec, font="Georgia")
    add_text(slide, cx+int(I(0.12)), cy+top_pad+int(I(0.72)),
             CW-int(I(0.24)), int(I(0.20)), sub,
             size=10, italic=True, color=ec, font="Georgia")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – 立春  雨水  惊蛰  春分
# ─────────────────────────────────────────────────────────────────────────────
def slide1(prs):
    s = blank(prs)
    # Title zone for each cell: ~(0..I(1.30)) wide, ~(I(0.22)..I(1.16)) tall
    # Keep this zone clear of foreground elements; only bg shapes allowed

    # ── 立春 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, C["spring_bg"], sky=C["spring_sky"], gnd=C["spring_gnd"])
    cx,cy = cell_xy(0,0)
    # BG: rolling hills (OK behind text)
    add_oval(s, cx-int(I(0.20)), cy+int(CH*0.34), int(I(1.50)), int(I(0.46)), rgb(0xB0,0xD8,0x90))
    add_oval(s, cx+int(I(0.80)), cy+int(CH*0.38), int(I(1.80)), int(I(0.44)), rgb(0xA0,0xD0,0x80))
    # clouds – right side of sky (avoid title zone)
    add_oval(s, cx+int(I(1.40)), cy+int(I(0.06)), int(I(0.56)), int(I(0.20)), C["white"])
    add_oval(s, cx+int(I(1.56)), cy+int(I(0.02)), int(I(0.36)), int(I(0.18)), C["white"])
    # sun – top-right corner
    sun(s, cx+CW-int(I(0.44)), cy+int(I(0.10)), int(I(0.26)))
    # house – right side, on hill
    house(s, cx+int(I(1.70)), cy+int(CH*0.32), scale=1.0)
    # bare tree – right of center (avoid title zone left)
    tree_bare(s, cx+int(I(1.36)), cy+int(CH*0.22), scale=1.5)
    for budx,budy in [(0.00,0.04),(0.16,0.00),(0.28,0.08),(-0.06,0.12),(0.10,0.18)]:
        add_oval(s, cx+int(I(1.36+budx)), cy+int(CH*0.22)-int(I(0.50))+int(I(budy)),
                 int(I(0.09)), int(I(0.07)), C["pink"])
    # smaller bare tree – far right edge
    tree_bare(s, cx+CW-int(I(0.10)), cy+int(CH*0.30), scale=1.1)
    for budx,budy in [(0.00,0.04),(0.12,0.00)]:
        add_oval(s, cx+CW-int(I(0.10))+int(I(budx)), cy+int(CH*0.30)-int(I(0.36))+int(I(budy)),
                 int(I(0.07)), int(I(0.05)), C["pink"])
    # big ox – center, below title
    ox(s, cx+int(I(0.20)), cy+int(CH*0.52), scale=1.8)
    # child sitting on ox back (positioned on saddle area)
    person(s, cx+int(I(0.50)), cy+int(CH*0.26), rgb(0xFF,0x88,0x66), scale=0.75)
    # kite – top right area (clear of title)
    kite(s, cx+int(I(1.60)), cy+int(I(0.02)), scale=0.80)
    add_rect(s, cx+int(I(1.69)), cy+int(I(0.28)), int(I(0.02)), int(I(0.50)), C["gold"])
    # birds – scattered across sky, avoid title zone
    for bx,by in [(1.36,0.06),(1.56,0.03),(1.86,0.08),(2.08,0.04)]:
        add_rect(s, cx+int(I(bx)), cy+int(I(by)), int(I(0.12)), int(I(0.04)), C["ink"])
    # seedlings and flowers in foreground
    for sx in [0.06, 0.22, 0.40, 0.58, 0.96, 1.14, 1.32, 1.50, 1.68, 1.86, 2.04, 2.20]:
        seedling(s, cx+int(I(sx)), cy+int(CH*0.72), scale=0.90)
    for fx,fc in [(0.10,C["pink"]),(0.36,C["warm_yellow"]),(0.66,C["rose"]),
                  (0.94,C["lavender"]),(1.24,C["pink"]),(1.54,C["rose"]),
                  (1.84,C["warm_yellow"]),(2.12,C["pink"])]:
        add_oval(s, cx+int(I(fx)), cy+int(CH*0.68), int(I(0.10)), int(I(0.08)), fc)
        add_oval(s, cx+int(I(fx+0.03)), cy+int(CH*0.66), int(I(0.04)), int(I(0.04)), C["warm_yellow"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,0,0,"立春","Start of Spring","Spring begins. It gets warmer.")
    cell_border(s,0,0)

    # ── 雨水 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, C["spring_bg"], sky=C["spring_sky"], gnd=C["spring_gnd"])
    cx,cy = cell_xy(1,0)
    # rain cloud – right side of sky (avoid title left)
    cloud_cx = cx+int(I(1.60)); cloud_cy = cy+int(I(0.10))
    add_oval(s, cloud_cx-int(I(0.50)), cloud_cy, int(I(1.20)), int(I(0.46)), rgb(0xBB,0xCC,0xDD))
    add_oval(s, cloud_cx-int(I(0.20)), cloud_cy-int(I(0.20)), int(I(0.70)), int(I(0.40)), rgb(0xCC,0xDD,0xEE))
    add_oval(s, cloud_cx+int(I(0.20)), cloud_cy-int(I(0.10)), int(I(0.50)), int(I(0.34)), rgb(0xDD,0xEE,0xFF))
    # rain drops – only under cloud area
    rain_drops(s, cx+int(I(1.30)), cy+int(I(0.48)), int(I(1.10)), int(I(0.90)), n=22)
    # BG hill
    add_oval(s, cx+int(I(0.40)), cy+int(CH*0.38), int(I(2.10)), int(I(0.38)), rgb(0xA0,0xCC,0x80))
    # puddle on ground
    add_oval(s, cx+int(I(0.06)), cy+int(CH*0.72), int(I(2.40)), int(I(0.20)), rgb(0x88,0xBB,0xEE))
    for rpx in [0.24, 0.68, 1.14, 1.60, 2.00]:
        add_oval(s, cx+int(I(rpx)), cy+int(CH*0.74), int(I(0.14)), int(I(0.07)), rgb(0xAA,0xCC,0xEE))
    # umbrella + person – center-right (clear of title+subtitle)
    umbrella(s, cx+int(I(1.14)), cy+int(CH*0.28), scale=1.5)
    person(s, cx+int(I(0.96)), cy+int(CH*0.40), rgb(0x44,0xAA,0xFF), skirt=True, scale=1.25)
    # pond with lotuses – right side
    add_oval(s, cx+int(I(1.60)), cy+int(CH*0.54), int(I(0.94)), int(I(0.28)), rgb(0x88,0xBB,0xDD))
    lotus(s, cx+int(I(1.62)), cy+int(CH*0.44), scale=1.10)
    lotus(s, cx+int(I(2.06)), cy+int(CH*0.48), scale=0.85)
    lotus(s, cx+int(I(1.84)), cy+int(CH*0.52), scale=0.65)
    # frog on lily pad – far right
    add_oval(s, cx+int(I(1.44)), cy+int(CH*0.62), int(I(0.22)), int(I(0.10)), C["dark_green"])
    frog(s, cx+int(I(1.46)), cy+int(CH*0.54), scale=0.75)
    # seedlings in foreground
    for sx in [0.06, 0.22, 0.40, 0.58, 1.10, 1.30, 1.56, 1.80, 2.04, 2.26]:
        seedling(s, cx+int(I(sx)), cy+int(CH*0.72), scale=0.80)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,1,0,"雨水","Rain Water","More rain. Everything grows.")
    cell_border(s,1,0)

    # ── 惊蛰 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["pale_yellow"], sky=rgb(0xCC,0xEE,0xAA), gnd=C["spring_gnd"])
    cx,cy = cell_xy(0,1)
    # thunder clouds – right side of sky (avoid title zone)
    add_oval(s, cx+int(I(0.80)), cy+int(I(0.02)), int(I(1.60)), int(I(0.52)), rgb(0x88,0x99,0xAA))
    add_oval(s, cx+int(I(1.10)), cy-int(I(0.04)), int(I(0.90)), int(I(0.46)), rgb(0x77,0x88,0x99))
    add_oval(s, cx+int(I(1.60)), cy+int(I(0.04)), int(I(0.70)), int(I(0.38)), rgb(0x99,0xAA,0xBB))
    # sun peeking from far right
    sun(s, cx+CW-int(I(0.22)), cy+int(I(0.04)), int(I(0.20)))
    # lightning bolt – right of center
    bx2,by2 = cx+int(I(1.40)), cy+int(I(0.42))
    add_rect(s, bx2, by2, int(I(0.12)), int(I(0.28)), C["warm_yellow"])
    add_rect(s, bx2-int(I(0.12)), by2+int(I(0.20)), int(I(0.30)), int(I(0.10)), C["warm_yellow"])
    add_rect(s, bx2+int(I(0.04)), by2+int(I(0.28)), int(I(0.12)), int(I(0.26)), C["warm_yellow"])
    # second smaller bolt
    bx3,by3 = cx+int(I(1.80)), cy+int(I(0.40))
    add_rect(s, bx3, by3, int(I(0.08)), int(I(0.16)), C["warm_yellow"])
    add_rect(s, bx3-int(I(0.08)), by3+int(I(0.12)), int(I(0.20)), int(I(0.06)), C["warm_yellow"])
    add_rect(s, bx3+int(I(0.03)), by3+int(I(0.16)), int(I(0.08)), int(I(0.14)), C["warm_yellow"])
    # two background hills for depth
    add_oval(s, cx-int(I(0.20)), cy+int(CH*0.40), int(I(1.50)), int(I(0.44)), rgb(0xBB,0xDD,0x88))
    add_oval(s, cx+int(I(0.70)), cy+int(CH*0.44), int(I(1.90)), int(I(0.40)), rgb(0xCC,0xDD,0x99))
    # tree – right side
    tree_leafy(s, cx+int(I(1.96)), cy+int(CH*0.28), scale=1.2, lc=C["grass"])
    # big ladybug – center-left, below title
    ladybug(s, cx+int(I(0.16)), cy+int(CH*0.52), scale=1.9)
    # big frog – right of center
    frog(s, cx+int(I(1.40)), cy+int(CH*0.52), scale=1.4)
    # worm peeking from ground – center
    add_oval(s, cx+int(I(0.80)), cy+int(CH*0.66), int(I(0.16)), int(I(0.24)), rgb(0xCC,0x88,0x44))
    add_oval(s, cx+int(I(0.82)), cy+int(CH*0.58), int(I(0.12)), int(I(0.12)), rgb(0xCC,0x88,0x44))
    add_oval(s, cx+int(I(0.85)), cy+int(CH*0.56), int(I(0.04)), int(I(0.04)), C["ink"])
    add_oval(s, cx+int(I(0.89)), cy+int(CH*0.56), int(I(0.04)), int(I(0.04)), C["ink"])
    # snail – between ladybug and worm
    add_oval(s, cx+int(I(0.56)), cy+int(CH*0.64), int(I(0.16)), int(I(0.12)), rgb(0xCC,0x99,0x66))
    add_oval(s, cx+int(I(0.62)), cy+int(CH*0.58), int(I(0.08)), int(I(0.08)), rgb(0xDD,0xAA,0x77))
    add_oval(s, cx+int(I(0.52)), cy+int(CH*0.62), int(I(0.18)), int(I(0.14)), rgb(0xBB,0x88,0x44))
    # dense seedlings
    for i,(sxi,syi) in enumerate([(0.06,0.70),(0.24,0.66),(0.42,0.68),(0.60,0.64),
                                   (0.96,0.66),(1.14,0.62),(1.32,0.68),(1.50,0.64),
                                   (1.72,0.66),(1.90,0.62),(2.08,0.68),(2.24,0.64)]):
        seedling(s, cx+int(I(sxi)), cy+int(CH*syi), scale=1.0+i*0.02)
    # butterflies – right side
    butterfly(s, cx+int(I(1.14)), cy+int(CH*0.36), scale=1.1)
    butterfly(s, cx+int(I(1.70)), cy+int(CH*0.40), scale=0.8)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,0,1,"惊蛰","Awakening of Insects","Thunder wakes sleeping insects.")
    cell_border(s,0,1)

    # ── 春分 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["spring_bg"], sky=C["spring_sky"], gnd=C["spring_gnd"])
    cx,cy = cell_xy(1,1)
    # sun – right of title
    sun(s, cx+CW-int(I(0.44)), cy+int(I(0.08)), int(I(0.28)))
    # cloud – right side
    add_oval(s, cx+int(I(1.40)), cy+int(I(0.02)), int(I(0.54)), int(I(0.20)), C["white"])
    add_oval(s, cx+int(I(1.56)), cy-int(I(0.02)), int(I(0.34)), int(I(0.16)), C["white"])
    # balance scale – far right, clear of title (compact)
    pole_x = cx+int(I(1.90)); pole_y = cy+int(I(0.22))
    add_rect(s, pole_x, pole_y, int(I(0.06)), int(I(0.56)), C["dark_brown"])
    add_rect(s, pole_x-int(I(0.36)), pole_y+int(I(0.03)), int(I(0.78)), int(I(0.04)), C["dark_brown"])
    # sun pan (left of pole)
    add_oval(s, pole_x-int(I(0.40)), pole_y-int(I(0.02)), int(I(0.24)), int(I(0.14)), C["warm_yellow"])
    add_oval(s, pole_x-int(I(0.34)), pole_y-int(I(0.16)), int(I(0.14)), int(I(0.14)), C["warm_yellow"])
    # moon pan (right of pole)
    add_oval(s, pole_x+int(I(0.18)), pole_y-int(I(0.02)), int(I(0.24)), int(I(0.14)), rgb(0x88,0xCC,0xFF))
    add_oval(s, pole_x+int(I(0.24)), pole_y-int(I(0.16)), int(I(0.14)), int(I(0.14)), rgb(0xCC,0xDD,0xFF))
    # BG flower hill
    add_oval(s, cx-int(I(0.10)), cy+int(CH*0.32), int(I(2.70)), int(I(0.44)), rgb(0xB0,0xD8,0x98))
    # flower meadow – front row (below title zone y>CH*0.44)
    flower_colors = [C["rose"],C["lavender"],C["pink"],C["warm_yellow"],C["rose"],
                     C["lavender"],C["pink"],C["warm_yellow"],C["rose"],C["lavender"]]
    for i,fc2 in enumerate(flower_colors):
        fxx=cx+int(I(0.06+i*0.24)); fyy=cy+int(CH*0.50)-int(I(0.04*(i%2)))
        add_rect(s, fxx+int(I(0.07)), fyy, int(I(0.05)), int(I(0.30)), C["grass"])
        add_oval(s, fxx, fyy-int(I(0.18)), int(I(0.20)), int(I(0.18)), fc2)
        add_oval(s, fxx+int(I(0.05)), fyy-int(I(0.14)), int(I(0.10)), int(I(0.10)), C["warm_yellow"])
    # back row of smaller flowers (below subtitle zone)
    for i2,fc3 in enumerate([C["pink"],C["rose"],C["lavender"],C["warm_yellow"],
                              C["pink"],C["rose"],C["lavender"],C["warm_yellow"],C["rose"]]):
        fxx2=cx+int(I(0.14+i2*0.26)); fyy2=cy+int(CH*0.42)-int(I(0.03*(i2%2)))
        add_rect(s, fxx2+int(I(0.04)), fyy2, int(I(0.03)), int(I(0.20)), C["dark_green"])
        add_oval(s, fxx2, fyy2-int(I(0.12)), int(I(0.14)), int(I(0.12)), fc3)
        add_oval(s, fxx2+int(I(0.03)), fyy2-int(I(0.08)), int(I(0.08)), int(I(0.08)), C["warm_yellow"])
    # big butterfly – right side
    butterfly(s, cx+int(I(1.86)), cy+int(CH*0.32), scale=1.4)
    # smaller butterfly – center
    butterfly(s, cx+int(I(1.16)), cy+int(CH*0.26), scale=0.85)
    # girl in flower field – below title text
    person(s, cx+int(I(0.10)), cy+int(CH*0.48), rgb(0xFF,0xAA,0xCC), skirt=True, scale=1.2)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,1,1,"春分","Spring Equinox","Day and night are equal length.")
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – 清明  谷雨  立夏  小满
# ─────────────────────────────────────────────────────────────────────────────
def slide2(prs):
    s = blank(prs)

    # ── 清明 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, rgb(0xCC,0xD8,0xE4), sky=rgb(0xDD,0xEE,0xFF), gnd=C["spring_gnd"])
    cx,cy = cell_xy(0,0)
    add_oval(s, cx+int(I(0.08)), cy+int(CH*0.28), int(I(1.82)), int(I(0.82)), C["pink"])
    tree_leafy(s, cx+int(I(0.55)), cy+int(CH*0.20), scale=0.90, lc=rgb(0x88,0xCC,0x44))
    house(s, cx+int(I(0.88)), cy+int(CH*0.22), scale=0.88)
    kite(s, cx+int(I(0.20)), cy+int(I(0.18)), scale=1.0)
    add_rect(s, cx+int(I(0.31)), cy+int(I(0.46)), int(I(0.02)), int(I(0.92)), C["gold"])
    person(s, cx+int(I(0.24)), cy+int(CH*0.60), rgb(0xFF,0x88,0xAA), skirt=True, scale=0.90)
    # second kite in background
    kite(s, cx+int(I(0.78)), cy+int(I(0.10)), scale=0.60)
    add_rect(s, cx+int(I(0.87)), cy+int(I(0.26)), int(I(0.02)), int(I(0.60)), C["gold"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,0,0,"清明","Clear and Bright","Sunny day for outings.")
    cell_border(s,0,0)

    # ── 谷雨 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, rgb(0xCC,0xD8,0xE4), sky=rgb(0xAA,0xCC,0xFF), gnd=C["spring_gnd"])
    cx,cy = cell_xy(1,0)
    rain_drops(s, cx+int(I(0.08)), cy+int(I(0.30)), CW-int(I(0.16)), int(I(0.82)), n=24)
    # crop field with seedlings
    wheat_field(s, cx+int(I(0.08)), cy+int(CH*0.38), CW-int(I(0.16)), int(CH*0.30),
                rows=4, cols=7, color=C["grass"])
    person(s, cx+CW-int(I(0.62)), cy+int(CH*0.42), rgb(0x55,0x99,0xFF),
           hat=True, hat_color=C["red"], scale=0.85)
    # small label box
    lb_x = cx+int(I(0.80)); lb_y = cy+int(CH*0.37)
    add_rect(s, lb_x, lb_y, int(I(0.84)), int(I(0.24)), rgb(0xFF,0xFF,0xFF))
    add_text(s, lb_x+int(I(0.04)), lb_y+int(I(0.04)), int(I(0.76)), int(I(0.18)),
             "Rain helps crops grow well.", size=9, color=C["dark_green"], font="Georgia")
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["spring_gnd"])
    cell_header(s,1,0,"谷雨","Grain Rain","Rain helps crops grow well.")
    cell_border(s,1,0)

    # ── 立夏 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["summer_bg"], sky=C["summer_sky"], gnd=C["summer_gnd"])
    cx,cy = cell_xy(0,1)
    sun(s, cx+CW-int(I(0.42)), cy+int(I(0.20)), int(I(0.24)))
    tree_leafy(s, cx+int(I(1.40)), cy+int(CH*0.28), scale=1.22, lc=C["grass"])
    add_oval(s, cx+int(I(0.04)), cy+int(CH*0.62), int(I(1.06)), int(I(0.30)), C["pond"])
    # bridge arch
    for bri in range(9):
        add_rect(s, cx+int(I(0.05+bri*0.12)), cy+int(CH*0.66),
                 int(I(0.10)), int(I(0.05)), rgb(0xCC,0xBB,0xAA))
    lotus(s, cx+int(I(0.28)), cy+int(CH*0.54), scale=0.82)
    lotus(s, cx+int(I(0.70)), cy+int(CH*0.58), scale=0.72)
    lotus(s, cx+int(I(0.50)), cy+int(CH*0.66), scale=0.58)
    person(s, cx+int(I(1.02)), cy+int(CH*0.50), rgb(0x44,0xAA,0xFF), scale=0.90)
    add_rect(s, cx+int(I(1.27)), cy+int(CH*0.26), int(I(0.04)), int(I(0.56)), C["dark_brown"])
    add_oval(s, cx+int(I(1.14)), cy+int(CH*0.14), int(I(0.30)), int(I(0.22)), rgb(0xCC,0xEE,0xFF))
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["summer_gnd"])
    cell_header(s,0,1,"立夏","Start of Summer","Summer comes. Plants grow fast.")
    cell_border(s,0,1)

    # ── 小满 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["summer_bg"], sky=C["summer_sky"], gnd=C["summer_gnd"])
    cx,cy = cell_xy(1,1)
    sun(s, cx+int(I(0.38)), cy+int(I(0.18)), int(I(0.28)))
    # hill silhouette
    add_oval(s, cx+int(I(0.58)), cy+int(CH*0.34), int(I(1.24)), int(I(0.58)), C["pale_yellow"])
    water_wheel(s, cx+CW-int(I(0.60)), cy+int(CH*0.42), scale=1.12)
    add_oval(s, cx+int(I(0.78)), cy+int(CH*0.68), int(I(1.34)), int(I(0.22)), C["pond"])
    # water stream bubbles
    for wb in range(5):
        add_oval(s, cx+int(I(0.88+wb*0.26)), cy+int(CH*0.74),
                 int(I(0.08)), int(I(0.06)), rgb(0xAA,0xDD,0xFF))
    wheat_field(s, cx+int(I(0.06)), cy+int(CH*0.52), int(I(1.56)), int(CH*0.28),
                rows=2, cols=9, color=C["warm_yellow"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["summer_gnd"])
    cell_header(s,1,1,"小满","Grain Full","Grains start to become full.")
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – 芒种  夏至  小暑  大暑
# ─────────────────────────────────────────────────────────────────────────────
def slide3(prs):
    s = blank(prs)

    # ── 芒种 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, C["pale_yellow"], sky=C["summer_sky"], gnd=C["autumn_gnd"])
    cx,cy = cell_xy(0,0)
    for bx,by in [(0.28,0.08),(0.50,0.06),(0.72,0.09),(1.02,0.07),(1.30,0.08)]:
        add_rect(s, cx+int(I(bx)), cy+int(I(by)), int(I(0.12)), int(I(0.04)), C["ink"])
    # layered wheat hills (red waves like original)
    wavy_hills(s, cx+int(I(0.06)), cy+int(CH*0.30), CW-int(I(0.12)), C["wheat"], n=5)
    wavy_hills(s, cx+int(I(0.06)), cy+int(CH*0.42), CW-int(I(0.12)), C["orange"], n=5)
    for ri in range(5):
        add_oval(s, cx+int(I(0.08+ri*0.40)), cy+int(CH*0.36)+int(I(0.06*ri)),
                 int(I(0.64)), int(I(0.28)), rgb(0xEE,0x44,0x22))
    house(s, cx+int(I(0.18)), cy+int(CH*0.36), scale=0.72)
    egret(s, cx+int(I(0.16)), cy+int(CH*0.62), scale=1.0)
    wheat_field(s, cx+int(I(0.48)), cy+int(CH*0.64), CW-int(I(0.58)), int(CH*0.18),
                rows=1, cols=8, color=C["wheat"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["autumn_gnd"])
    cell_header(s,0,0,"芒种","Grain in Ear","Time for planting and harvesting.")
    cell_border(s,0,0)

    # ── 夏至 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, C["pale_yellow"], sky=C["summer_sky"], gnd=C["summer_gnd"])
    cx,cy = cell_xy(1,0)
    sun(s, cx+int(I(1.02)), cy+int(I(0.16)), int(I(0.32)))
    add_oval(s, cx+int(I(0.18)), cy+int(CH*0.38), CW-int(I(0.36)), int(CH*0.42), C["pond"])
    for i3,(lx3,ly3,sc3) in enumerate([(0.24,0.42,1.0),(0.82,0.36,0.88),
                                         (1.44,0.44,0.92),(1.92,0.38,0.76)]):
        lotus(s, cx+int(I(lx3)), cy+int(CH*ly3), scale=sc3)
    frog(s, cx+CW-int(I(0.60)), cy+int(CH*0.70), scale=1.02)
    add_oval(s, cx+CW-int(I(0.82)), cy+int(CH*0.72),
             int(I(0.72)), int(I(0.24)), C["dark_green"])
    # lily pads
    for lpx,lpy in [(0.22,0.62),(0.68,0.68),(1.22,0.64)]:
        add_oval(s, cx+int(I(lpx)), cy+int(CH*lpy), int(I(0.24)), int(I(0.12)), C["dark_green"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["summer_gnd"])
    cell_header(s,1,0,"夏至","Summer Solstice","Longest day of the year.")
    cell_border(s,1,0)

    # ── 小暑 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["purple"], sky=rgb(0x50,0x28,0x78), gnd=C["teal"])
    cx,cy = cell_xy(0,1)
    add_oval(s, cx-int(I(0.20)), cy+int(CH*0.44), CW+int(I(0.40)), int(CH*0.52), C["teal"])
    tree_leafy(s, cx+int(I(0.44)), cy+int(CH*0.24), scale=0.96,
               lc=C["burnt_org"], tc=C["dark_brown"])
    # rabbit
    add_oval(s, cx+int(I(0.20)), cy+int(CH*0.70), int(I(0.22)), int(I(0.14)), rgb(0xF0,0xF0,0xF0))
    add_oval(s, cx+int(I(0.24)), cy+int(CH*0.62), int(I(0.14)), int(I(0.12)), rgb(0xF0,0xF0,0xF0))
    add_oval(s, cx+int(I(0.25)), cy+int(CH*0.58), int(I(0.04)), int(I(0.08)), rgb(0xF8,0xAA,0xBB))
    add_oval(s, cx+int(I(0.31)), cy+int(CH*0.58), int(I(0.04)), int(I(0.08)), rgb(0xF8,0xAA,0xBB))
    for ffx2,ffy2 in [(1.18,0.22),(1.62,0.18),(0.78,0.16),(1.90,0.26),(1.40,0.12)]:
        firefly(s, cx+int(I(ffx2)), cy+int(CH*ffy2), scale=0.90)
    # heat shimmer wavy lines
    for wi in range(4):
        add_oval(s, cx+int(I(0.90+wi*0.32)), cy+int(CH*0.38),
                 int(I(0.22)), int(I(0.06)), rgb(0xFF,0xFF,0x88))
    cell_header(s,0,1,"小暑","Minor Heat","It becomes very hot.",
                zc=C["white"], ec=rgb(0xEE,0xEE,0xFF))
    cell_border(s,0,1)

    # ── 大暑 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["purple"], sky=rgb(0x40,0x18,0x68), gnd=C["teal"])
    cx,cy = cell_xy(1,1)
    add_oval(s, cx-int(I(0.20)), cy+int(CH*0.44), CW+int(I(0.40)), int(CH*0.52), C["teal"])
    sun(s, cx+CW-int(I(0.48)), cy+int(I(0.14)), int(I(0.30)))
    # tent / awning
    add_rect(s, cx+int(I(1.02)), cy+int(I(0.18)), int(I(0.82)), int(I(0.36)), C["warm_yellow"])
    add_rect(s, cx+int(I(0.98)), cy+int(I(0.14)), int(I(0.90)), int(I(0.07)), C["deep_blue"])
    add_rect(s, cx+int(I(0.98)), cy+int(I(0.18)), int(I(0.06)), int(I(0.36)), C["deep_blue"])
    add_rect(s, cx+int(I(1.82)), cy+int(I(0.18)), int(I(0.06)), int(I(0.36)), C["deep_blue"])
    # girl holding watermelon
    person(s, cx+int(I(0.28)), cy+int(CH*0.46), C["rose"], skirt=True, scale=1.02)
    watermelon(s, cx+int(I(0.50)), cy+int(CH*0.52), scale=0.96)
    # ground watermelons
    for wxi in [0.08, 1.18, 1.52, 1.82]:
        watermelon(s, cx+int(I(wxi)), cy+int(CH*0.76), scale=0.64)
    # sweat drops on girl
    for swi in range(3):
        add_oval(s, cx+int(I(0.44+swi*0.12)), cy+int(CH*0.44),
                 int(I(0.06)), int(I(0.10)), rgb(0x88,0xCC,0xFF))
    cell_header(s,1,1,"大暑","Major Heat","Hottest time of the year.",
                zc=C["white"], ec=rgb(0xEE,0xEE,0xFF))
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – 立秋  处暑  白露  秋分
# ─────────────────────────────────────────────────────────────────────────────
def slide4(prs):
    s = blank(prs)

    # ── 立秋 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, C["autumn_bg"], sky=C["autumn_sky"], gnd=C["autumn_gnd"])
    cx,cy = cell_xy(0,0)
    sun(s, cx+CW-int(I(0.46)), cy+int(I(0.20)), int(I(0.22)))
    tree_leafy(s, cx+int(I(0.36)), cy+int(CH*0.22), scale=1.02, lc=C["burnt_org"])
    tree_leafy(s, cx+int(I(1.62)), cy+int(CH*0.20), scale=0.88, lc=C["warm_yellow"])
    for flx,fly,fc2 in [(0.52,0.46,C["burnt_org"]),(0.88,0.36,C["warm_yellow"]),
                         (1.32,0.50,C["orange"]),(1.74,0.42,C["burnt_org"]),
                         (0.68,0.52,C["warm_yellow"]),(1.52,0.38,C["orange"])]:
        add_oval(s, cx+int(I(flx)), cy+int(CH*fly), int(I(0.10)), int(I(0.08)), fc2)
    # easel scene
    person(s, cx+int(I(0.10)), cy+int(CH*0.50), C["lavender"], skirt=True, scale=0.90)
    add_rect(s, cx+int(I(0.40)), cy+int(CH*0.36), int(I(0.26)), int(I(0.24)), C["white"])
    add_rect(s, cx+int(I(0.38)), cy+int(CH*0.34), int(I(0.30)), int(I(0.04)), C["dark_brown"])
    add_rect(s, cx+int(I(0.36)), cy+int(CH*0.60), int(I(0.04)), int(I(0.24)), C["dark_brown"])
    add_rect(s, cx+int(I(0.62)), cy+int(CH*0.60), int(I(0.04)), int(I(0.24)), C["dark_brown"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["autumn_gnd"])
    cell_header(s,0,0,"立秋","Start of Autumn","Autumn begins. It cools down.")
    cell_border(s,0,0)

    # ── 处暑 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, C["autumn_bg"], sky=C["autumn_sky"], gnd=C["autumn_gnd"])
    cx,cy = cell_xy(1,0)
    sun(s, cx+int(I(0.40)), cy+int(I(0.20)), int(I(0.22)))
    wavy_hills(s, cx+int(I(0.06)), cy+int(CH*0.42), CW-int(I(0.12)),
               rgb(0xFF,0xCC,0x55), n=4)
    person(s, cx+int(I(0.96)), cy+int(CH*0.44), C["burnt_org"],
           hat=True, hat_color=C["warm_yellow"], scale=0.96)
    add_oval(s, cx+int(I(1.20)), cy+int(CH*0.54), int(I(0.42)), int(I(0.30)), C["brown"])
    for itx,itc in [(1.22,C["orange"]),(1.32,C["red"]),(1.46,C["warm_yellow"])]:
        add_oval(s, cx+int(I(itx)), cy+int(CH*0.50), int(I(0.12)), int(I(0.14)), itc)
    scarecrow(s, cx+int(I(1.68)), cy+int(CH*0.36), scale=0.92)
    # haystacks
    for hsx in [0.10, 0.40]:
        add_oval(s, cx+int(I(hsx)), cy+int(CH*0.60), int(I(0.30)), int(I(0.18)), C["wheat"])
        add_oval(s, cx+int(I(hsx+0.04)), cy+int(CH*0.56), int(I(0.22)), int(I(0.12)), C["gold"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["autumn_gnd"])
    cell_header(s,1,0,"处暑","End of Heat","Hot days are over.")
    cell_border(s,1,0)

    # ── 白露 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["off_white"], sky=rgb(0xCC,0xDD,0xEE), gnd=rgb(0xAA,0xCC,0x88))
    cx,cy = cell_xy(0,1)
    # morning mist
    add_oval(s, cx-int(I(0.10)), cy+int(CH*0.42), CW+int(I(0.20)), int(I(0.42)), rgb(0xEE,0xF4,0xFF))
    tree_leafy(s, cx+int(I(0.56)), cy+int(CH*0.20), scale=1.0, lc=C["grass"])
    # dew drops (blue teardrops)
    for dwx,dwy in [(0.26,0.32),(0.44,0.26),(0.68,0.38),(0.92,0.30),
                     (1.14,0.36),(1.36,0.24),(1.60,0.32),(1.82,0.28)]:
        add_oval(s, cx+int(I(dwx)), cy+int(CH*dwy), int(I(0.10)), int(I(0.14)), rgb(0xAA,0xCC,0xFF))
        add_oval(s, cx+int(I(dwx+0.02)), cy+int(CH*dwy)+int(I(0.02)),
                 int(I(0.04)), int(I(0.04)), C["white"])
    # geese V formation
    for gi,(gbx,gby) in enumerate([(0.28,0.09),(0.52,0.07),(0.76,0.05),
                                    (1.00,0.07),(1.24,0.09)]):
        add_oval(s, cx+int(I(gbx)), cy+int(CH*gby), int(I(0.16)), int(I(0.08)), rgb(0x88,0x88,0x99))
        add_oval(s, cx+int(I(gbx+0.04)), cy+int(CH*gby)-int(I(0.06)),
                 int(I(0.08)), int(I(0.08)), rgb(0x88,0x88,0x99))
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), rgb(0xAA,0xCC,0x88))
    cell_header(s,0,1,"白露","White Dew","Dew appears on the grass at dawn.")
    cell_border(s,0,1)

    # ── 秋分 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["autumn_bg"], sky=C["autumn_sky"], gnd=C["autumn_gnd"])
    cx,cy = cell_xy(1,1)
    # moon
    add_oval(s, cx+CW-int(I(0.54)), cy+int(I(0.18)), int(I(0.30)), int(I(0.30)), rgb(0xFF,0xF0,0xCC))
    # mid-autumn mooncake on table
    add_oval(s, cx+int(I(0.72)), cy+int(CH*0.60), int(I(0.38)), int(I(0.22)), C["warm_yellow"])
    add_oval(s, cx+int(I(0.78)), cy+int(CH*0.58), int(I(0.26)), int(I(0.18)), C["burnt_org"])
    # persimmon tree
    persimmon_tree(s, cx+int(I(1.52)), cy+int(CH*0.30), scale=1.0)
    # balance scale
    pole_x2=cx+int(I(1.04)); pole_y2=cy+int(I(0.24))
    add_rect(s, pole_x2, pole_y2, int(I(0.06)), int(I(0.55)), C["dark_brown"])
    add_rect(s, pole_x2-int(I(0.40)), pole_y2+int(I(0.02)), int(I(0.86)), int(I(0.04)), C["dark_brown"])
    add_oval(s, pole_x2-int(I(0.44)), pole_y2-int(I(0.02)), int(I(0.24)), int(I(0.14)), C["warm_yellow"])
    add_oval(s, pole_x2-int(I(0.38)), pole_y2-int(I(0.16)), int(I(0.14)), int(I(0.14)), C["warm_yellow"])
    add_oval(s, pole_x2+int(I(0.22)), pole_y2-int(I(0.02)), int(I(0.24)), int(I(0.14)), rgb(0xCC,0xDD,0xFF))
    add_oval(s, pole_x2+int(I(0.28)), pole_y2-int(I(0.16)), int(I(0.14)), int(I(0.14)), rgb(0xCC,0xDD,0xFF))
    # family watching moon (3 figures)
    for fi,(fxi2,fc3) in enumerate([(0.10,rgb(0xAA,0x88,0xFF)),
                                     (0.34,rgb(0xFF,0x88,0xAA)),
                                     (0.56,rgb(0xAA,0xCC,0xFF))]):
        person(s, cx+int(I(fxi2)), cy+int(CH*0.50), fc3, skirt=(fi==1), scale=0.80+fi*0.04)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["autumn_gnd"])
    cell_header(s,1,1,"秋分","Autumn Equinox","Equal day and night again.")
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – 寒露  霜降  立冬  小雪
# ─────────────────────────────────────────────────────────────────────────────
def slide5(prs):
    s = blank(prs)

    # ── 寒露 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, C["off_white"], sky=rgb(0xCC,0xDD,0xEE), gnd=rgb(0xAA,0xCC,0x88))
    cx,cy = cell_xy(0,0)
    cotton_plant(s, cx+int(I(0.22)), cy+int(CH*0.12), scale=1.02)
    cotton_plant(s, cx+int(I(0.58)), cy+int(CH*0.20), scale=0.88)
    cotton_plant(s, cx+int(I(1.60)), cy+int(CH*0.16), scale=0.75)
    person(s, cx+int(I(0.84)), cy+int(CH*0.44), C["rose"], skirt=True, scale=0.90)
    # basket of cotton
    add_oval(s, cx+int(I(1.10)), cy+int(CH*0.60), int(I(0.32)), int(I(0.22)), C["brown"])
    add_oval(s, cx+int(I(1.12)), cy+int(CH*0.56), int(I(0.28)), int(I(0.14)), C["white"])
    for dwi,(dwx2,dwy2) in enumerate([(0.30,0.22),(0.44,0.16),(0.70,0.28),(0.96,0.20)]):
        add_oval(s, cx+int(I(dwx2)), cy+int(CH*dwy2), int(I(0.08)), int(I(0.12)), rgb(0xAA,0xCC,0xFF))
        add_oval(s, cx+int(I(dwx2+0.02)), cy+int(CH*dwy2)+int(I(0.02)),
                 int(I(0.04)), int(I(0.04)), C["white"])
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), rgb(0xAA,0xCC,0x88))
    cell_header(s,0,0,"寒露","Cold Dew","Dew turns cold.")
    cell_border(s,0,0)

    # ── 霜降 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, C["off_white"], sky=rgb(0xBB,0xCC,0xDD), gnd=rgb(0x88,0xAA,0x66))
    cx,cy = cell_xy(1,0)
    persimmon_tree(s, cx+int(I(0.10)), cy+int(CH*0.08), scale=1.18)
    # frost sparkles
    for fi3 in range(14):
        add_oval(s, cx+int(I(0.14+fi3*0.16)), cy+int(CH*0.80),
                 int(I(0.06)), int(I(0.04)), C["white"])
        add_oval(s, cx+int(I(0.18+fi3*0.16)), cy+int(CH*0.84),
                 int(I(0.04)), int(I(0.03)), rgb(0xCC,0xDD,0xFF))
    # table with persimmons
    add_rect(s, cx+int(I(0.74)), cy+int(CH*0.60), int(I(0.62)), int(I(0.06)), rgb(0xBB,0x99,0x77))
    add_rect(s, cx+int(I(0.76)), cy+int(CH*0.66), int(I(0.06)), int(I(0.18)), rgb(0xAA,0x88,0x66))
    add_rect(s, cx+int(I(1.26)), cy+int(CH*0.66), int(I(0.06)), int(I(0.18)), rgb(0xAA,0x88,0x66))
    for oi in range(4):
        add_oval(s, cx+int(I(0.82+oi*0.14)), cy+int(CH*0.54),
                 int(I(0.12)), int(I(0.14)), C["orange"])
    person(s, cx+int(I(0.92)), cy+int(CH*0.38), rgb(0x88,0xAA,0xDD), scale=0.88)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), rgb(0x88,0xAA,0x66))
    cell_header(s,1,0,"霜降","Frost's Descent","Frost comes. It gets cold.")
    cell_border(s,1,0)

    # ── 立冬 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["winter_bg"], sky=C["winter_sky"], gnd=C["pale_yellow"])
    cx,cy = cell_xy(0,1)
    # family (dad+daughter+mum)
    person(s, cx+int(I(0.10)), cy+int(CH*0.28), rgb(0x88,0xAA,0xDD), scale=1.02)
    person(s, cx+int(I(0.44)), cy+int(CH*0.34), rgb(0xFF,0xAA,0xCC), skirt=True, scale=0.80)
    person(s, cx+int(I(0.70)), cy+int(CH*0.28), rgb(0xAA,0xBB,0xDD), scale=1.02)
    # winter feast table
    add_rect(s, cx+int(I(0.06)), cy+int(CH*0.66), CW-int(I(0.12)), int(I(0.08)), C["warm_yellow"])
    add_rect(s, cx+int(I(0.10)), cy+int(CH*0.74), int(I(0.08)), int(I(0.14)), C["brown"])
    add_rect(s, cx+CW-int(I(0.18)), cy+int(CH*0.74), int(I(0.08)), int(I(0.14)), C["brown"])
    for di2,(dish_x,dish_c) in enumerate([(0.12,C["orange"]),(0.34,C["rose"]),(0.56,C["grass"]),
                                           (0.78,C["warm_yellow"]),(1.00,C["red"]),
                                           (1.22,C["teal"]),(1.44,C["purple"]),(1.66,C["burnt_org"])]):
        add_oval(s, cx+int(I(dish_x)), cy+int(CH*0.60), int(I(0.22)), int(I(0.16)), dish_c)
        add_oval(s, cx+int(I(dish_x+0.02)), cy+int(CH*0.58), int(I(0.18)), int(I(0.08)),
                 rgb(0xFF,0xFF,0xFF))
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["pale_yellow"])
    cell_header(s,0,1,"立冬","Start of Winter","Winter comes.")
    cell_border(s,0,1)

    # ── 小雪 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["winter_bg"], sky=C["winter_sky"], gnd=C["winter_gnd"])
    cx,cy = cell_xy(1,1)
    wavy_hills(s, cx+int(I(0.06)), cy+int(CH*0.55), CW-int(I(0.12)), C["white"], n=4)
    tree_pine(s, cx+int(I(0.36)), cy+int(CH*0.26), scale=1.02)
    tree_pine(s, cx+int(I(0.86)), cy+int(CH*0.20), scale=1.18)
    tree_pine(s, cx+int(I(1.46)), cy+int(CH*0.28), scale=0.92)
    # snow on treetops
    for tpx,tpy in [(0.26,0.25),(0.74,0.19),(1.36,0.27)]:
        add_oval(s, cx+int(I(tpx)), cy+int(CH*tpy), int(I(0.40)), int(I(0.12)), C["white"])
    snow_flakes(s, cx+int(I(0.08)), cy+int(I(0.10)), CW-int(I(0.16)), int(CH*0.48), n=14)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["winter_gnd"])
    cell_header(s,1,1,"小雪","Minor Snow","Light snow falls.")
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – 大雪  冬至  小寒  大寒
# ─────────────────────────────────────────────────────────────────────────────
def slide6(prs):
    s = blank(prs)

    # ── 大雪 (0,0) ────────────────────────────────────────────────────────────
    cell_bg(s,0,0, C["winter_bg"], sky=rgb(0xAA,0xBB,0xCC), gnd=C["winter_gnd"])
    cx,cy = cell_xy(0,0)
    snowy_house(s, cx+int(I(0.14)), cy+int(CH*0.36), scale=1.06)
    tree_pine(s, cx+int(I(1.44)), cy+int(CH*0.28), scale=0.96)
    tree_pine(s, cx+int(I(1.72)), cy+int(CH*0.34), scale=0.78)
    # snow on pines
    for tpx2,tpy2 in [(1.34,0.27),(1.64,0.33)]:
        add_oval(s, cx+int(I(tpx2)), cy+int(CH*tpy2), int(I(0.38)), int(I(0.10)), C["white"])
    # snowman beside house
    snowman(s, cx+int(I(1.08)), cy+int(CH*0.40), scale=0.90)
    # snow drifts
    add_oval(s, cx+int(I(0.06)), cy+int(CH*0.72), CW-int(I(0.12)), int(I(0.26)), C["white"])
    snow_flakes(s, cx+int(I(0.06)), cy+int(I(0.10)), CW-int(I(0.12)), int(CH*0.40), n=16)
    # road/path through snow
    add_oval(s, cx+int(I(0.50)), cy+int(CH*0.70), int(I(0.44)), int(I(0.18)), rgb(0xCC,0xCC,0xDD))
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["winter_gnd"])
    cell_header(s,0,0,"大雪","Major Snow","Heavy snow covers the ground.")
    cell_border(s,0,0)

    # ── 冬至 (1,0) ────────────────────────────────────────────────────────────
    cell_bg(s,1,0, rgb(0xDD,0xEE,0xFF), sky=C["winter_sky"], gnd=C["pale_yellow"])
    cx,cy = cell_xy(1,0)
    # room with window
    add_rect(s, cx+int(I(1.36)), cy+int(I(0.22)), int(I(0.42)), int(I(0.40)), rgb(0x88,0xBB,0xDD))
    add_rect(s, cx+int(I(1.36)), cy+int(I(0.22)), int(I(0.07)), int(I(0.40)), C["warm_yellow"])
    add_rect(s, cx+int(I(1.71)), cy+int(I(0.22)), int(I(0.07)), int(I(0.40)), C["warm_yellow"])
    add_rect(s, cx+int(I(1.36)), cy+int(I(0.56)), int(I(0.42)), int(I(0.04)), C["warm_yellow"])
    # table with dumplings (冬至吃饺子)
    add_rect(s, cx+int(I(0.52)), cy+int(CH*0.62), int(I(0.90)), int(I(0.07)), C["warm_yellow"])
    add_rect(s, cx+int(I(0.56)), cy+int(CH*0.69), int(I(0.07)), int(I(0.18)), C["brown"])
    add_rect(s, cx+int(I(1.30)), cy+int(CH*0.69), int(I(0.07)), int(I(0.18)), C["brown"])
    # dumpling row
    for di3 in range(5):
        dx3=cx+int(I(0.60+di3*0.16)); dy3=cy+int(CH*0.56)
        add_oval(s, dx3, dy3, int(I(0.14)), int(I(0.10)), rgb(0xFF,0xFF,0xDD))
        add_oval(s, dx3+int(I(0.04)), dy3-int(I(0.04)), int(I(0.06)), int(I(0.06)), C["white"])
    # girl at table eating
    person(s, cx+int(I(0.56)), cy+int(CH*0.42), rgb(0xFF,0x88,0xAA), skirt=True, scale=0.88)
    # mum serving
    person(s, cx+int(I(1.42)), cy+int(CH*0.38), rgb(0xFF,0xBB,0x88), skirt=True, scale=0.92)
    # bowl in hands
    add_oval(s, cx+int(I(1.64)), cy+int(CH*0.52), int(I(0.16)), int(I(0.10)), rgb(0xFF,0xFF,0xDD))
    # longest night label
    add_rect(s, cx+int(I(0.08)), cy+int(I(0.20)), int(I(0.80)), int(I(0.20)), rgb(0x22,0x22,0x44))
    add_text(s, cx+int(I(0.10)), cy+int(I(0.22)), int(I(0.76)), int(I(0.18)),
             "Longest night of the year", size=10, color=C["white"], font="Georgia")
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["pale_yellow"])
    cell_header(s,1,0,"冬至","Winter Solstice","Longest night of the year.")
    cell_border(s,1,0)

    # ── 小寒 (0,1) ────────────────────────────────────────────────────────────
    cell_bg(s,0,1, C["winter_bg"], sky=rgb(0x88,0xAA,0xCC), gnd=C["winter_gnd"])
    cx,cy = cell_xy(0,1)
    # plum blossom tree (梅花 – signature of 小寒)
    # trunk
    add_rect(s, cx+int(I(0.96)), cy+int(CH*0.30), int(I(0.08)), int(I(0.44)), C["dark_brown"])
    for bx3,by3,bw3,bh3 in [(-I(0.28),CH*0.28,I(0.32),I(0.06)),
                               ( I(0.06),CH*0.18,I(0.28),I(0.06)),
                               (-I(0.44),CH*0.20,I(0.24),I(0.05)),
                               ( I(0.08),CH*0.40,I(0.22),I(0.05))]:
        add_rect(s, cx+int(I(1.00)+bx3), cy+int(by3), int(bw3), int(bh3), C["dark_brown"])
    # plum blossoms (pink/red)
    for pmx,pmy in [(0.56,0.20),(0.70,0.14),(0.84,0.22),(1.02,0.12),(1.16,0.18),
                     (1.28,0.14),(0.62,0.30),(0.96,0.32),(1.22,0.28)]:
        for a2 in range(0,360,72):
            rad2=math.radians(a2)
            add_oval(s, cx+int(I(pmx))+int(I(0.07)*math.cos(rad2))-int(I(0.04)),
                     cy+int(CH*pmy)+int(I(0.07)*math.sin(rad2))-int(I(0.04)),
                     int(I(0.08)), int(I(0.08)),
                     C["rose"] if (pmx+pmy)%0.2 < 0.1 else C["pink"])
        add_oval(s, cx+int(I(pmx))-int(I(0.03)), cy+int(CH*pmy)-int(I(0.03)),
                 int(I(0.06)), int(I(0.06)), C["warm_yellow"])
    snow_flakes(s, cx+int(I(0.06)), cy+int(I(0.10)), CW-int(I(0.12)), int(CH*0.50), n=10)
    # child in winter coat
    person(s, cx+int(I(0.22)), cy+int(CH*0.52), rgb(0x88,0xBB,0xFF), scale=0.92)
    # breath puff
    add_oval(s, cx+int(I(0.44)), cy+int(CH*0.46), int(I(0.14)), int(I(0.10)),
             rgb(0xDD,0xEE,0xFF))
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["winter_gnd"])
    cell_header(s,0,1,"小寒","Minor Cold","The cold deepens. Plum blooms.")
    cell_border(s,0,1)

    # ── 大寒 (1,1) ────────────────────────────────────────────────────────────
    cell_bg(s,1,1, C["winter_bg"], sky=rgb(0x66,0x88,0xAA), gnd=C["winter_gnd"])
    cx,cy = cell_xy(1,1)
    # dramatic cold sky gradient (dark top)
    add_oval(s, cx-int(I(0.20)), cy-int(I(0.20)), CW+int(I(0.40)), int(CH*0.60),
             rgb(0x44,0x66,0x88))
    # frozen lake / ice cracking pattern
    add_oval(s, cx+int(I(0.08)), cy+int(CH*0.52), CW-int(I(0.16)), int(CH*0.24),
             rgb(0xCC,0xDD,0xEE))
    for cri in range(5):
        add_rect(s, cx+int(I(0.18+cri*0.40)), cy+int(CH*0.52),
                 int(I(0.02)), int(CH*0.24), rgb(0xAA,0xBB,0xCC))
    add_rect(s, cx+int(I(0.08)), cy+int(CH*0.60), CW-int(I(0.16)), int(I(0.02)),
             rgb(0xAA,0xBB,0xCC))
    # icicles hanging from top
    for icx in range(7):
        ix=cx+int(I(0.30+icx*0.28)); iy=cy+int(I(0.06))
        add_rect(s, ix, iy, int(I(0.05)), int(I(0.12+icx%3*0.04)), rgb(0xAA,0xCC,0xFF))
        add_oval(s, ix-int(I(0.01)), iy+int(I(0.10+icx%3*0.04)),
                 int(I(0.07)), int(I(0.05)), rgb(0xCC,0xEE,0xFF))
    # snowman in heavy coat
    snowman(s, cx+int(I(0.26)), cy+int(CH*0.34), scale=0.94)
    # two children ice-skating
    person(s, cx+int(I(1.08)), cy+int(CH*0.52), rgb(0xFF,0xAA,0xCC), skirt=True, scale=0.85)
    person(s, cx+int(I(1.38)), cy+int(CH*0.52), rgb(0x88,0xBB,0xFF), scale=0.85)
    snow_flakes(s, cx+int(I(0.06)), cy+int(I(0.06)), CW-int(I(0.12)), int(CH*0.48), n=14)
    grass_strip(s, cx+BW, cy+CH-int(I(0.24)), CW-BW*2, int(I(0.22)), C["winter_gnd"])
    cell_header(s,1,1,"大寒","Major Cold","Coldest time of the year.",
                zc=C["white"], ec=rgb(0xCC,0xDD,0xFF))
    cell_border(s,1,1)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    builders = [
        ("Slide 1: 立春 雨水 惊蛰 春分",   slide1),
        ("Slide 2: 清明 谷雨 立夏 小满",   slide2),
        ("Slide 3: 芒种 夏至 小暑 大暑",   slide3),
        ("Slide 4: 立秋 处暑 白露 秋分",   slide4),
        ("Slide 5: 寒露 霜降 立冬 小雪",   slide5),
        ("Slide 6: 大雪 冬至 小寒 大寒",   slide6),
    ]
    for name, fn in builders:
        print(f"  {name}")
        fn(prs)
    prs.save(OUT)
    print(f"\n✅  Saved: {OUT}")

if __name__ == "__main__":
    build()
