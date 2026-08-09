"""First-class PPT tools for SafeClaw `/ppt` mode.

SoT: docs/features/ppt-mode/methodology.md
"""

from __future__ import annotations

import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

_DECK_ID_RE = re.compile(r"^[a-zA-Z0-9][a-zA-Z0-9_-]{0,63}$")
_THEMES = {
    "default": {"bg": "FFFFFF", "title": "1E293B", "body": "334155", "accent": "2563EB"},
    "dark": {"bg": "0F172A", "title": "F8FAFC", "body": "CBD5E1", "accent": "38BDF8"},
    "warm": {"bg": "FFFBEB", "title": "78350F", "body": "92400E", "accent": "D97706"},
}


class PptToolError(ValueError):
    """Fail Fast error for PPT tools."""


@dataclass
class ImageRef:
    workspace_rel: str
    left_in: float = 5.5
    top_in: float = 1.5
    width_in: float = 4.0


@dataclass
class SlideState:
    title: str
    bullets: List[str] = field(default_factory=list)
    notes: str = ""
    layout: str = "title_bullets"
    images: List[ImageRef] = field(default_factory=list)


@dataclass
class DeckState:
    deck_id: str
    title: str
    theme_id: str = "default"
    slides: List[SlideState] = field(default_factory=list)
    current_version: Optional[int] = None
    dirty: bool = True


# session_id -> deck_id -> DeckState
_STORE: Dict[str, Dict[str, DeckState]] = {}


def clear_ppt_store() -> None:
    """Test helper."""
    _STORE.clear()


def _require_deck_id(deck_id: str) -> str:
    d = (deck_id or "").strip()
    if not d or not _DECK_ID_RE.match(d):
        raise PptToolError(
            "[ppt] Invalid deck_id\n"
            "  Expected: alphanumeric/underscore/hyphen, 1–64 chars, start alnum\n"
            f"  Actual: {deck_id!r}"
        )
    return d


def _ppt_root(workspace_dir: Path) -> Path:
    root = (workspace_dir / "ppt").resolve()
    root.mkdir(parents=True, exist_ok=True)
    return root


def _safe_under(root: Path, candidate: Path) -> Path:
    resolved = candidate.resolve()
    try:
        resolved.relative_to(root.resolve())
    except ValueError as e:
        raise PptToolError(
            "[ppt] Path escapes allowed root\n"
            f"  Root: {root}\n"
            f"  Path: {resolved}"
        ) from e
    return resolved


def _get_deck(session_id: str, deck_id: str) -> DeckState:
    sid = session_id or "_default"
    decks = _STORE.get(sid) or {}
    deck = decks.get(deck_id)
    if deck is None:
        raise PptToolError(
            "[ppt] Deck not initialized\n"
            f"  session_id: {sid}\n"
            f"  deck_id: {deck_id}\n"
            "  Hint: call safe_claw_ppt_deck_init first"
        )
    return deck


def _put_deck(session_id: str, deck: DeckState) -> None:
    sid = session_id or "_default"
    _STORE.setdefault(sid, {})[deck.deck_id] = deck


def _pptx_path(workspace_dir: Path, deck_id: str, version: int) -> Path:
    return _ppt_root(workspace_dir) / f"{deck_id}_v{version}.pptx"


def _preview_dir(workspace_dir: Path, deck_id: str, version: int) -> Path:
    return _ppt_root(workspace_dir) / "previews" / f"{deck_id}_v{version}"


def _next_version(workspace_dir: Path, deck_id: str) -> int:
    root = _ppt_root(workspace_dir)
    existing = list(root.glob(f"{deck_id}_v*.pptx"))
    nums: List[int] = []
    for p in existing:
        m = re.search(r"_v(\d+)\.pptx$", p.name)
        if m:
            nums.append(int(m.group(1)))
    return (max(nums) + 1) if nums else 1


def _serialize_deck(deck: DeckState) -> Dict[str, Any]:
    return {
        "deck_id": deck.deck_id,
        "title": deck.title,
        "theme_id": deck.theme_id,
        "current_version": deck.current_version,
        "dirty": deck.dirty,
        "slide_count": len(deck.slides),
        "slides": [
            {
                "index": i + 1,
                "title": s.title,
                "bullets": s.bullets,
                "notes": s.notes,
                "layout": s.layout,
                "images": [img.workspace_rel for img in s.images],
            }
            for i, s in enumerate(deck.slides)
        ],
    }


def _build_pptx(deck: DeckState, out_path: Path, workspace_dir: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise PptToolError(
            "[ppt] python-pptx is required\n"
            "  Hint: conda activate safe_claw && pip install python-pptx\n"
            f"  Error: {e}"
        ) from e

    theme = _THEMES.get(deck.theme_id, _THEMES["default"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(hex6: str) -> RGBColor:
        return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))

    for slide_state in deck.slides:
        slide = prs.slides.add_slide(blank)
        # background
        fill = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height,
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = _rgb(theme["bg"])
        fill.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = slide_state.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = _rgb(theme["title"])
        p.alignment = PP_ALIGN.LEFT

        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(7.5), Inches(5))
        btf = body.text_frame
        btf.word_wrap = True
        btf.clear()
        for i, bullet in enumerate(slide_state.bullets or [" "]):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = _rgb(theme["body"])

        if slide_state.notes:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = slide_state.notes
            except Exception as notes_err:
                logger.warning("[ppt] Could not set notes: %s", notes_err)

        for img in slide_state.images:
            src = _safe_under(workspace_dir.resolve(), (workspace_dir / img.workspace_rel).resolve())
            if not src.is_file():
                raise PptToolError(
                    "[ppt] Image not found for image_place\n"
                    f"  Path: {src}"
                )
            slide.shapes.add_picture(
                str(src),
                Inches(img.left_in),
                Inches(img.top_in),
                width=Inches(img.width_in),
            )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists():
        raise PptToolError(
            "[ppt] Refusing to overwrite existing version (create-only)\n"
            f"  Path: {out_path}\n"
            "  Hint: save_version always picks a new _vN"
        )
    prs.save(str(out_path))


def probe_preview_engine() -> str:
    """Return engine name or raise PptToolError."""
    try:
        import spire.presentation  # noqa: F401

        return "spire"
    except ImportError:
        pass
    try:
        import aspose.slides  # noqa: F401

        return "aspose"
    except ImportError:
        pass
    raise PptToolError(
        "[ppt] No preview renderer available\n"
        "  Tried: spire.presentation, aspose.slides\n"
        "  Hint: install one in conda env safe_claw"
    )


def render_preview(pptx_path: Path, out_dir: Path) -> List[Path]:
    engine = probe_preview_engine()
    out_dir.mkdir(parents=True, exist_ok=True)
    pngs: List[Path] = []

    if engine == "spire":
        from spire.presentation import Presentation

        pres = Presentation()
        try:
            pres.LoadFromFile(str(pptx_path))
            for i in range(pres.Slides.Count):
                image = pres.Slides[i].SaveAsImage()
                out = out_dir / f"slide_{i + 1:02d}.png"
                image.Save(str(out))
                pngs.append(out)
        finally:
            pres.Dispose()
    else:
        import aspose.slides as slides
        import aspose.pydrawing as drawing

        with slides.Presentation(str(pptx_path)) as pres:
            for i, slide in enumerate(pres.slides):
                image = slide.get_image(1.0, 1.0)
                out = out_dir / f"slide_{i + 1:02d}.png"
                image.save(str(out), drawing.imaging.ImageFormat.png)
                pngs.append(out)

    if not pngs:
        raise PptToolError(
            "[ppt] Preview produced zero slides\n"
            f"  pptx: {pptx_path}\n"
            f"  engine: {engine}"
        )
    return pngs


def build_ppt_tools(
    workspace_dir: Path,
    session_id: str = "_default",
    *,
    on_preview: Optional[Callable[[Dict[str, Any]], None]] = None,
) -> List[Any]:
    """Build LangChain tools for `/ppt` mode only."""
    from langchain_core.tools import tool

    ws = Path(workspace_dir)
    sid = session_id or "_default"

    @tool
    def safe_claw_ppt_deck_init(
        deck_id: str,
        title: str,
        theme_id: str = "default",
    ) -> str:
        """Initialize a PPT deck session (required before slide edits)."""
        did = _require_deck_id(deck_id)
        theme = (theme_id or "default").strip()
        if theme not in _THEMES:
            raise PptToolError(
                "[ppt] Unknown theme_id\n"
                f"  Expected: one of {sorted(_THEMES)}\n"
                f"  Actual: {theme_id!r}"
            )
        t = (title or "").strip()
        if not t:
            raise PptToolError("[ppt] deck title is required\n  Actual: empty")
        deck = DeckState(deck_id=did, title=t, theme_id=theme, slides=[], dirty=True)
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_deck_inspect(deck_id: str) -> str:
        """Return current deck structure JSON (pages, theme, version, dirty)."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        return json.dumps(_serialize_deck(deck), ensure_ascii=False, indent=2)

    @tool
    def safe_claw_ppt_slide_upsert(
        deck_id: str,
        slide_index: int,
        title: str,
        bullets: str = "",
        notes: str = "",
        layout: str = "title_bullets",
    ) -> str:
        """Create or replace a slide. slide_index is 1-based. bullets: newline-separated."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        t = (title or "").strip()
        if not t:
            raise PptToolError("[ppt] slide title is required\n  Actual: empty")
        idx = int(slide_index)
        if idx < 1:
            raise PptToolError(
                f"[ppt] slide_index must be >= 1\n  Actual: {slide_index}"
            )
        bullet_list = [b.strip() for b in (bullets or "").split("\n") if b.strip()]
        slide = SlideState(
            title=t,
            bullets=bullet_list,
            notes=(notes or "").strip(),
            layout=(layout or "title_bullets").strip(),
        )
        # Allow append at len+1 or replace existing
        if idx == len(deck.slides) + 1:
            deck.slides.append(slide)
        elif 1 <= idx <= len(deck.slides):
            prev_imgs = deck.slides[idx - 1].images
            slide.images = prev_imgs
            deck.slides[idx - 1] = slide
        else:
            raise PptToolError(
                "[ppt] slide_index out of range\n"
                f"  slide_count: {len(deck.slides)}\n"
                f"  Allowed: 1..{len(deck.slides)} (replace) or {len(deck.slides) + 1} (append)\n"
                f"  Actual: {idx}"
            )
        deck.dirty = True
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_slide_remove(deck_id: str, slide_index: int) -> str:
        """Remove a slide from deck session state (not FS delete). Cannot remove the last slide."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        if len(deck.slides) <= 1:
            raise PptToolError(
                "[ppt] Cannot remove the last slide\n"
                f"  deck_id: {deck.deck_id}\n"
                f"  slide_count: {len(deck.slides)}"
            )
        idx = int(slide_index)
        if idx < 1 or idx > len(deck.slides):
            raise PptToolError(
                f"[ppt] slide_index out of range\n  Actual: {idx} count={len(deck.slides)}"
            )
        deck.slides.pop(idx - 1)
        deck.dirty = True
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_slide_reorder(deck_id: str, order: str) -> str:
        """Reorder slides. order: comma-separated 1-based indices, e.g. '3,1,2'."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        parts = [p.strip() for p in (order or "").split(",") if p.strip()]
        try:
            indices = [int(p) for p in parts]
        except ValueError as e:
            raise PptToolError(
                f"[ppt] Invalid reorder order\n  Actual: {order!r}"
            ) from e
        n = len(deck.slides)
        if sorted(indices) != list(range(1, n + 1)):
            raise PptToolError(
                "[ppt] Reorder must be a permutation of 1..N\n"
                f"  N: {n}\n"
                f"  Actual: {indices}"
            )
        deck.slides = [deck.slides[i - 1] for i in indices]
        deck.dirty = True
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_theme_apply(deck_id: str, theme_id: str) -> str:
        """Apply a named theme (default|dark|warm)."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        theme = (theme_id or "").strip()
        if theme not in _THEMES:
            raise PptToolError(
                "[ppt] Unknown theme_id\n"
                f"  Expected: one of {sorted(_THEMES)}\n"
                f"  Actual: {theme_id!r}"
            )
        deck.theme_id = theme
        deck.dirty = True
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_image_place(
        deck_id: str,
        slide_index: int,
        workspace_rel_path: str,
        left_in: float = 5.5,
        top_in: float = 1.5,
        width_in: float = 4.0,
    ) -> str:
        """Place a workspace-relative image on a slide (applied on next save_version)."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        idx = int(slide_index)
        if idx < 1 or idx > len(deck.slides):
            raise PptToolError(
                f"[ppt] slide_index out of range for image_place\n  Actual: {idx}"
            )
        rel = (workspace_rel_path or "").strip().lstrip("/")
        if not rel:
            raise PptToolError("[ppt] workspace_rel_path is required")
        src = _safe_under(ws.resolve(), (ws / rel).resolve())
        if not src.is_file():
            raise PptToolError(f"[ppt] Image file not found\n  Path: {src}")
        deck.slides[idx - 1].images.append(
            ImageRef(
                workspace_rel=rel,
                left_in=float(left_in),
                top_in=float(top_in),
                width_in=float(width_in),
            )
        )
        deck.dirty = True
        _put_deck(sid, deck)
        return json.dumps({"ok": True, "deck": _serialize_deck(deck)}, ensure_ascii=False)

    @tool
    def safe_claw_ppt_save_version(deck_id: str) -> str:
        """Serialize deck to a new workspace/ppt/<deck_id>_vN.pptx (never overwrites)."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        if not deck.slides:
            raise PptToolError(
                "[ppt] Cannot save empty deck\n"
                "  Hint: safe_claw_ppt_slide_upsert at least one slide"
            )
        version = _next_version(ws, deck.deck_id)
        out = _pptx_path(ws, deck.deck_id, version)
        _build_pptx(deck, out, ws)
        deck.current_version = version
        deck.dirty = False
        _put_deck(sid, deck)
        rel = str(out.relative_to(ws.resolve()))
        payload = {
            "ok": True,
            "deck_id": deck.deck_id,
            "version": version,
            "pptx_path": rel,
            "slide_count": len(deck.slides),
        }
        return json.dumps(payload, ensure_ascii=False)

    @tool
    def safe_claw_ppt_preview(deck_id: str, version: int = 0) -> str:
        """Render a saved pptx version to PNG previews. Requires prior save_version. Fail Fast if dirty unsaved or no engine."""
        deck = _get_deck(sid, _require_deck_id(deck_id))
        if deck.dirty:
            raise PptToolError(
                "[ppt] Deck has unsaved changes — call safe_claw_ppt_save_version before preview\n"
                f"  deck_id: {deck.deck_id}\n"
                f"  dirty: true"
            )
        ver = int(version) if int(version) > 0 else deck.current_version
        if not ver:
            raise PptToolError(
                "[ppt] No saved version to preview\n"
                "  Hint: safe_claw_ppt_save_version first"
            )
        pptx = _pptx_path(ws, deck.deck_id, int(ver))
        if not pptx.is_file():
            raise PptToolError(f"[ppt] pptx not found\n  Path: {pptx}")
        out_dir = _preview_dir(ws, deck.deck_id, int(ver))
        pngs = render_preview(pptx, out_dir)
        rel_pptx = str(pptx.relative_to(ws.resolve()))
        preview_rels = [str(p.relative_to(ws.resolve())) for p in pngs]
        preview_urls = [f"/api/workspace-file?path={r}" for r in preview_rels]
        payload = {
            "ok": True,
            "type": "ppt_preview",
            "deck_id": deck.deck_id,
            "version": int(ver),
            "pptx_path": rel_pptx,
            "slide_count": len(pngs),
            "preview_paths": preview_rels,
            "preview_urls": preview_urls,
        }
        if on_preview:
            on_preview(payload)
        return json.dumps(payload, ensure_ascii=False)

    @tool
    def safe_claw_ppt_list_versions(deck_id: str) -> str:
        """List saved versions and whether preview PNGs exist."""
        did = _require_deck_id(deck_id)
        root = _ppt_root(ws)
        versions = []
        for p in sorted(root.glob(f"{did}_v*.pptx")):
            m = re.search(r"_v(\d+)\.pptx$", p.name)
            if not m:
                continue
            ver = int(m.group(1))
            prev = _preview_dir(ws, did, ver)
            pngs = sorted(prev.glob("slide_*.png")) if prev.is_dir() else []
            versions.append(
                {
                    "version": ver,
                    "pptx_path": str(p.relative_to(ws.resolve())),
                    "preview_ready": bool(pngs),
                    "slide_count": len(pngs),
                    "preview_urls": [
                        f"/api/workspace-file?path={png.relative_to(ws.resolve())}"
                        for png in pngs
                    ],
                }
            )
        return json.dumps(
            {"deck_id": did, "versions": versions},
            ensure_ascii=False,
            indent=2,
        )

    return [
        safe_claw_ppt_deck_init,
        safe_claw_ppt_deck_inspect,
        safe_claw_ppt_slide_upsert,
        safe_claw_ppt_slide_remove,
        safe_claw_ppt_slide_reorder,
        safe_claw_ppt_theme_apply,
        safe_claw_ppt_image_place,
        safe_claw_ppt_save_version,
        safe_claw_ppt_preview,
        safe_claw_ppt_list_versions,
    ]


PPT_TOOL_NAMES = frozenset(
    {
        "safe_claw_ppt_deck_init",
        "safe_claw_ppt_deck_inspect",
        "safe_claw_ppt_slide_upsert",
        "safe_claw_ppt_slide_remove",
        "safe_claw_ppt_slide_reorder",
        "safe_claw_ppt_theme_apply",
        "safe_claw_ppt_image_place",
        "safe_claw_ppt_save_version",
        "safe_claw_ppt_preview",
        "safe_claw_ppt_list_versions",
    }
)
