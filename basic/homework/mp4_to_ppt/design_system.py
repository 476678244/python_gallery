"""
Scandinavian Flat Design System for 24 Solar Terms
统一视觉系统规范
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(r,g,b): return RGBColor(r,g,b)
def I(v): return Inches(v)

# ============================================================================
# 1. 色板系统 - Scandinavian Soft Palette
# ============================================================================
PALETTE = dict(
    # 主色调
    spring_green   = rgb(0xB8, 0xD4, 0xA8),  # 柔和春绿
    warm_yellow    = rgb(0xE8, 0xD4, 0x88),  # 温暖黄
    soft_coral     = rgb(0xE8, 0xA0, 0x88),  # 柔珊瑚
    rain_blue      = rgb(0x98, 0xB8, 0xD8),  # 雨蓝
    flower_pink    = rgb(0xE8, 0xA8, 0xB8),  # 花粉
    sky_mint       = rgb(0xA8, 0xD8, 0xC8),  # 薄荷天
    wheat_gold     = rgb(0xD8, 0xC0, 0x78),  # 麦金
    autumn_amber   = rgb(0xD8, 0xA8, 0x68),  # 秋琥珀
    winter_blue    = rgb(0xB8, 0xC8, 0xD8),  # 冬蓝
    frost_white    = rgb(0xF0, 0xF4, 0xF8),  # 霜白
    
    # 中性色
    paper_white    = rgb(0xFA, 0xF8, 0xF4),  # 纸白背景
    ink_dark       = rgb(0x2A, 0x28, 0x26),  # 深墨 - 中文标题
    ink_brown      = rgb(0x5A, 0x4A, 0x3A),  # 深棕 - 英文副标题
    warm_gray      = rgb(0x8A, 0x80, 0x78),  # 暖灰
    light_gray     = rgb(0xC8, 0xC4, 0xC0),  # 浅灰
    
    # 季节背景
    spring_bg      = rgb(0xE0, 0xEC, 0xD8),  # 春背景
    spring_sky     = rgb(0xC8, 0xE0, 0xE8),  # 春天空
    spring_gnd     = rgb(0xA8, 0xC8, 0x98),  # 春地面
    
    summer_bg      = rgb(0xE8, 0xF0, 0xD8),  # 夏背景
    summer_sky     = rgb(0xB8, 0xD8, 0xE8),  # 夏天空
    summer_gnd     = rgb(0x90, 0xC0, 0x88),  # 夏地面
    
    autumn_bg      = rgb(0xF0, 0xE0, 0xC8),  # 秋背景
    autumn_sky     = rgb(0xF8, 0xE8, 0xC0),  # 秋天空
    autumn_gnd     = rgb(0xC8, 0xA8, 0x70),  # 秋地面
    
    winter_bg      = rgb(0xE0, 0xE8, 0xF0),  # 冬背景
    winter_sky     = rgb(0xC8, 0xD8, 0xE8),  # 冬天空
    winter_gnd     = rgb(0xD8, 0xE4, 0xF0),  # 冬地面
)

# ============================================================================
# 2. 布局系统 - Layout Grid & Safe Areas
# ============================================================================
# 基于 13.33" x 7.5" 的 2x2 网格，每个 cell 约 6.665" x 3.75"
# 安全区定义：
# - 左边距: 48px ≈ 0.33 inches (I(0.33))
# - 上边距: 48px ≈ 0.33 inches (I(0.33))
# - 标题区域高度: ~1.2 inches

SAFE_MARGIN_LEFT = I(0.40)    # 约 58px
SAFE_MARGIN_TOP = I(0.33)     # 约 48px
TITLE_ZONE_HEIGHT = I(1.20)   # 标题安全区高度
SUBTITLE_OFFSET = I(0.48)     # 英文副标题偏移
DESC_OFFSET = I(0.78)         # 描述文字偏移

# 统一地平线位置 (相对 cell 底部)
HORIZON_Y_RATIO = 0.78        # 地平线固定在 78% 高度处

# 统一太阳位置 (右上角黄金分割点)
SUN_POS_X_RATIO = 0.75        # 距离 cell 右边 25%
SUN_POS_Y_RATIO = 0.15        # 距离 cell 顶部 15%

# ============================================================================
# 3. 字体系统 - Typography System
# ============================================================================
FONTS = dict(
    chinese = "Source Han Sans SC",  # 思源黑体
    english = "Nunito",               # 圆润几何无衬线
    backup_zh = "Microsoft YaHei",
    backup_en = "Segoe UI",
)

TYPOGRAPHY = dict(
    # 中文节气标题
    title_zh = dict(
        size = 28,           # 稍小但更稳
        bold = True,
        color = PALETTE["ink_dark"],
        font = FONTS["chinese"],
    ),
    # 英文节气名
    title_en = dict(
        size = 13,
        bold = False,
        italic = True,
        color = PALETTE["ink_brown"],
        font = FONTS["english"],
    ),
    # 描述文字
    description = dict(
        size = 10,
        bold = False,
        italic = False,
        color = PALETTE["warm_gray"],
        font = FONTS["english"],
    ),
)

# ============================================================================
# 4. 描边与阴影规则
# ============================================================================
STROKE_WIDTH = 0              # 统一无描边
SHADOW_OPACITY = 0             # 统一无阴影 (纯扁平)

# ============================================================================
# 5. 构图系统 - Composition Rules
# ============================================================================
class Composition:
    """统一构图系统 - 提供标准化的元素放置规则"""
    
    @staticmethod
    def get_title_zone(cell_x, cell_y):
        """获取标题安全区位置"""
        return (
            cell_x + int(SAFE_MARGIN_LEFT),
            cell_y + int(SAFE_MARGIN_TOP),
            int(I(1.50)),  # 标题区域宽度
            int(TITLE_ZONE_HEIGHT)
        )
    
    @staticmethod
    def get_sun_position(cell_x, cell_y, cell_w, cell_h):
        """统一太阳位置 - 右上角黄金分割区"""
        return (
            cell_x + int(cell_w * SUN_POS_X_RATIO),
            cell_y + int(cell_h * SUN_POS_Y_RATIO)
        )
    
    @staticmethod
    def get_horizon_y(cell_y, cell_h):
        """统一地平线 Y 坐标"""
        return cell_y + int(cell_h * HORIZON_Y_RATIO)
    
    @staticmethod
    def get_ground_zone(cell_x, cell_y, cell_w, cell_h):
        """获取地面区域"""
        horizon_y = Composition.get_horizon_y(cell_y, cell_h)
        ground_h = cell_y + cell_h - horizon_y
        return (cell_x, horizon_y, cell_w, ground_h)

# ============================================================================
# 6. 绘图辅助函数
# ============================================================================
def add_rect(slide, x, y, w, h, fill, stroke=False):
    """添加矩形 - 统一无描边"""
    s = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if stroke:
        s.line.color.rgb = PALETTE["light_gray"]
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_oval(slide, x, y, w, h, fill, stroke=False):
    """添加椭圆 - 统一无描边"""
    s = slide.shapes.add_shape(9, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if stroke:
        s.line.color.rgb = PALETTE["light_gray"]
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, style_key="title_zh", align=PP_ALIGN.LEFT):
    """添加文字 - 使用预定义样式"""
    style = TYPOGRAPHY[style_key]
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(style["size"])
    r.font.bold = style["bold"]
    r.font.italic = style.get("italic", False)
    r.font.name = style["font"]
    r.font.color.rgb = style["color"]
    return tb

def add_cell_header(slide, cell_x, cell_y, zh, en, sub, comp=None):
    """添加统一的 cell 标题区"""
    comp = comp or Composition()
    tx, ty, tw, th = comp.get_title_zone(cell_x, cell_y)
    
    # 中文标题 - Bold
    add_text(slide, tx, ty, tw, int(I(0.40)), zh, "title_zh")
    
    # 英文标题 - Medium Italic, 深棕色
    add_text(slide, tx, ty + int(I(0.44)), tw, int(I(0.24)), en, "title_en")
    
    # 描述文字 - 浅灰
    add_text(slide, tx, ty + int(I(0.68)), tw, int(I(0.24)), sub, "description")

def push_back(slide, shp):
    """将形状移到最底层"""
    sp = shp._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)
